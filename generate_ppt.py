from __future__ import annotations

import argparse
import math
import re
import tomllib
from copy import deepcopy
from dataclasses import dataclass, field
from functools import lru_cache
from pathlib import Path
from typing import Sequence

from openpyxl import load_workbook
from openai import OpenAI
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.oxml.ns import qn
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt

from survey_stats import (
    RoleDefinition,
    TEMPLATE_DEFINITIONS,
    format_value,
    get_effective_role_definition,
)
from summary_table import SUMMARY_ROW_DEFINITIONS, normalize_text

VALID_SECTION_MODES = ("auto", "template", "summary")
DEFAULT_FILE_PATTERN = "*.xlsx"
DEFAULT_SHEET_NAME_MODE = "first"
DEFAULT_SECTION_MODE = "auto"

HEADER_FILL_COLOR = "BF1E4B"
OVERALL_FILL_COLOR = "D87A96"
SECTION_FILL_COLOR = "E49AAF"
BODY_FILL_COLOR = "F4E8EA"
BORDER_COLOR = "FFFFFF"
HEADER_TEXT_COLOR = "FFFFFF"
BODY_TEXT_COLOR = "4D5874"
DEFAULT_NOTES_TARGET_CHARS = 300
DEFAULT_NOTES_TEMPERATURE = 0.4
DEFAULT_NOTES_MAX_TOKENS = 500
DEFAULT_NOTES_CHECKPOINT_CHARS = 80
CATEGORY_LABEL_PREFIX_RE = re.compile(r"^[一二三四五六七八九十百零]+、")


@dataclass(frozen=True)
class TableRegion:
    left: float
    top: float
    width: float
    height: float

    def emu(self) -> tuple[int, int, int, int]:
        return (
            Inches(self.left),
            Inches(self.top),
            Inches(self.width),
            Inches(self.height),
        )


@dataclass(frozen=True)
class PptLayoutConfig:
    summary_table: TableRegion = field(
        default_factory=lambda: TableRegion(0.73, 1.45, 11.87, 0.56)
    )
    detail_single_table: TableRegion = field(
        default_factory=lambda: TableRegion(0.73, 2.10, 11.87, 4.95)
    )
    detail_left_table: TableRegion = field(
        default_factory=lambda: TableRegion(0.73, 2.10, 5.78, 4.95)
    )
    detail_right_table: TableRegion = field(
        default_factory=lambda: TableRegion(6.82, 2.10, 5.78, 4.95)
    )


@dataclass(frozen=True)
class LlmNotesConfig:
    enabled: bool = False
    env_path: Path = Path(".env")
    system_role_path: Path = Path("system_role.md")
    target_chars: int = DEFAULT_NOTES_TARGET_CHARS
    temperature: float = DEFAULT_NOTES_TEMPERATURE
    max_tokens: int = DEFAULT_NOTES_MAX_TOKENS
    checkpoint_chars: int = DEFAULT_NOTES_CHECKPOINT_CHARS


@dataclass(frozen=True)
class LlmRuntimeConfig:
    client: object
    model: str
    system_role: str
    target_chars: int
    temperature: float
    max_tokens: int
    checkpoint_chars: int


@dataclass(frozen=True)
class CategoryIntroSlideConfig:
    ppt_path: Path
    slide_number: int


@dataclass(frozen=True)
class PptBatchConfig:
    template_path: Path
    input_dir: Path
    output_ppt: Path
    file_pattern: str = DEFAULT_FILE_PATTERN
    sheet_name_mode: str = DEFAULT_SHEET_NAME_MODE
    sheet_name: str | None = None
    blank_display: str = ""
    title_suffix: str = ""
    section_mode: str = DEFAULT_SECTION_MODE
    max_single_table_rows: int = 18
    max_split_table_rows: int = 19
    sort_files: bool = True
    layout: PptLayoutConfig = field(default_factory=PptLayoutConfig)
    llm_notes: LlmNotesConfig = field(default_factory=LlmNotesConfig)
    body_font_size_pt: float = 10.5
    header_font_size_pt: float = 11.0
    summary_font_size_pt: float = 12.0
    template_slide_index: int = 0
    category_intro_slides: dict[str, CategoryIntroSlideConfig] = field(default_factory=dict)


@dataclass(frozen=True)
class SectionBlock:
    heading: str
    rows: tuple[tuple[str, float | None, float | None], ...]

    @property
    def row_count(self) -> int:
        return len(self.rows)


@dataclass(frozen=True)
class DetailLayout:
    is_split: bool
    single_rows: tuple[tuple[str, float | None, float | None], ...] = ()
    left_blocks: tuple[SectionBlock, ...] = ()
    right_blocks: tuple[SectionBlock, ...] = ()


@dataclass(frozen=True)
class WorkbookDisplayMeta:
    sort_index: int
    alias_index: int
    title: str
    category_label: str | None = None


def build_parser() -> argparse.ArgumentParser:
    parser = argparse.ArgumentParser(description="根据 Excel 批量生成 PPT")
    parser.add_argument("--config", type=Path, help="TOML 配置文件路径")
    parser.add_argument("--template-path", type=Path, help="PPT 模板路径")
    parser.add_argument("--input-dir", type=Path, help="Excel 输入目录")
    parser.add_argument("--output-ppt", type=Path, help="输出 PPT 路径")
    parser.add_argument(
        "--section-mode",
        choices=VALID_SECTION_MODES,
        help="二级标题识别口径：auto/template/summary",
    )
    parser.add_argument("--blank-display", help="空值显示文本，默认空字符串")
    parser.add_argument(
        "--dry-run",
        action="store_true",
        help="只校验输入和布局，不写出 PPT",
    )
    return parser


def load_batch_config(
    config_path: Path,
    *,
    template_path: Path | None = None,
    input_dir: Path | None = None,
    output_ppt: Path | None = None,
    section_mode: str | None = None,
    blank_display: str | None = None,
) -> PptBatchConfig:
    config_dir = config_path.parent
    raw = tomllib.loads(config_path.read_text(encoding="utf-8"))
    layout = load_layout_config(raw.get("layout", {}))
    llm_notes = load_llm_notes_config(config_dir, raw.get("llm_notes", {}))
    category_intro_slides = load_category_intro_slides_config(
        config_dir,
        raw.get("category_intro_slides"),
    )

    effective_section_mode = normalize_section_mode(section_mode or raw.get("section_mode"))
    return PptBatchConfig(
        template_path=resolve_config_path(config_dir, template_path or raw["template_path"]),
        input_dir=resolve_config_path(config_dir, input_dir or raw["input_dir"]),
        output_ppt=resolve_config_path(config_dir, output_ppt or raw["output_ppt"]),
        file_pattern=str(raw.get("file_pattern", DEFAULT_FILE_PATTERN)),
        sheet_name_mode=str(raw.get("sheet_name_mode", DEFAULT_SHEET_NAME_MODE)),
        sheet_name=raw.get("sheet_name"),
        blank_display=blank_display if blank_display is not None else str(raw.get("blank_display", "")),
        title_suffix=str(raw.get("title_suffix", "")),
        section_mode=effective_section_mode,
        max_single_table_rows=int(raw.get("max_single_table_rows", 18)),
        max_split_table_rows=int(raw.get("max_split_table_rows", 19)),
        sort_files=bool(raw.get("sort_files", True)),
        layout=layout,
        llm_notes=llm_notes,
        body_font_size_pt=float(raw.get("body_font_size_pt", 10.5)),
        header_font_size_pt=float(raw.get("header_font_size_pt", 11.0)),
        summary_font_size_pt=float(raw.get("summary_font_size_pt", 12.0)),
        template_slide_index=int(raw.get("template_slide_index", 0)),
        category_intro_slides=category_intro_slides,
    )


def load_llm_notes_config(config_dir: Path, raw: object) -> LlmNotesConfig:
    if raw is None:
        return LlmNotesConfig()
    if not isinstance(raw, dict):
        raise ValueError("llm_notes 必须是对象")
    return LlmNotesConfig(
        enabled=bool(raw.get("enabled", False)),
        env_path=resolve_config_path(config_dir, raw.get("env_path", ".env")),
        system_role_path=resolve_config_path(
            config_dir,
            raw.get("system_role_path", "system_role.md"),
        ),
        target_chars=int(raw.get("target_chars", DEFAULT_NOTES_TARGET_CHARS)),
        temperature=float(raw.get("temperature", DEFAULT_NOTES_TEMPERATURE)),
        max_tokens=int(raw.get("max_tokens", DEFAULT_NOTES_MAX_TOKENS)),
        checkpoint_chars=int(raw.get("checkpoint_chars", DEFAULT_NOTES_CHECKPOINT_CHARS)),
    )


def load_category_intro_slides_config(
    config_dir: Path,
    raw: object,
) -> dict[str, CategoryIntroSlideConfig]:
    if raw is None:
        return {}
    if not isinstance(raw, dict):
        raise ValueError("category_intro_slides 必须是对象")

    intro_slides: dict[str, CategoryIntroSlideConfig] = {}
    for category_label, item in raw.items():
        if not isinstance(category_label, str) or not category_label.strip():
            raise ValueError("category_intro_slides 的键必须是非空客户大类名称")
        if not isinstance(item, dict):
            raise ValueError(f"{category_label} 的章节页配置必须是对象")
        if "ppt_path" not in item:
            raise ValueError(f"{category_label} 缺少 ppt_path")
        if "slide_number" not in item:
            raise ValueError(f"{category_label} 缺少 slide_number")

        slide_number = int(item["slide_number"])
        if slide_number < 1:
            raise ValueError(f"{category_label} 的 slide_number 必须从 1 开始")

        intro_slides[category_label] = CategoryIntroSlideConfig(
            ppt_path=resolve_config_path(config_dir, item["ppt_path"]),
            slide_number=slide_number,
        )
    return intro_slides


def resolve_config_path(config_dir: Path, raw_path: str | Path) -> Path:
    path = Path(raw_path)
    if path.is_absolute():
        return path
    return config_dir / path


def load_layout_config(raw: dict[str, object]) -> PptLayoutConfig:
    return PptLayoutConfig(
        summary_table=load_table_region(
            raw.get("summary_table"),
            TableRegion(0.73, 1.45, 11.87, 0.56),
        ),
        detail_single_table=load_table_region(
            raw.get("detail_single_table"),
            TableRegion(0.73, 2.10, 11.87, 4.95),
        ),
        detail_left_table=load_table_region(
            raw.get("detail_left_table"),
            TableRegion(0.73, 2.10, 5.78, 4.95),
        ),
        detail_right_table=load_table_region(
            raw.get("detail_right_table"),
            TableRegion(6.82, 2.10, 5.78, 4.95),
        ),
    )


def load_table_region(raw: object, defaults: TableRegion) -> TableRegion:
    if raw is None:
        return defaults
    if not isinstance(raw, dict):
        raise ValueError("layout 下的表格区域必须是对象")
    return TableRegion(
        left=float(raw.get("left", defaults.left)),
        top=float(raw.get("top", defaults.top)),
        width=float(raw.get("width", defaults.width)),
        height=float(raw.get("height", defaults.height)),
    )


def normalize_section_mode(section_mode: str | None) -> str:
    normalized = str(section_mode or DEFAULT_SECTION_MODE).strip().lower()
    if normalized not in VALID_SECTION_MODES:
        raise ValueError(f"section_mode 仅支持: {', '.join(VALID_SECTION_MODES)}")
    return normalized


def strip_category_label_prefix(category_label: str) -> str:
    return CATEGORY_LABEL_PREFIX_RE.sub("", category_label).strip()


@lru_cache(maxsize=1)
def build_workbook_display_lookup() -> dict[str, tuple[int, int, str, str]]:
    lookup: dict[str, tuple[int, int, str, str]] = {}
    for definition_index, definition in enumerate(SUMMARY_ROW_DEFINITIONS):
        display_key = normalize_text(definition.display_name)
        if display_key and display_key not in lookup:
            lookup[display_key] = (
                definition_index,
                -1,
                definition.category_label,
                definition.display_name,
            )

        for alias_index, alias in enumerate(definition.source_aliases):
            alias_key = normalize_text(alias)
            if alias_key and alias_key not in lookup:
                lookup[alias_key] = (
                    definition_index,
                    alias_index,
                    definition.category_label,
                    definition.display_name,
                )
    return lookup


def resolve_workbook_display_meta(workbook_name: str) -> WorkbookDisplayMeta:
    entry = build_workbook_display_lookup().get(normalize_text(workbook_name))
    if entry is None:
        return WorkbookDisplayMeta(
            sort_index=len(SUMMARY_ROW_DEFINITIONS),
            alias_index=0,
            title=workbook_name,
            category_label=None,
        )

    definition_index, alias_index, category_label, display_name = entry
    return WorkbookDisplayMeta(
        sort_index=definition_index,
        alias_index=max(alias_index, 0),
        title=f"{strip_category_label_prefix(category_label)}——{display_name}",
        category_label=category_label,
    )


def build_partial_output_path(output_ppt: Path) -> Path:
    return output_ppt.with_name(f"{output_ppt.stem}.partial{output_ppt.suffix}")


def discover_input_files(config: PptBatchConfig) -> list[Path]:
    if not config.input_dir.exists():
        raise FileNotFoundError(f"输入目录不存在: {config.input_dir}")
    files = list(config.input_dir.glob(config.file_pattern))
    files = [path for path in files if path.is_file()]
    if config.sort_files:
        files.sort(
            key=lambda path: (
                resolve_workbook_display_meta(path.stem).sort_index,
                resolve_workbook_display_meta(path.stem).alias_index,
                path.name,
            )
        )
    if not files:
        raise FileNotFoundError(
            f"{config.input_dir} 下没有匹配 {config.file_pattern} 的 Excel 文件"
        )
    return files


def read_report_rows(
    workbook_path: Path,
    *,
    sheet_name_mode: str = DEFAULT_SHEET_NAME_MODE,
    sheet_name: str | None = None,
) -> list[tuple[str, float | None, float | None]]:
    workbook = load_workbook(workbook_path, data_only=True, read_only=True)
    try:
        if sheet_name_mode == "first":
            worksheet = workbook[workbook.sheetnames[0]]
        elif sheet_name_mode == "named":
            if not sheet_name:
                raise ValueError("sheet_name_mode=named 时必须提供 sheet_name")
            worksheet = workbook[sheet_name]
        else:
            raise ValueError("sheet_name_mode 仅支持 first 或 named")

        rows = list(worksheet.iter_rows(values_only=True))
    finally:
        workbook.close()

    if not rows:
        raise ValueError(f"{workbook_path} 为空工作簿")

    header = tuple("" if value is None else str(value).strip() for value in rows[0][:3])
    if header[:3] != ("指标", "满意度", "重要性"):
        raise ValueError(
            f"{workbook_path} 表头不符合预期，需为 指标/满意度/重要性，实际为: {header}"
        )

    report_rows: list[tuple[str, float | None, float | None]] = []
    for raw_row in rows[1:]:
        label = "" if raw_row[0] is None else str(raw_row[0]).strip()
        if not label:
            continue
        report_rows.append((label, to_optional_float(raw_row[1]), to_optional_float(raw_row[2])))

    if not report_rows:
        raise ValueError(f"{workbook_path} 没有可用数据行")
    return report_rows


def to_optional_float(value: object) -> float | None:
    if value is None:
        return None
    if isinstance(value, float):
        if math.isnan(value):
            return None
        return value
    try:
        return float(value)
    except (TypeError, ValueError):
        return None


def resolve_section_definition(
    role_name: str,
    rows: Sequence[tuple[str, float | None, float | None]],
    *,
    section_mode: str = DEFAULT_SECTION_MODE,
) -> RoleDefinition | None:
    base_definition = next(
        (
            template
            for template in TEMPLATE_DEFINITIONS.values()
            if template.role_name == role_name
        ),
        None,
    )
    if base_definition is None:
        return None

    normalized_mode = normalize_section_mode(section_mode)
    if normalized_mode == "template":
        return get_effective_role_definition(base_definition, "template")
    if normalized_mode == "summary":
        return get_effective_role_definition(base_definition, "summary")

    candidates: list[RoleDefinition] = []
    for mode in ("template", "summary"):
        candidate = get_effective_role_definition(base_definition, mode)
        if candidate.sections not in [existing.sections for existing in candidates]:
            candidates.append(candidate)

    labels = {row[0] for row in rows}
    return max(
        candidates,
        key=lambda definition: sum(
            1 for section in definition.sections if section.name in labels
        ),
    )


def build_section_blocks(
    rows: Sequence[tuple[str, float | None, float | None]],
    role_definition: RoleDefinition | None,
) -> list[SectionBlock]:
    if not rows:
        return []

    start_index = 1 if role_definition and rows[0][0] == role_definition.role_name else 0
    detail_rows = rows[start_index:]
    if not detail_rows:
        return []

    if role_definition is None:
        return [SectionBlock(heading=detail_rows[0][0], rows=tuple(detail_rows))]

    section_names = {section.name for section in role_definition.sections}
    blocks: list[SectionBlock] = []
    current_rows: list[tuple[str, float | None, float | None]] = []
    current_heading: str | None = None

    for row in detail_rows:
        label = row[0]
        if label in section_names:
            if current_rows:
                blocks.append(SectionBlock(heading=current_heading or current_rows[0][0], rows=tuple(current_rows)))
            current_heading = label
            current_rows = [row]
            continue

        if not current_rows:
            current_heading = current_heading or label
            current_rows = [row]
            continue

        current_rows.append(row)

    if current_rows:
        blocks.append(SectionBlock(heading=current_heading or current_rows[0][0], rows=tuple(current_rows)))

    return blocks


def filter_empty_satisfaction_sections(
    detail_rows: Sequence[tuple[str, float | None, float | None]],
    role_definition: RoleDefinition | None,
) -> list[tuple[str, float | None, float | None]]:
    if not detail_rows or role_definition is None:
        return list(detail_rows)

    filtered_blocks: list[SectionBlock] = []
    for block in build_section_blocks(detail_rows, role_definition):
        metric_rows = block.rows[1:]
        if metric_rows and all(satisfaction is None for _, satisfaction, _ in metric_rows):
            continue
        filtered_blocks.append(block)
    return flatten_blocks(filtered_blocks)


def choose_detail_layout(
    *,
    detail_rows: Sequence[tuple[str, float | None, float | None]],
    role_definition: RoleDefinition | None,
    max_single_table_rows: int,
    max_split_table_rows: int,
) -> DetailLayout:
    if len(detail_rows) <= max_single_table_rows:
        return DetailLayout(is_split=False, single_rows=tuple(detail_rows))

    section_blocks = build_section_blocks(detail_rows, role_definition)
    if len(section_blocks) < 2:
        raise ValueError("数据超出单表容量，但无法按二级标题拆分为左右双表")

    candidates: list[tuple[int, int, int]] = []
    for split_index in range(1, len(section_blocks)):
        left_count = sum(block.row_count for block in section_blocks[:split_index])
        right_count = sum(block.row_count for block in section_blocks[split_index:])
        if left_count <= max_split_table_rows and right_count <= max_split_table_rows:
            candidates.append((abs(left_count - right_count), max(left_count, right_count), split_index))

    if not candidates:
        raise ValueError(
            "按二级标题拆分后仍超出左右双表容量，请调大 max_split_table_rows 或调整模板布局"
        )

    _, _, split_index = min(candidates)
    return DetailLayout(
        is_split=True,
        left_blocks=tuple(section_blocks[:split_index]),
        right_blocks=tuple(section_blocks[split_index:]),
    )


def flatten_blocks(blocks: Sequence[SectionBlock]) -> list[tuple[str, float | None, float | None]]:
    rows: list[tuple[str, float | None, float | None]] = []
    for block in blocks:
        rows.extend(block.rows)
    return rows


def format_report_value(value: float | None, *, blank_display: str = "") -> str:
    if value is None:
        return blank_display
    formatted = format_value(value)
    return formatted if formatted != "" else blank_display


def load_env_file(path: Path) -> dict[str, str]:
    if not path.exists():
        raise FileNotFoundError(f".env 文件不存在: {path}")
    values: dict[str, str] = {}
    for raw_line in path.read_text(encoding="utf-8").splitlines():
        line = raw_line.strip()
        if not line or line.startswith("#"):
            continue
        if line.startswith("export "):
            line = line[7:].strip()
        if "=" not in line:
            continue
        key, value = line.split("=", 1)
        key = key.strip()
        value = value.strip()
        if value[:1] == value[-1:] and value[:1] in {"'", '"'}:
            value = value[1:-1]
        values[key] = value
    return values


def read_required_env_value(values: dict[str, str], *keys: str) -> str:
    for key in keys:
        value = values.get(key)
        if value:
            return value
    raise ValueError(f".env 缺少必要配置: {' / '.join(keys)}")


def prepare_llm_runtime(
    config: LlmNotesConfig,
    *,
    client_factory=OpenAI,
) -> LlmRuntimeConfig:
    env_values = load_env_file(config.env_path)
    api_key = read_required_env_value(env_values, "OPENAI_API_KEY", "LLM_API_KEY")
    model = read_required_env_value(env_values, "OPENAI_MODEL", "LLM_MODEL")
    base_url = env_values.get("OPENAI_BASE_URL") or env_values.get("LLM_BASE_URL")
    timeout_raw = env_values.get("OPENAI_TIMEOUT") or env_values.get("LLM_TIMEOUT")
    temperature_raw = env_values.get("OPENAI_TEMPERATURE") or env_values.get("LLM_TEMPERATURE")
    timeout = float(timeout_raw) if timeout_raw else None
    temperature = float(temperature_raw) if temperature_raw else config.temperature

    client_kwargs = {"api_key": api_key}
    if base_url:
        client_kwargs["base_url"] = base_url
    if timeout is not None:
        client_kwargs["timeout"] = timeout

    client = client_factory(**client_kwargs)
    system_role = config.system_role_path.read_text(encoding="utf-8").strip()
    if not system_role:
        raise ValueError(f"system role 文件为空: {config.system_role_path}")

    return LlmRuntimeConfig(
        client=client,
        model=model,
        system_role=system_role,
        target_chars=config.target_chars,
        temperature=temperature,
        max_tokens=config.max_tokens,
        checkpoint_chars=config.checkpoint_chars,
    )


def build_notes_prompt(
    *,
    title: str,
    report_rows: Sequence[tuple[str, float | None, float | None]],
    role_definition: RoleDefinition | None,
    target_chars: int,
) -> str:
    overall_label, overall_satisfaction, overall_importance = report_rows[0]
    section_names = (
        {section.name for section in role_definition.sections}
        if role_definition is not None
        else set()
    )
    table_lines = ["指标 | 满意度 | 重要性", "--- | --- | ---"]
    for label, satisfaction, importance in report_rows:
        if label != overall_label and (satisfaction is None or importance is None):
            continue
        row_type = "二级标题" if label in section_names else "指标"
        if label == overall_label:
            row_type = "总体"
        table_lines.append(
            f"{row_type}:{label} | {format_report_value(satisfaction)} | {format_report_value(importance)}"
        )

    min_chars = max(180, target_chars - 40)
    max_chars = target_chars + 40
    return (
        f"请基于以下客户满意度表格数据，撰写一段用于 PPT 备注页的中文分析描述。\n"
        f"要求：\n"
        f"1. 严格基于数据本身，不虚构原因，不编造样本量和同比环比。\n"
        f"2. 先概述总体满意度和重要性水平，再指出表现较好的部分与相对偏弱的部分。\n"
        f"3. 二级指标或三级指标若无有效分值，直接忽略，不要单独提及空值、未评价项或缺失项。\n"
        f"4. 语言面向管理层，简洁、正式、可直接用于备注页。\n"
        f"5. 只输出一段话，不要标题，不要项目符号。\n"
        f"6. 控制在约 {target_chars} 字，尽量落在 {min_chars}-{max_chars} 字之间。\n\n"
        f"页面标题：{title}\n"
        f"总体行：{overall_label}，满意度 {format_report_value(overall_satisfaction)}，重要性 {format_report_value(overall_importance)}\n"
        f"表格数据：\n" + "\n".join(table_lines)
    )


def extract_completion_text(response) -> str:
    choices = getattr(response, "choices", None) or []
    if not choices:
        return ""
    message = getattr(choices[0], "message", None)
    if message is None:
        return ""
    content = getattr(message, "content", "")
    if isinstance(content, str):
        return content.strip()
    if isinstance(content, list):
        parts: list[str] = []
        for item in content:
            text = getattr(item, "text", None)
            if text:
                parts.append(text)
        return "".join(parts).strip()
    return str(content).strip()


def extract_stream_chunk_text(chunk) -> str:
    choices = getattr(chunk, "choices", None) or []
    if not choices:
        return ""
    delta = getattr(choices[0], "delta", None)
    if delta is None:
        return ""
    content = getattr(delta, "content", None)
    if isinstance(content, str):
        return content
    if isinstance(content, list):
        parts: list[str] = []
        for item in content:
            text = getattr(item, "text", None)
            if text:
                parts.append(text)
        return "".join(parts)
    if content is None:
        return ""
    return str(content)


def generate_notes_text(
    *,
    title: str,
    report_rows: Sequence[tuple[str, float | None, float | None]],
    role_definition: RoleDefinition | None,
    runtime: LlmRuntimeConfig,
    on_text_update=None,
) -> str:
    prompt = build_notes_prompt(
        title=title,
        report_rows=report_rows,
        role_definition=role_definition,
        target_chars=runtime.target_chars,
    )
    messages = [
        {"role": "system", "content": runtime.system_role},
        {"role": "user", "content": prompt},
    ]

    stream = runtime.client.chat.completions.create(
        model=runtime.model,
        messages=messages,
        temperature=runtime.temperature,
        max_tokens=runtime.max_tokens,
        stream=True,
    )

    fragments: list[str] = []
    for chunk in stream:
        piece = extract_stream_chunk_text(chunk)
        if not piece:
            continue
        fragments.append(piece)
        current_text = "".join(fragments)
        if on_text_update is not None:
            on_text_update(current_text, False)
        print(piece, end="", flush=True)

    text = "".join(fragments).strip()
    if text:
        if on_text_update is not None:
            on_text_update(text, True)
        return text

    response = runtime.client.chat.completions.create(
        model=runtime.model,
        messages=messages,
        temperature=runtime.temperature,
        max_tokens=runtime.max_tokens,
    )
    text = extract_completion_text(response)
    if not text:
        raise ValueError(f"{title} 的备注页分析未返回有效文本")
    if on_text_update is not None:
        on_text_update(text, True)
    print(text, end="", flush=True)
    return text


def write_notes_text(slide, text: str) -> None:
    notes_slide = slide.notes_slide
    text_frame = notes_slide.notes_text_frame
    if text_frame is None:
        raise ValueError("备注页缺少可写入的文本框")
    text_frame.text = text


def ensure_parent_dir(path: Path) -> None:
    path.parent.mkdir(parents=True, exist_ok=True)


def save_presentation_checkpoint(presentation: Presentation, output_path: Path) -> None:
    ensure_parent_dir(output_path)
    presentation.save(str(output_path))


def remove_file_if_exists(path: Path | None) -> None:
    if path is None:
        return
    try:
        path.unlink()
    except FileNotFoundError:
        return


def find_title_shape(slide):
    if slide.shapes.title is not None:
        return slide.shapes.title
    for shape in slide.shapes:
        if getattr(shape, "is_placeholder", False):
            placeholder_type = str(shape.placeholder_format.type)
            if "TITLE" in placeholder_type:
                return shape
    raise ValueError("模板页缺少标题占位符")


def apply_title(slide, title: str) -> None:
    title_shape = find_title_shape(slide)
    title_shape.text = title
    if title_shape.has_text_frame:
        title_shape.text_frame.word_wrap = True


def remove_slide(presentation: Presentation, slide_index: int) -> None:
    slide_id_list = presentation.slides._sldIdLst
    slide_id = list(slide_id_list)[slide_index]
    presentation.part.drop_rel(slide_id.rId)
    slide_id_list.remove(slide_id)


def clear_all_slides(presentation: Presentation) -> None:
    for slide_index in range(len(presentation.slides) - 1, -1, -1):
        remove_slide(presentation, slide_index)


def get_blank_slide_layout(presentation: Presentation):
    for layout in presentation.slide_layouts:
        if layout.name in {"空白", "Blank"}:
            return layout
    return min(presentation.slide_layouts, key=lambda layout: len(layout.placeholders))


def find_slide_layout_by_name(presentation: Presentation, layout_name: str | None):
    if layout_name:
        for layout in presentation.slide_layouts:
            if layout.name == layout_name:
                return layout
    return None


def shape_has_external_relationship(shape) -> bool:
    xml = shape._element.xml
    return 'r:embed="' in xml or 'r:link="' in xml or 'r:id="' in xml


def assign_unique_shape_ids(shape_element, start_id: int) -> int:
    next_id = start_id
    for element in shape_element.iter():
        if element.tag == qn("p:cNvPr"):
            element.set("id", str(next_id))
            next_id += 1
    return next_id


def remove_slide_placeholders(slide) -> None:
    for shape in list(slide.shapes):
        if not getattr(shape, "is_placeholder", False):
            continue
        sp_tree = shape._element.getparent()
        if sp_tree is not None:
            sp_tree.remove(shape._element)


def copy_supported_shapes(source_slide, destination_slide) -> None:
    sp_tree = destination_slide.shapes._spTree
    next_shape_id = destination_slide.shapes._next_shape_id
    copied_count = 0

    for shape in source_slide.shapes:
        if shape_has_external_relationship(shape):
            continue
        new_element = deepcopy(shape._element)
        next_shape_id = assign_unique_shape_ids(new_element, next_shape_id)
        sp_tree.insert_element_before(new_element, "p:extLst")
        copied_count += 1

    if copied_count == 0:
        raise ValueError("章节页未找到可复制的可见形状")


def insert_category_intro_slide(
    presentation: Presentation,
    intro_config: CategoryIntroSlideConfig,
    *,
    source_presentations: dict[Path, Presentation],
) -> None:
    source_path = intro_config.ppt_path.resolve()
    if not source_path.exists():
        raise FileNotFoundError(f"章节页模板不存在: {source_path}")

    source_presentation = source_presentations.get(source_path)
    if source_presentation is None:
        source_presentation = Presentation(str(source_path))
        source_presentations[source_path] = source_presentation

    slide_index = intro_config.slide_number - 1
    if slide_index < 0 or slide_index >= len(source_presentation.slides):
        raise IndexError(
            f"{source_path} 不存在第 {intro_config.slide_number} 页章节页"
        )

    source_slide = source_presentation.slides[slide_index]
    destination_layout = find_slide_layout_by_name(
        presentation,
        source_slide.slide_layout.name if source_slide.slide_layout is not None else None,
    )
    if destination_layout is None:
        destination_layout = get_blank_slide_layout(presentation)

    destination_slide = presentation.slides.add_slide(destination_layout)
    remove_slide_placeholders(destination_slide)
    copy_supported_shapes(source_slide, destination_slide)


def generate_presentation(
    config: PptBatchConfig,
    *,
    dry_run: bool = False,
    llm_client_factory=OpenAI,
) -> Path:
    files = discover_input_files(config)
    presentation = Presentation(str(config.template_path))
    llm_runtime = None
    partial_output_path = None
    if config.llm_notes.enabled and not dry_run:
        llm_runtime = prepare_llm_runtime(config.llm_notes, client_factory=llm_client_factory)
        partial_output_path = build_partial_output_path(config.output_ppt)

    if config.template_slide_index >= len(presentation.slides):
        raise IndexError("template_slide_index 超出模板页数量")

    template_slide = presentation.slides[config.template_slide_index]
    template_layout = template_slide.slide_layout
    clear_all_slides(presentation)
    total_files = len(files)
    inserted_intro_categories: set[str] = set()
    source_presentations: dict[Path, Presentation] = {}

    try:
        for index, workbook_path in enumerate(files):
            workbook_meta = resolve_workbook_display_meta(workbook_path.stem)
            category_label = workbook_meta.category_label
            if category_label and category_label not in inserted_intro_categories:
                intro_config = config.category_intro_slides.get(category_label)
                if intro_config is not None:
                    insert_category_intro_slide(
                        presentation,
                        intro_config,
                        source_presentations=source_presentations,
                    )
                inserted_intro_categories.add(category_label)

            slide = presentation.slides.add_slide(template_layout)
            render_workbook_slide(
                slide,
                workbook_path,
                config,
                presentation=presentation,
                checkpoint_output_path=partial_output_path,
                llm_runtime=llm_runtime,
                slide_index=index + 1,
                total_slides=total_files,
            )
            if partial_output_path is not None:
                save_presentation_checkpoint(presentation, partial_output_path)
                print(
                    f"[{index + 1}/{total_files}] 已保存检查点：{partial_output_path.name}",
                    flush=True,
                )

        if dry_run:
            return config.output_ppt

        ensure_parent_dir(config.output_ppt)
        presentation.save(str(config.output_ppt))
        remove_file_if_exists(partial_output_path)
        return config.output_ppt
    except BaseException:
        if not dry_run and partial_output_path is not None:
            try:
                save_presentation_checkpoint(presentation, partial_output_path)
                print(
                    f"生成中断，已保存当前检查点：{partial_output_path}",
                    flush=True,
                )
            except Exception as checkpoint_error:
                print(f"生成中断，保存检查点失败：{checkpoint_error}", flush=True)
        raise


def render_workbook_slide(
    slide,
    workbook_path: Path,
    config: PptBatchConfig,
    *,
    presentation: Presentation | None = None,
    checkpoint_output_path: Path | None = None,
    llm_runtime: LlmRuntimeConfig | None = None,
    slide_index: int | None = None,
    total_slides: int | None = None,
) -> None:
    title = resolve_workbook_display_meta(workbook_path.stem).title + config.title_suffix
    apply_title(slide, title)

    report_rows = read_report_rows(
        workbook_path,
        sheet_name_mode=config.sheet_name_mode,
        sheet_name=config.sheet_name,
    )
    overall_row = report_rows[0]
    detail_rows = report_rows[1:]
    role_definition = resolve_section_definition(
        workbook_path.stem,
        report_rows,
        section_mode=config.section_mode,
    )
    detail_rows = filter_empty_satisfaction_sections(detail_rows, role_definition)

    render_table(
        slide,
        config.layout.summary_table,
        [overall_row],
        blank_display=config.blank_display,
        section_names=set(),
        overall_label=overall_row[0],
        header_font_size_pt=config.header_font_size_pt,
        body_font_size_pt=config.summary_font_size_pt,
    )

    if detail_rows:
        detail_layout = choose_detail_layout(
            detail_rows=detail_rows,
            role_definition=role_definition,
            max_single_table_rows=config.max_single_table_rows,
            max_split_table_rows=config.max_split_table_rows,
        )

        if detail_layout.is_split:
            left_section_names = {block.heading for block in detail_layout.left_blocks}
            right_section_names = {block.heading for block in detail_layout.right_blocks}
            render_table(
                slide,
                config.layout.detail_left_table,
                flatten_blocks(detail_layout.left_blocks),
                blank_display=config.blank_display,
                section_names=left_section_names,
                overall_label=None,
                header_font_size_pt=config.header_font_size_pt,
                body_font_size_pt=config.body_font_size_pt,
            )
            render_table(
                slide,
                config.layout.detail_right_table,
                flatten_blocks(detail_layout.right_blocks),
                blank_display=config.blank_display,
                section_names=right_section_names,
                overall_label=None,
                header_font_size_pt=config.header_font_size_pt,
                body_font_size_pt=config.body_font_size_pt,
            )
        else:
            section_names = (
                {section.name for section in role_definition.sections} if role_definition else set()
            )
            render_table(
                slide,
                config.layout.detail_single_table,
                list(detail_layout.single_rows),
                blank_display=config.blank_display,
                section_names=section_names,
                overall_label=None,
                header_font_size_pt=config.header_font_size_pt,
                body_font_size_pt=config.body_font_size_pt,
            )

    if llm_runtime is not None:
        progress_prefix = ""
        if slide_index is not None and total_slides is not None:
            progress_prefix = f"[{slide_index}/{total_slides}] "
        print(f"{progress_prefix}正在生成备注页分析：{title}", flush=True)
        print(f"{progress_prefix}流式输出：", end="", flush=True)
        checkpoint_state = {"last_saved_length": 0}

        def handle_notes_update(current_text: str, is_final: bool) -> None:
            write_notes_text(slide, current_text)
            if (
                presentation is not None
                and checkpoint_output_path is not None
                and (
                    is_final
                    or len(current_text) - checkpoint_state["last_saved_length"]
                    >= llm_runtime.checkpoint_chars
                )
            ):
                save_presentation_checkpoint(presentation, checkpoint_output_path)
                checkpoint_state["last_saved_length"] = len(current_text)

        notes_text = generate_notes_text(
            title=title,
            report_rows=[overall_row, *detail_rows],
            role_definition=role_definition,
            runtime=llm_runtime,
            on_text_update=handle_notes_update,
        )
        print("", flush=True)
        write_notes_text(slide, notes_text)
        print(
            f"{progress_prefix}备注页分析完成：{title}（{len(notes_text)}字）",
            flush=True,
        )


def render_table(
    slide,
    region: TableRegion,
    rows: Sequence[tuple[str, float | None, float | None]],
    *,
    blank_display: str,
    section_names: set[str],
    overall_label: str | None,
    header_font_size_pt: float,
    body_font_size_pt: float,
) -> None:
    left, top, width, height = region.emu()
    shape = slide.shapes.add_table(len(rows) + 1, 3, left, top, width, height)
    table = shape.table
    set_column_widths(table, width)

    headers = ("指标", "满意度", "重要性")
    for column_index, header in enumerate(headers):
        cell = table.cell(0, column_index)
        set_cell_text(
            cell,
            header,
            bold=True,
            align=PP_ALIGN.CENTER,
            font_size_pt=header_font_size_pt,
            fill_color=HEADER_FILL_COLOR,
            text_color=HEADER_TEXT_COLOR,
        )

    for row_index, row in enumerate(rows, start=1):
        label, satisfaction, importance = row
        if overall_label and label == overall_label:
            fill_color = OVERALL_FILL_COLOR
            bold = True
        elif label in section_names:
            fill_color = SECTION_FILL_COLOR
            bold = True
        else:
            fill_color = BODY_FILL_COLOR
            bold = False

        values = (
            label,
            format_report_value(satisfaction, blank_display=blank_display),
            format_report_value(importance, blank_display=blank_display),
        )
        for column_index, value in enumerate(values):
            cell = table.cell(row_index, column_index)
            set_cell_text(
                cell,
                value,
                bold=bold,
                align=PP_ALIGN.CENTER,
                font_size_pt=body_font_size_pt,
                fill_color=fill_color,
                text_color=BODY_TEXT_COLOR,
            )


def set_column_widths(table, total_width: int) -> None:
    first_col = int(total_width * 0.62)
    second_col = int(total_width * 0.19)
    third_col = total_width - first_col - second_col
    table.columns[0].width = first_col
    table.columns[1].width = second_col
    table.columns[2].width = third_col


def set_cell_text(
    cell,
    text: str,
    *,
    bold: bool,
    align,
    font_size_pt: float,
    fill_color: str,
    text_color: str,
) -> None:
    cell.fill.solid()
    cell.fill.fore_color.rgb = RGBColor.from_string(fill_color)
    set_cell_border(cell, BORDER_COLOR)
    cell.vertical_anchor = MSO_ANCHOR.MIDDLE
    cell.margin_left = Pt(4)
    cell.margin_right = Pt(4)
    cell.margin_top = Pt(2)
    cell.margin_bottom = Pt(2)

    text_frame = cell.text_frame
    text_frame.clear()
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = align
    run = paragraph.add_run()
    run.text = text
    run.font.bold = bold
    run.font.size = Pt(font_size_pt)
    run.font.color.rgb = RGBColor.from_string(text_color)


def set_cell_border(cell, color: str, width: int = 12700) -> None:
    tc = cell._tc
    tc_pr = tc.get_or_add_tcPr()

    for edge in ("a:lnL", "a:lnR", "a:lnT", "a:lnB"):
        line = tc_pr.find(qn(edge))
        if line is None:
            line = OxmlElement(edge)
            tc_pr.append(line)

        line.set("w", str(width))
        line.set("cap", "flat")
        line.set("cmpd", "sng")
        line.set("algn", "ctr")

        solid_fill = line.find(qn("a:solidFill"))
        if solid_fill is None:
            solid_fill = OxmlElement("a:solidFill")
            line.append(solid_fill)

        srgb = solid_fill.find(qn("a:srgbClr"))
        if srgb is None:
            srgb = OxmlElement("a:srgbClr")
            solid_fill.append(srgb)
        srgb.set("val", color)

        prst_dash = line.find(qn("a:prstDash"))
        if prst_dash is None:
            prst_dash = OxmlElement("a:prstDash")
            line.append(prst_dash)
        prst_dash.set("val", "solid")

        round_node = line.find(qn("a:round"))
        if round_node is None:
            round_node = OxmlElement("a:round")
            line.append(round_node)

        head_end = line.find(qn("a:headEnd"))
        if head_end is None:
            head_end = OxmlElement("a:headEnd")
            line.append(head_end)
        head_end.set("type", "none")
        head_end.set("w", "med")
        head_end.set("len", "med")

        tail_end = line.find(qn("a:tailEnd"))
        if tail_end is None:
            tail_end = OxmlElement("a:tailEnd")
            line.append(tail_end)
        tail_end.set("type", "none")
        tail_end.set("w", "med")
        tail_end.set("len", "med")


def build_default_config_from_args(args: argparse.Namespace) -> PptBatchConfig:
    if not args.template_path or not args.input_dir or not args.output_ppt:
        raise ValueError("未提供 --config 时，必须同时提供 --template-path / --input-dir / --output-ppt")
    return PptBatchConfig(
        template_path=args.template_path,
        input_dir=args.input_dir,
        output_ppt=args.output_ppt,
        blank_display=args.blank_display or "",
        section_mode=normalize_section_mode(args.section_mode),
    )


def main() -> None:
    parser = build_parser()
    args = parser.parse_args()

    if args.config:
        config = load_batch_config(
            args.config,
            template_path=args.template_path,
            input_dir=args.input_dir,
            output_ppt=args.output_ppt,
            section_mode=args.section_mode,
            blank_display=args.blank_display,
        )
    else:
        config = build_default_config_from_args(args)

    output_path = generate_presentation(config, dry_run=args.dry_run)
    if args.dry_run:
        print(f"校验完成，未写出文件: {output_path}")
    else:
        print(f"PPT 已生成: {output_path}")


if __name__ == "__main__":
    main()
