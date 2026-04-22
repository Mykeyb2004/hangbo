from __future__ import annotations

import io
import tempfile
import unittest
from contextlib import redirect_stdout
from pathlib import Path
from unittest import mock

import generate_ppt as generate_ppt_module
from openpyxl import Workbook
from pptx.enum.shapes import MSO_SHAPE_TYPE
from pptx import Presentation
from pptx.oxml.ns import qn
from pptx.util import Pt

from generate_ppt import (
    BODY_FILL_COLOR,
    BODY_TEXT_COLOR,
    BORDER_COLOR,
    build_notes_prompt,
    build_llm_failure_log_path,
    CategoryIntroSlideConfig,
    ChartPageConfig,
    CHART_TEXTBOX_FIRST_LINE_INDENT_PT,
    CHART_TEXTBOX_FONT_SIZE_PT,
    CHART_TEXTBOX_FONT_NAME,
    CHART_TEXTBOX_LINE_SPACING,
    HEADER_FILL_COLOR,
    HEADER_TEXT_COLOR,
    LlmNotesConfig,
    OVERALL_FILL_COLOR,
    PptBatchConfig,
    PptLayoutConfig,
    SECTION_FILL_COLOR,
    TableRegion,
    build_partial_output_path,
    build_section_blocks,
    choose_detail_layout,
    discover_input_files,
    format_report_value,
    generate_presentation,
    load_batch_config,
    render_table,
    render_chart_textbox,
    resolve_chart_textbox_style,
    resolve_section_definition,
    resolve_workbook_display_meta,
)
from survey_customer_category_rules import CustomerCategoryRule


def create_report_workbook(path: Path, rows: list[tuple[object, object, object]]) -> None:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = path.stem
    for row in rows:
        worksheet.append(list(row))
    workbook.save(path)


def collect_slide_texts(slide) -> list[str]:
    texts: list[str] = []
    for shape in slide.shapes:
        if getattr(shape, "has_text_frame", False):
            for paragraph in shape.text_frame.paragraphs:
                text = paragraph.text.strip()
                if text:
                    texts.append(text)
        if hasattr(shape, "shapes"):
            for nested_shape in shape.shapes:
                if not getattr(nested_shape, "has_text_frame", False):
                    continue
                for paragraph in nested_shape.text_frame.paragraphs:
                    text = paragraph.text.strip()
                    if text:
                        texts.append(text)
    return texts


class FakeMessage:
    def __init__(self, content: str) -> None:
        self.content = content


class FakeChoice:
    def __init__(self, content: str) -> None:
        self.message = FakeMessage(content)


class FakeCompletionResponse:
    def __init__(self, content: str) -> None:
        self.choices = [FakeChoice(content)]


class FakeDelta:
    def __init__(self, content: str | None) -> None:
        self.content = content


class FakeStreamChoice:
    def __init__(self, content: str | None) -> None:
        self.delta = FakeDelta(content)


class FakeStreamChunk:
    def __init__(self, content: str | None) -> None:
        self.choices = [FakeStreamChoice(content)]


class FakeChatCompletions:
    RESPONSE_TEXT = (
        "总体判断：本页客户体验主要由会展服务支撑，配套与动线体验形成拖累。\n"
        "亮点：会展服务表现突出，接待引导服务保持满分。\n"
        "关注点：园区停车方便和展会路线安排偏弱，且路线安排与重要性存在明显差距。"
    )

    def __init__(self, outer) -> None:
        self.outer = outer

    def create(self, **kwargs):
        self.outer.create_calls.append(kwargs)
        if kwargs.get("stream"):
            return [
                FakeStreamChunk("总体判断：本页客户体验主要由会展服务支撑，"),
                FakeStreamChunk("配套与动线体验形成拖累。\n"),
                FakeStreamChunk("亮点：会展服务表现突出，"),
                FakeStreamChunk("接待引导服务保持满分。\n"),
                FakeStreamChunk("关注点：园区停车方便和展会路线安排偏弱，"),
                FakeStreamChunk("且路线安排与重要性存在明显差距。"),
                FakeStreamChunk(None),
            ]
        return FakeCompletionResponse(self.RESPONSE_TEXT)


class FakeChat:
    def __init__(self, outer) -> None:
        self.completions = FakeChatCompletions(outer)


class FakeOpenAI:
    instances: list["FakeOpenAI"] = []

    def __init__(self, **kwargs) -> None:
        self.kwargs = kwargs
        self.create_calls: list[dict[str, object]] = []
        self.chat = FakeChat(self)
        self.__class__.instances.append(self)


class FakeInterruptedChatCompletions:
    FIRST_RESPONSE_TEXT = (
        "第一页分析已经完成，整体表现较稳，重点服务指标保持在较高水平，"
        "建议继续关注优势环节的稳定性，并结合后续数据观察波动。"
    )

    def __init__(self, outer) -> None:
        self.outer = outer
        self.stream_calls = 0

    def create(self, **kwargs):
        self.outer.create_calls.append(kwargs)
        if not kwargs.get("stream"):
            return FakeCompletionResponse(self.FIRST_RESPONSE_TEXT)

        self.stream_calls += 1
        if self.stream_calls == 1:
            return [
                FakeStreamChunk("第一页分析已经完成，整体表现较稳，"),
                FakeStreamChunk("重点服务指标保持在较高水平，"),
                FakeStreamChunk("建议继续关注优势环节的稳定性，并结合后续数据观察波动。"),
                FakeStreamChunk(None),
            ]

        def interrupted_stream():
            yield FakeStreamChunk("第二页分析进行到一半，")
            raise KeyboardInterrupt()

        return interrupted_stream()


class FakeInterruptedChat:
    def __init__(self, outer) -> None:
        self.completions = FakeInterruptedChatCompletions(outer)


class FakeInterruptedOpenAI:
    instances: list["FakeInterruptedOpenAI"] = []

    def __init__(self, **kwargs) -> None:
        self.kwargs = kwargs
        self.create_calls: list[dict[str, object]] = []
        self.chat = FakeInterruptedChat(self)
        self.__class__.instances.append(self)


class FakeFlakyChatCompletions:
    RESPONSE_TEXT = (
        "总体判断：专业观众整体体验由会展服务支撑，硬件与配套服务形成拖累。\n"
        "亮点：接待引导服务和工作人员仪容仪表保持高分。\n"
        "关注点：园区停车方便和餐饮服务偏弱，需优先补强。"
    )

    def __init__(self, outer) -> None:
        self.outer = outer
        self.stream_calls = 0

    def create(self, **kwargs):
        self.outer.create_calls.append(kwargs)
        if not kwargs.get("stream"):
            return FakeCompletionResponse(self.RESPONSE_TEXT)

        self.stream_calls += 1
        if self.stream_calls < 3:
            raise RuntimeError("transient llm failure")

        return [
            FakeStreamChunk("总体判断：专业观众整体体验由会展服务支撑，"),
            FakeStreamChunk("硬件与配套服务形成拖累。\n"),
            FakeStreamChunk("亮点：接待引导服务和工作人员仪容仪表保持高分。\n"),
            FakeStreamChunk("关注点：园区停车方便和餐饮服务偏弱，需优先补强。"),
            FakeStreamChunk(None),
        ]


class FakeFlakyChat:
    def __init__(self, outer) -> None:
        self.completions = FakeFlakyChatCompletions(outer)


class FakeFlakyOpenAI:
    instances: list["FakeFlakyOpenAI"] = []

    def __init__(self, **kwargs) -> None:
        self.kwargs = kwargs
        self.create_calls: list[dict[str, object]] = []
        self.chat = FakeFlakyChat(self)
        self.__class__.instances.append(self)


class FakeAlwaysFailChatCompletions:
    def __init__(self, outer) -> None:
        self.outer = outer

    def create(self, **kwargs):
        self.outer.create_calls.append(kwargs)
        raise RuntimeError("upstream llm unavailable")


class FakeAlwaysFailChat:
    def __init__(self, outer) -> None:
        self.completions = FakeAlwaysFailChatCompletions(outer)


class FakeAlwaysFailOpenAI:
    instances: list["FakeAlwaysFailOpenAI"] = []

    def __init__(self, **kwargs) -> None:
        self.kwargs = kwargs
        self.create_calls: list[dict[str, object]] = []
        self.chat = FakeAlwaysFailChat(self)
        self.__class__.instances.append(self)


class GeneratePptTest(unittest.TestCase):
    def setUp(self) -> None:
        FakeOpenAI.instances.clear()
        FakeInterruptedOpenAI.instances.clear()
        FakeFlakyOpenAI.instances.clear()
        FakeAlwaysFailOpenAI.instances.clear()

    def test_system_role_file_no_longer_requires_markdown_hierarchy(self) -> None:
        repo_root = Path(__file__).resolve().parents[1]
        system_role_text = (repo_root / "system_role.md").read_text(encoding="utf-8")

        self.assertNotIn("注意层级用markdown标注", system_role_text)
        self.assertNotIn("第一层级到第四层级分别为：一、（一）1. （1）", system_role_text)

    def test_system_role_guides_ppt_short_comments_without_report_template(self) -> None:
        repo_root = Path(__file__).resolve().parents[1]
        system_role_text = (repo_root / "system_role.md").read_text(encoding="utf-8")

        self.assertIn("PPT备注页、图表页说明或管理短评", system_role_text)
        self.assertIn("全表数据是判断依据，不是覆盖清单", system_role_text)
        self.assertIn("避免在多页报告中反复使用相同开头、连接句和收尾句", system_role_text)
        self.assertNotIn("当用户要求分析、总结或报告类内容时，优先按以下逻辑组织", system_role_text)
        self.assertNotIn("执行摘要", system_role_text)
        self.assertIn("总体判断：", system_role_text)
        self.assertIn("关注点：", system_role_text)

    def test_build_notes_prompt_limits_gap_analysis_to_secondary_and_tertiary_metrics(self) -> None:
        rows = [
            ("专业观众", 9.93, 10.0),
            ("会展服务", 9.86, 9.90),
            ("工作人员仪容仪表", 10.0, 10.0),
            ("硬件设施", 9.31, 9.89),
            ("展会路线安排", 9.20, 9.80),
        ]

        prompt = build_notes_prompt(
            title="会展客户——专业观众",
            report_rows=rows,
            role_definition=resolve_section_definition("专业观众", rows),
            target_chars=300,
        )

        self.assertIn(
            "如需进行满意度与重要性差异分析，仅针对二级、三级指标，不对总体行或页面标题做此类分析。",
            prompt,
        )
        self.assertNotIn("可以做适当的进行满意度与重要性差异分析。", prompt)
        self.assertIn(
            "目标是让管理层在较短时间内抓住重点，而不是完整复述表格。",
            prompt,
        )
        self.assertIn("严格按以下结构逐行输出", prompt)
        self.assertIn("总体判断：", prompt)
        self.assertIn("关注点：", prompt)
        self.assertIn("不逐项复述表格，不追求覆盖全面", prompt)
        self.assertIn("避免多页报告中反复使用相同开头、连接句和收尾句", prompt)
        self.assertIn("不要使用“本页”", prompt)
        self.assertIn("不要使用“整体评价较高”", prompt)
        self.assertIn("总体:专业观众 | 9.93 | 10", prompt)
        self.assertIn("二级标题:会展服务 | 9.86 | 9.9", prompt)
        self.assertIn("指标:工作人员仪容仪表 | 10 | 10", prompt)

    def test_project_pipeline_defaults_use_balanced_ppt_llm_notes_settings(self) -> None:
        from pipeline_config import load_pipeline_defaults

        repo_root = Path(__file__).resolve().parents[1]
        defaults = load_pipeline_defaults(repo_root / "pipeline.defaults.toml")

        self.assertEqual(defaults.ppt.llm_notes.target_chars, 120)
        self.assertEqual(defaults.ppt.llm_notes.temperature, 0.4)
        self.assertEqual(defaults.ppt.llm_notes.max_tokens, 200)

    def test_build_section_blocks_groups_rows_by_second_level_titles(self) -> None:
        rows = [
            ("专业观众", 9.93, 10.0),
            ("会展服务", 10.0, 10.0),
            ("工作人员仪容仪表", 10.0, 10.0),
            ("工作人员服务态度", 10.0, 10.0),
            ("硬件设施", 9.91, 10.0),
            ("展会路线安排", 9.93, 10.0),
        ]

        role_definition = resolve_section_definition("专业观众", rows)
        section_blocks = build_section_blocks(rows, role_definition)

        self.assertEqual([block.heading for block in section_blocks], ["会展服务", "硬件设施"])
        self.assertEqual([len(block.rows) for block in section_blocks], [3, 2])

    def test_choose_detail_layout_splits_into_two_tables_without_breaking_sections(self) -> None:
        rows = [
            ("专业观众", 9.93, 10.0),
            ("会展服务", 10.0, 10.0),
            ("工作人员仪容仪表", 10.0, 10.0),
            ("工作人员服务态度", 10.0, 10.0),
            ("工作人员业务技能", 10.0, 10.0),
            ("接待引导服务", 10.0, 10.0),
            ("硬件设施", 9.91, 10.0),
            ("展会路线安排", 9.93, 10.0),
            ("园区停车方便", 10.0, 10.0),
            ("交通便利，容易到达", 10.0, 10.0),
            ("标识标牌清晰", 9.73, 10.0),
            ("设施设备齐全", 9.89, 10.0),
            ("展厅使用情况", 9.84, 10.0),
            ("参展环境", 10.0, 10.0),
            ("配套服务", 9.8, 10.0),
            ("餐饮服务", 9.4, 10.0),
            ("客房服务", None, None),
            ("安保服务", 10.0, 10.0),
            ("保洁服务", 10.0, 10.0),
            ("智慧场馆", 10.0, 10.0),
            ("杭州国博APP", 10.0, 10.0),
            ("室内导航系统", None, None),
            ("寻车系统", None, None),
            ("云上看馆", None, None),
        ]

        role_definition = resolve_section_definition("专业观众", rows)
        detail_layout = choose_detail_layout(
            detail_rows=rows[1:],
            role_definition=role_definition,
            max_single_table_rows=18,
            max_split_table_rows=19,
        )

        self.assertTrue(detail_layout.is_split)
        self.assertEqual(
            [block.heading for block in detail_layout.left_blocks],
            ["会展服务", "硬件设施"],
        )
        self.assertEqual(
            [block.heading for block in detail_layout.right_blocks],
            ["配套服务", "智慧场馆"],
        )

    def test_format_report_value_hides_empty_values(self) -> None:
        self.assertEqual(format_report_value(None, blank_display=""), "")
        self.assertEqual(format_report_value(9.50, blank_display=""), "9.5")

    def test_discover_input_files_uses_customer_group_order(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            input_dir = temp_path / "input"
            input_dir.mkdir()

            for file_name in (
                "自助餐.xlsx",
                "未知客户.xlsx",
                "会议主承办.xlsx",
                "展览主承办.xlsx",
                "参展商.xlsx",
            ):
                create_report_workbook(
                    input_dir / file_name,
                    [
                        ("指标", "满意度", "重要性"),
                        (Path(file_name).stem, 9.9, 9.8),
                    ],
                )

            config = PptBatchConfig(
                template_path=Path("templates/template.pptx"),
                input_dir=input_dir,
                output_ppt=temp_path / "report.pptx",
            )

            files = discover_input_files(config)

            self.assertEqual(
                [path.name for path in files],
                [
                    "展览主承办.xlsx",
                    "参展商.xlsx",
                    "会议主承办.xlsx",
                    "自助餐.xlsx",
                    "未知客户.xlsx",
                ],
            )

    def test_discover_input_files_follows_display_rules_single_source(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            input_dir = temp_path / "input"
            input_dir.mkdir()

            for file_name in ("自助餐.xlsx", "展览主承办.xlsx"):
                create_report_workbook(
                    input_dir / file_name,
                    [
                        ("指标", "满意度", "重要性"),
                        (Path(file_name).stem, 9.9, 9.8),
                    ],
                )

            config = PptBatchConfig(
                template_path=Path("templates/template.pptx"),
                input_dir=input_dir,
                output_ppt=temp_path / "report.pptx",
            )
            reversed_rules = (
                CustomerCategoryRule(
                    name="自助餐",
                    customer_group="二、餐饮客户",
                    customer_category="自助餐",
                    source_file_name="餐饮.xlsx",
                    sequence_number=1,
                ),
                CustomerCategoryRule(
                    name="展览主承办",
                    customer_group="一、会展客户",
                    customer_category="展览活动主（承）办",
                    source_file_name="展览.xlsx",
                    sequence_number=2,
                ),
            )

            with mock.patch.object(
                generate_ppt_module,
                "DISPLAY_ORDERED_CUSTOMER_CATEGORY_RULES",
                reversed_rules,
            ):
                generate_ppt_module.build_workbook_display_lookup.cache_clear()
                files = discover_input_files(config)

            generate_ppt_module.build_workbook_display_lookup.cache_clear()
            self.assertEqual(
                [path.name for path in files],
                ["自助餐.xlsx", "展览主承办.xlsx"],
            )

    def test_resolve_workbook_display_meta_forces_shared_display_name(self) -> None:
        banquet_meta = resolve_workbook_display_meta("酒店宴会")
        buffet_meta = resolve_workbook_display_meta("酒店自助餐")

        self.assertEqual(banquet_meta.title, "酒店客户——酒店餐饮客户")
        self.assertEqual(buffet_meta.title, "酒店客户——酒店餐饮客户")

    def test_resolve_section_definition_supports_aggregate_hotel_catering_workbook(self) -> None:
        rows = [
            ("酒店餐饮客户", 9.7, 9.5),
            ("餐饮服务", 9.9, 9.8),
            ("菜品温度", 9.9, 9.8),
            ("菜肴品种", 9.8, 9.7),
            ("硬件设施", 9.5, 9.4),
            ("园区停车方便", 9.5, 9.4),
            ("交通便利，容易到达", 9.4, 9.3),
            ("智慧场馆", 9.2, 9.1),
            ("杭州国博APP", 9.2, 9.1),
        ]

        role_definition = resolve_section_definition("酒店餐饮客户", rows)

        self.assertIsNotNone(role_definition)
        self.assertEqual(
            [section.name for section in role_definition.sections],
            ["餐饮服务", "硬件设施", "智慧场馆"],
        )

    def test_load_batch_config_parses_category_intro_slides(self) -> None:
        repo_root = Path(__file__).resolve().parents[1]
        template_path = repo_root / "templates" / "template.pptx"
        chapter_template_path = repo_root / "templates" / "chapter.pptx"
        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            config_path = temp_path / "ppt-config.toml"
            config_path.write_text(
                "\n".join(
                    [
                        f'template_path = "{template_path}"',
                        'input_dir = "input"',
                        'output_ppt = "output/report.pptx"',
                        "",
                        '[category_intro_slides."一、会展客户"]',
                        f'ppt_path = "{chapter_template_path}"',
                        "slide_number = 3",
                        "",
                        '[category_intro_slides."五、酒店客户"]',
                        f'ppt_path = "{chapter_template_path}"',
                        "slide_number = 5",
                        "",
                        "[chart_page]",
                        "enabled = true",
                        'placeholder_text = "图表解读待补充"',
                        "image_dpi = 180",
                    ]
                ),
                encoding="utf-8",
            )

            config = load_batch_config(config_path)

            self.assertEqual(
                config.category_intro_slides["一、会展客户"].ppt_path,
                chapter_template_path,
            )
            self.assertEqual(config.category_intro_slides["一、会展客户"].slide_number, 3)
            self.assertEqual(config.category_intro_slides["五、酒店客户"].slide_number, 5)
            self.assertTrue(config.chart_page.enabled)
            self.assertEqual(config.chart_page.placeholder_text, "图表解读待补充")
            self.assertEqual(config.chart_page.image_dpi, 180)

    def test_ppt_batch_config_defaults_use_14pt_table_fonts(self) -> None:
        config = PptBatchConfig(
            template_path=Path("templates/template.pptx"),
            input_dir=Path("input"),
            output_ppt=Path("output/report.pptx"),
        )

        self.assertEqual(config.body_font_size_pt, 14.0)
        self.assertEqual(config.header_font_size_pt, 14.0)
        self.assertEqual(config.summary_font_size_pt, 14.0)
        self.assertEqual(config.layout.summary_table.height, 0.62)
        self.assertEqual(config.layout.detail_single_table.height, 5.25)
        self.assertEqual(config.layout.detail_left_table.height, 5.25)
        self.assertEqual(config.layout.detail_right_table.height, 5.25)

    def test_render_table_uses_kaiti_for_text_and_times_for_numbers(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])

        render_table(
            slide,
            TableRegion(0.73, 1.45, 11.87, 0.62),
            [("专业观众", 9.93, 10.0)],
            blank_display="",
            section_names=set(),
            overall_label=None,
            header_font_size_pt=14.0,
            body_font_size_pt=14.0,
        )

        table = next(shape.table for shape in slide.shapes if getattr(shape, "has_table", False))
        header_run = table.cell(0, 0).text_frame.paragraphs[0].runs[0]
        label_run = table.cell(1, 0).text_frame.paragraphs[0].runs[0]
        satisfaction_run = table.cell(1, 1).text_frame.paragraphs[0].runs[0]
        importance_run = table.cell(1, 2).text_frame.paragraphs[0].runs[0]

        self.assertEqual(header_run.font.size, Pt(14.0))
        self.assertEqual(label_run.font.size, Pt(14.0))
        self.assertEqual(satisfaction_run.font.size, Pt(14.0))
        self.assertEqual(importance_run.font.size, Pt(14.0))
        self.assertEqual(header_run.font.name, "楷体")
        self.assertEqual(label_run.font.name, "楷体")
        self.assertEqual(satisfaction_run.font.name, "Times New Roman")
        self.assertEqual(importance_run.font.name, "Times New Roman")

    def test_generate_presentation_creates_single_and_double_table_slides(self) -> None:
        repo_root = Path(__file__).resolve().parents[1]
        template_path = repo_root / "templates" / "template.pptx"

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            input_dir = temp_path / "input"
            input_dir.mkdir()
            output_path = temp_path / "report.pptx"

            create_report_workbook(
                input_dir / "专业观众.xlsx",
                [
                    ("指标", "满意度", "重要性"),
                    ("专业观众", 9.93, 10.0),
                    ("会展服务", 10.0, 10.0),
                    ("工作人员仪容仪表", 10.0, 10.0),
                    ("工作人员服务态度", 10.0, 10.0),
                    ("工作人员业务技能", 10.0, 10.0),
                    ("接待引导服务", 10.0, 10.0),
                    ("硬件设施", 9.91, 10.0),
                    ("展会路线安排", 9.93, 10.0),
                    ("园区停车方便", 10.0, 10.0),
                    ("交通便利，容易到达", 10.0, 10.0),
                    ("标识标牌清晰", 9.73, 10.0),
                    ("设施设备齐全", 9.89, 10.0),
                    ("展厅使用情况", 9.84, 10.0),
                    ("参展环境", 10.0, 10.0),
                    ("配套服务", 9.8, 10.0),
                    ("餐饮服务", 9.4, 10.0),
                    ("客房服务", None, None),
                    ("安保服务", 10.0, 10.0),
                    ("保洁服务", 10.0, 10.0),
                    ("智慧场馆", 10.0, 10.0),
                    ("杭州国博APP", 10.0, 10.0),
                    ("室内导航系统", None, None),
                    ("寻车系统", None, None),
                    ("云上看馆", None, None),
                ],
            )
            create_report_workbook(
                input_dir / "自助餐.xlsx",
                [
                    ("指标", "满意度", "重要性"),
                    ("自助餐", 9.80, 9.60),
                    ("餐饮服务", 9.85, 9.72),
                    ("菜品口味", 9.90, 9.60),
                    ("菜品丰富度", 9.80, 9.60),
                    ("补菜及时性", None, None),
                    ("硬件设施", 9.70, 9.40),
                    ("环境卫生", 9.70, 9.40),
                    ("桌椅舒适度", 9.70, 9.40),
                ],
            )

            config = PptBatchConfig(
                template_path=template_path,
                input_dir=input_dir,
                output_ppt=output_path,
                blank_display="",
                max_single_table_rows=10,
                max_split_table_rows=19,
                layout=PptLayoutConfig(),
            )

            generate_presentation(config)

            self.assertTrue(output_path.exists())

            presentation = Presentation(output_path)
            self.assertEqual(len(presentation.slides), 2)

            slide_tables = {}
            for slide in presentation.slides:
                title = slide.shapes.title.text
                table_count = sum(1 for shape in slide.shapes if getattr(shape, "has_table", False))
                table_texts = [
                    "\n".join(
                        cell.text
                        for row in shape.table.rows
                        for cell in row.cells
                    )
                    for shape in slide.shapes
                    if getattr(shape, "has_table", False)
                ]
                slide_tables[title] = (table_count, table_texts)

            professional_audience_title = "会展客户——专业观众"
            buffet_title = "餐饮客户——自助餐"
            self.assertEqual(slide_tables[professional_audience_title][0], 3)
            self.assertEqual(slide_tables[buffet_title][0], 2)
            self.assertTrue(any("会展服务" in text for text in slide_tables[professional_audience_title][1]))
            self.assertTrue(any("智慧场馆" in text for text in slide_tables[professional_audience_title][1]))
            self.assertTrue(any("补菜及时性" in text for text in slide_tables[buffet_title][1]))
            self.assertTrue(
                any("\n\n" in text or text.endswith("\n") for text in slide_tables[professional_audience_title][1])
            )

            first_slide_tables = [
                shape.table
                for shape in presentation.slides[0].shapes
                if getattr(shape, "has_table", False)
            ]
            summary_table = first_slide_tables[0]
            detail_table = first_slide_tables[1]

            self.assertEqual(str(summary_table.cell(0, 0).fill.fore_color.rgb), HEADER_FILL_COLOR)
            self.assertEqual(str(summary_table.cell(1, 0).fill.fore_color.rgb), OVERALL_FILL_COLOR)
            self.assertEqual(str(detail_table.cell(1, 0).fill.fore_color.rgb), SECTION_FILL_COLOR)
            self.assertEqual(str(detail_table.cell(2, 0).fill.fore_color.rgb), BODY_FILL_COLOR)
            self.assertEqual(
                str(summary_table.cell(0, 0).text_frame.paragraphs[0].runs[0].font.color.rgb),
                HEADER_TEXT_COLOR,
            )
            self.assertEqual(
                str(detail_table.cell(2, 0).text_frame.paragraphs[0].runs[0].font.color.rgb),
                BODY_TEXT_COLOR,
            )

            left_border = detail_table.cell(2, 0)._tc.tcPr.find(qn("a:lnL"))
            border_color = left_border.find(qn("a:solidFill")).find(qn("a:srgbClr")).get("val")
            self.assertEqual(border_color, BORDER_COLOR)

    def test_generate_presentation_inserts_category_intro_slides_once_per_matching_category(self) -> None:
        repo_root = Path(__file__).resolve().parents[1]
        template_path = repo_root / "templates" / "template.pptx"
        chapter_template_path = repo_root / "templates" / "chapter.pptx"

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            input_dir = temp_path / "input"
            input_dir.mkdir()
            output_path = temp_path / "report-with-chapters.pptx"

            for file_name in ("参展商.xlsx", "会议主承办.xlsx", "自助餐.xlsx"):
                create_report_workbook(
                    input_dir / file_name,
                    [
                        ("指标", "满意度", "重要性"),
                        (Path(file_name).stem, 9.9, 9.8),
                    ],
                )

            config = PptBatchConfig(
                template_path=template_path,
                input_dir=input_dir,
                output_ppt=output_path,
                category_intro_slides={
                    "一、会展客户": CategoryIntroSlideConfig(
                        ppt_path=chapter_template_path,
                        slide_number=3,
                    ),
                    "二、餐饮客户": CategoryIntroSlideConfig(
                        ppt_path=chapter_template_path,
                        slide_number=4,
                    ),
                    "五、酒店客户": CategoryIntroSlideConfig(
                        ppt_path=chapter_template_path,
                        slide_number=5,
                    ),
                },
            )

            generate_presentation(config)

            presentation = Presentation(output_path)
            slide_texts = [collect_slide_texts(slide) for slide in presentation.slides]

            self.assertEqual(len(presentation.slides), 5)
            self.assertIn("会展区客户满意度", slide_texts[0])
            self.assertEqual(len(presentation.slides[0].placeholders), 0)
            self.assertEqual(presentation.slides[1].shapes.title.text, "会展客户——参展商")
            self.assertEqual(presentation.slides[2].shapes.title.text, "会展客户——会议活动主（承）办")
            self.assertIn("餐饮区客户满意度", slide_texts[3])
            self.assertEqual(len(presentation.slides[3].placeholders), 0)
            self.assertEqual(presentation.slides[4].shapes.title.text, "餐饮客户——自助餐")
            self.assertEqual(
                sum("会展区客户满意度" in texts for texts in slide_texts),
                1,
            )
            self.assertEqual(
                sum("餐饮区客户满意度" in texts for texts in slide_texts),
                1,
            )
            self.assertFalse(
                any("酒店区客户满意度及酒店暗访评分" in texts for texts in slide_texts),
            )

    def test_generate_presentation_appends_chart_slide_with_same_title(self) -> None:
        repo_root = Path(__file__).resolve().parents[1]
        template_path = repo_root / "templates" / "template.pptx"

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            input_dir = temp_path / "input"
            input_dir.mkdir()
            output_path = temp_path / "report-with-chart.pptx"

            create_report_workbook(
                input_dir / "专业观众.xlsx",
                [
                    ("指标", "满意度", "重要性"),
                    ("专业观众", 9.93, 10.0),
                    ("会展服务", 9.86, 9.90),
                    ("工作人员仪容仪表", 10.0, 10.0),
                    ("硬件设施", 9.31, 9.89),
                    ("园区停车方便", 9.67, 9.92),
                    ("配套服务", 9.47, 9.82),
                    ("餐饮服务", 8.9, 9.8),
                    ("智慧场馆", 8.5, 9.0),
                    ("杭州国博APP", 9.0, 8.0),
                ],
            )

            config = PptBatchConfig(
                template_path=template_path,
                input_dir=input_dir,
                output_ppt=output_path,
                chart_page=ChartPageConfig(
                    enabled=True,
                    placeholder_text="图表分析内容待补充。",
                    image_dpi=120,
                ),
            )

            generate_presentation(config)

            presentation = Presentation(output_path)
            self.assertEqual(len(presentation.slides), 2)
            self.assertEqual(presentation.slides[0].shapes.title.text, "会展客户——专业观众")
            self.assertEqual(presentation.slides[1].shapes.title.text, "会展客户——专业观众")
            self.assertTrue(
                any(shape.shape_type == MSO_SHAPE_TYPE.PICTURE for shape in presentation.slides[1].shapes)
            )
            self.assertIn(
                "图表分析内容待补充。",
                collect_slide_texts(presentation.slides[1]),
            )
            textbox_shape = next(
                shape
                for shape in presentation.slides[1].shapes
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
            )
            textbox_run = textbox_shape.text_frame.paragraphs[0].runs[0]
            self.assertEqual(textbox_run.font.name, CHART_TEXTBOX_FONT_NAME)
            self.assertEqual(textbox_run.font.size, Pt(CHART_TEXTBOX_FONT_SIZE_PT))
            self.assertEqual(textbox_shape.text_frame.paragraphs[0].line_spacing, CHART_TEXTBOX_LINE_SPACING)
            self.assertEqual(
                textbox_shape.text_frame.paragraphs[0]._p.pPr.get("indent"),
                str(Pt(CHART_TEXTBOX_FIRST_LINE_INDENT_PT)),
            )

    def test_render_chart_textbox_reduces_line_spacing_for_long_text(self) -> None:
        presentation = Presentation()
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        long_text = (
            "总体判断：会议主承办满意度主要由会展服务支撑，现场对接、工作人员服务及整体协同均处高分；"
            "硬件设施与智慧场馆相对偏弱，拉低综合体验，其中标识标牌、设施设备齐全度更需留意。\n"
            "亮点：会展服务表现最突出，现场对接协调沟通、工作人员仪容仪表和服务态度均为满分，"
            "配套服务中的客房服务、保洁服务也有较强支撑。\n"
            "关注点：硬件设施是会议主承办体验短板，标识标牌清晰度得分最低，设施设备齐全、"
            "休息空间、餐饮服务分值也偏低；同时会后回访低于会展服务板块多数触点。"
        )

        render_chart_textbox(
            slide,
            TableRegion(6.55, 1.58, 5.50, 5.10),
            long_text,
        )

        textbox_shape = next(
            shape for shape in slide.shapes if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
        )
        font_sizes = {
            run.font.size.pt
            for paragraph in textbox_shape.text_frame.paragraphs
            for run in paragraph.runs
            if run.text.strip()
        }

        self.assertEqual(len(font_sizes), 1)
        self.assertEqual(font_sizes.pop(), CHART_TEXTBOX_FONT_SIZE_PT)
        self.assertLess(
            textbox_shape.text_frame.paragraphs[0].line_spacing,
            CHART_TEXTBOX_LINE_SPACING,
        )

    def test_resolve_chart_textbox_style_reduces_spacing_for_meeting_page_text(self) -> None:
        text = (
            "总体判断：会议主承办满意度主要由会展服务和配套服务支撑，现场对接、工作人员服务、"
            "客房与保洁评价突出；硬件设施相对承压，尤其标识标牌、设施设备齐全度及交通流线拉低体验。\n"
            "亮点：会展服务中现场对接协调沟通、工作人员服务态度与仪容仪表均获满分，"
            "配套服务中的客房服务和保洁服务也处于高分。\n"
            "关注点：会议主承办对硬件设施更敏感，标识标牌清晰、设施设备齐全、交通便利与"
            "交通流线等指标满意度低于其重要性，餐饮服务得分也偏低。"
        )

        text_style = resolve_chart_textbox_style(
            text,
            TableRegion(6.55, 1.58, 5.50, 5.10),
        )

        self.assertLess(text_style.line_spacing, CHART_TEXTBOX_LINE_SPACING)

    def test_resolve_chart_textbox_style_uses_tighter_spacing_for_q1_meeting_chart_text(self) -> None:
        text = (
            "总体判断：会议主承办满意度主要由会展服务支撑，现场对接、人员服务态度与效率等环节"
            "评分突出；拖累项集中在硬件设施，尤其标识标牌、设施设备齐全度及休息空间，拉低整体体验。\n"
            "亮点：会展服务表现最强，现场对接协调沟通、工作人员仪容仪表和服务态度均为满分，"
            "配套服务中的客房服务、保洁服务也维持较好水平。\n"
            "关注点：硬件设施短板更集中，标识标牌清晰度为各项最低，设施设备齐全、交通流线及"
            "交通便利性与重要性相比仍有提升空间；配套服务中的餐饮服务评分也偏低。"
        )

        text_style = resolve_chart_textbox_style(
            text,
            TableRegion(6.55, 1.58, 5.50, 5.10),
        )

        self.assertLessEqual(text_style.line_spacing, 1.02)

    def test_generate_presentation_reuses_llm_notes_in_chart_slide_textbox(self) -> None:
        repo_root = Path(__file__).resolve().parents[1]
        template_path = repo_root / "templates" / "template.pptx"

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            input_dir = temp_path / "input"
            input_dir.mkdir()
            output_path = temp_path / "report-with-chart-and-notes.pptx"
            env_path = temp_path / ".env"
            system_role_path = temp_path / "system_role.md"

            env_path.write_text(
                "OPENAI_API_KEY=test-key\n"
                "OPENAI_BASE_URL=https://example.com/v1\n"
                "OPENAI_MODEL=fake-model\n",
                encoding="utf-8",
            )
            system_role_path.write_text("你是测试用分析助手。", encoding="utf-8")

            create_report_workbook(
                input_dir / "专业观众.xlsx",
                [
                    ("指标", "满意度", "重要性"),
                    ("专业观众", 9.93, 10.0),
                    ("会展服务", 9.86, 9.90),
                    ("工作人员仪容仪表", 10.0, 10.0),
                    ("硬件设施", 9.31, 9.89),
                    ("园区停车方便", 9.67, 9.92),
                    ("配套服务", 9.47, 9.82),
                    ("餐饮服务", 8.9, 9.8),
                ],
            )

            config = PptBatchConfig(
                template_path=template_path,
                input_dir=input_dir,
                output_ppt=output_path,
                chart_page=ChartPageConfig(
                    enabled=True,
                    placeholder_text="图表分析内容待补充。",
                    image_dpi=120,
                ),
                llm_notes=LlmNotesConfig(
                    enabled=True,
                    env_path=env_path,
                    system_role_path=system_role_path,
                    target_chars=300,
                    temperature=0.2,
                    max_tokens=400,
                ),
            )

            generate_presentation(config, llm_client_factory=FakeOpenAI)

            presentation = Presentation(output_path)
            notes_text = presentation.slides[0].notes_slide.notes_text_frame.text
            chart_texts = collect_slide_texts(presentation.slides[1])
            textbox_shape = next(
                shape
                for shape in presentation.slides[1].shapes
                if shape.shape_type == MSO_SHAPE_TYPE.TEXT_BOX
            )

            self.assertEqual(len(presentation.slides), 2)
            self.assertTrue(notes_text)
            for line in notes_text.splitlines():
                self.assertIn(line, chart_texts)
            self.assertEqual(len(textbox_shape.text_frame.paragraphs), 3)
            for paragraph in textbox_shape.text_frame.paragraphs:
                self.assertEqual(
                    paragraph._p.pPr.get("indent"),
                    str(Pt(CHART_TEXTBOX_FIRST_LINE_INDENT_PT)),
                )
            self.assertNotIn("图表分析内容待补充。", chart_texts)

    def test_generate_presentation_skips_sections_when_all_metric_satisfaction_values_are_empty(self) -> None:
        repo_root = Path(__file__).resolve().parents[1]
        template_path = repo_root / "templates" / "template.pptx"

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            input_dir = temp_path / "input"
            input_dir.mkdir()
            output_path = temp_path / "report.pptx"

            create_report_workbook(
                input_dir / "专业观众.xlsx",
                [
                    ("指标", "满意度", "重要性"),
                    ("专业观众", 9.93, 10.0),
                    ("会展服务", 10.0, 10.0),
                    ("工作人员仪容仪表", 10.0, 10.0),
                    ("工作人员服务态度", 9.9, 10.0),
                    ("智慧场馆", None, None),
                    ("杭州国博APP", None, 10.0),
                    ("室内导航系统", None, 10.0),
                    ("配套服务", 9.8, 10.0),
                    ("餐饮服务", 9.4, 10.0),
                ],
            )

            config = PptBatchConfig(
                template_path=template_path,
                input_dir=input_dir,
                output_ppt=output_path,
                blank_display="",
                max_single_table_rows=10,
                max_split_table_rows=19,
                layout=PptLayoutConfig(),
            )

            generate_presentation(config)

            presentation = Presentation(output_path)
            detail_tables = [
                shape.table
                for shape in presentation.slides[0].shapes
                if getattr(shape, "has_table", False)
            ][1:]
            detail_text = "\n".join(
                cell.text
                for table in detail_tables
                for row in table.rows
                for cell in row.cells
            )

            self.assertIn("会展服务", detail_text)
            self.assertIn("配套服务", detail_text)
            self.assertNotIn("智慧场馆", detail_text)
            self.assertNotIn("杭州国博APP", detail_text)
            self.assertNotIn("室内导航系统", detail_text)

    def test_generate_presentation_writes_llm_notes_from_env_and_system_role(self) -> None:
        repo_root = Path(__file__).resolve().parents[1]
        template_path = repo_root / "templates" / "template.pptx"

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            input_dir = temp_path / "input"
            input_dir.mkdir()
            output_path = temp_path / "notes-report.pptx"
            env_path = temp_path / ".env"
            system_role_path = temp_path / "system_role.md"

            env_path.write_text(
                "OPENAI_API_KEY=test-key\n"
                "OPENAI_BASE_URL=https://example.com/v1\n"
                "OPENAI_MODEL=fake-model\n"
                "OPENAI_TEMPERATURE=0.75\n",
                encoding="utf-8",
            )
            system_role_path.write_text("你是测试用分析助手。", encoding="utf-8")

            create_report_workbook(
                input_dir / "专业观众.xlsx",
                [
                    ("指标", "满意度", "重要性"),
                    ("专业观众", 9.93, 10.0),
                    ("会展服务", 10.0, 10.0),
                    ("工作人员仪容仪表", 10.0, 10.0),
                    ("智慧场馆", 10.0, 10.0),
                    ("室内导航系统", None, None),
                ],
            )

            config = PptBatchConfig(
                template_path=template_path,
                input_dir=input_dir,
                output_ppt=output_path,
                max_single_table_rows=3,
                max_split_table_rows=4,
                llm_notes=LlmNotesConfig(
                    enabled=True,
                    env_path=env_path,
                    system_role_path=system_role_path,
                    target_chars=300,
                    temperature=0.2,
                    max_tokens=400,
                ),
            )

            stdout_buffer = io.StringIO()
            with redirect_stdout(stdout_buffer):
                generate_presentation(config, llm_client_factory=FakeOpenAI)

            self.assertEqual(len(FakeOpenAI.instances), 1)
            fake_client = FakeOpenAI.instances[0]
            self.assertEqual(fake_client.kwargs["api_key"], "test-key")
            self.assertEqual(fake_client.kwargs["base_url"], "https://example.com/v1")
            self.assertEqual(fake_client.create_calls[0]["model"], "fake-model")
            self.assertTrue(fake_client.create_calls[0]["stream"])
            self.assertEqual(fake_client.create_calls[0]["temperature"], 0.75)
            self.assertEqual(fake_client.create_calls[0]["max_tokens"], 400)
            self.assertEqual(
                fake_client.create_calls[0]["messages"][0]["content"],
                "你是测试用分析助手。",
            )
            self.assertIn(
                "页面标题：会展客户——专业观众",
                fake_client.create_calls[0]["messages"][1]["content"],
            )
            self.assertNotIn("空值项：", fake_client.create_calls[0]["messages"][1]["content"])
            self.assertNotIn("室内导航系统", fake_client.create_calls[0]["messages"][1]["content"])

            presentation = Presentation(output_path)
            notes_text = presentation.slides[0].notes_slide.notes_text_frame.text
            self.assertIn("总体判断：", notes_text)
            self.assertIn("关注点：", notes_text)
            self.assertNotIn("本页", notes_text)
            self.assertIn("专业观众", notes_text)

            progress_output = stdout_buffer.getvalue()
            self.assertIn("[1/1] 正在生成备注页分析：会展客户——专业观众", progress_output)
            self.assertIn("[1/1] 流式输出：", progress_output)
            self.assertNotIn("本页", progress_output)
            self.assertIn("[1/1] 备注页分析完成：会展客户——专业观众", progress_output)

    def test_generate_presentation_retries_transient_llm_stream_failures(self) -> None:
        repo_root = Path(__file__).resolve().parents[1]
        template_path = repo_root / "templates" / "template.pptx"

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            input_dir = temp_path / "input"
            input_dir.mkdir()
            output_path = temp_path / "notes-report.pptx"
            env_path = temp_path / ".env"
            system_role_path = temp_path / "system_role.md"

            env_path.write_text(
                "OPENAI_API_KEY=test-key\n"
                "OPENAI_BASE_URL=https://example.com/v1\n"
                "OPENAI_MODEL=fake-model\n",
                encoding="utf-8",
            )
            system_role_path.write_text("你是测试用分析助手。", encoding="utf-8")

            create_report_workbook(
                input_dir / "专业观众.xlsx",
                [
                    ("指标", "满意度", "重要性"),
                    ("专业观众", 9.93, 10.0),
                    ("会展服务", 10.0, 10.0),
                    ("工作人员仪容仪表", 10.0, 10.0),
                    ("配套服务", 9.2, 9.8),
                    ("餐饮服务", 8.8, 9.6),
                ],
            )

            config = PptBatchConfig(
                template_path=template_path,
                input_dir=input_dir,
                output_ppt=output_path,
                llm_notes=LlmNotesConfig(
                    enabled=True,
                    env_path=env_path,
                    system_role_path=system_role_path,
                    target_chars=120,
                    temperature=0.2,
                    max_tokens=200,
                ),
            )

            generate_presentation(config, llm_client_factory=FakeFlakyOpenAI)

            self.assertEqual(len(FakeFlakyOpenAI.instances), 1)
            fake_client = FakeFlakyOpenAI.instances[0]
            self.assertEqual(len(fake_client.create_calls), 3)
            presentation = Presentation(output_path)
            notes_text = presentation.slides[0].notes_slide.notes_text_frame.text
            self.assertIn("总体判断：", notes_text)
            self.assertIn("关注点：", notes_text)

    def test_generate_presentation_falls_back_to_placeholder_and_writes_failure_log(self) -> None:
        repo_root = Path(__file__).resolve().parents[1]
        template_path = repo_root / "templates" / "template.pptx"

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            input_dir = temp_path / "input"
            input_dir.mkdir()
            output_path = temp_path / "notes-report.pptx"
            failure_log_path = build_llm_failure_log_path(output_path)
            env_path = temp_path / ".env"
            system_role_path = temp_path / "system_role.md"

            env_path.write_text(
                "OPENAI_API_KEY=test-key\n"
                "OPENAI_BASE_URL=https://example.com/v1\n"
                "OPENAI_MODEL=fake-model\n",
                encoding="utf-8",
            )
            system_role_path.write_text("你是测试用分析助手。", encoding="utf-8")

            create_report_workbook(
                input_dir / "专业观众.xlsx",
                [
                    ("指标", "满意度", "重要性"),
                    ("专业观众", 9.93, 10.0),
                    ("会展服务", 10.0, 10.0),
                    ("工作人员仪容仪表", 10.0, 10.0),
                    ("配套服务", 9.2, 9.8),
                    ("餐饮服务", 8.8, 9.6),
                ],
            )
            create_report_workbook(
                input_dir / "参展商.xlsx",
                [
                    ("指标", "满意度", "重要性"),
                    ("参展商", 9.58, 9.81),
                    ("会场服务", 10.0, 10.0),
                    ("工作人员服务态度", 10.0, 10.0),
                    ("硬件设施", 8.0, 9.5),
                    ("园区停车方便", 8.0, 9.6),
                ],
            )

            config = PptBatchConfig(
                template_path=template_path,
                input_dir=input_dir,
                output_ppt=output_path,
                chart_page=ChartPageConfig(
                    enabled=True,
                    placeholder_text="图表分析内容待补充。",
                    image_dpi=120,
                ),
                llm_notes=LlmNotesConfig(
                    enabled=True,
                    env_path=env_path,
                    system_role_path=system_role_path,
                    target_chars=120,
                    temperature=0.2,
                    max_tokens=200,
                ),
            )

            stdout_buffer = io.StringIO()
            with (
                mock.patch.object(
                    generate_ppt_module,
                    "DEFAULT_LLM_STREAM_RETRY_TIMEOUT_SECONDS",
                    0.0,
                ),
                redirect_stdout(stdout_buffer),
            ):
                generate_presentation(config, llm_client_factory=FakeAlwaysFailOpenAI)

            presentation = Presentation(output_path)
            self.assertEqual(len(presentation.slides), 4)
            first_notes = presentation.slides[0].notes_slide.notes_text_frame.text
            second_notes = presentation.slides[2].notes_slide.notes_text_frame.text
            self.assertIn("大模型分析暂未生成", first_notes)
            self.assertIn("大模型分析暂未生成", second_notes)
            self.assertIn(
                "大模型分析暂未生成",
                "\n".join(collect_slide_texts(presentation.slides[1])),
            )
            self.assertTrue(failure_log_path.exists())
            failure_log_text = failure_log_path.read_text(encoding="utf-8")
            self.assertIn("会展客户——专业观众", failure_log_text)
            self.assertIn("会展客户——参展商", failure_log_text)
            self.assertIn("RuntimeError", failure_log_text)

            progress_output = stdout_buffer.getvalue()
            self.assertIn("备注页分析失败，已使用占位文本", progress_output)
            self.assertIn("LLM 失败清单已写入", progress_output)

    def test_generate_presentation_preserves_checkpoint_when_llm_is_interrupted(self) -> None:
        repo_root = Path(__file__).resolve().parents[1]
        template_path = repo_root / "templates" / "template.pptx"

        with tempfile.TemporaryDirectory() as temp_dir:
            temp_path = Path(temp_dir)
            input_dir = temp_path / "input"
            input_dir.mkdir()
            output_path = temp_path / "notes-report.pptx"
            partial_path = build_partial_output_path(output_path)
            env_path = temp_path / ".env"
            system_role_path = temp_path / "system_role.md"

            env_path.write_text(
                "OPENAI_API_KEY=test-key\n"
                "OPENAI_BASE_URL=https://example.com/v1\n"
                "OPENAI_MODEL=fake-model\n",
                encoding="utf-8",
            )
            system_role_path.write_text("你是测试用分析助手。", encoding="utf-8")

            create_report_workbook(
                input_dir / "第一页.xlsx",
                [
                    ("指标", "满意度", "重要性"),
                    ("第一页", 9.93, 10.0),
                    ("会展服务", 10.0, 10.0),
                    ("工作人员仪容仪表", 10.0, 10.0),
                ],
            )
            create_report_workbook(
                input_dir / "第二页.xlsx",
                [
                    ("指标", "满意度", "重要性"),
                    ("第二页", 9.83, 9.5),
                    ("会展服务", 9.7, 9.4),
                    ("工作人员服务态度", 9.6, 9.3),
                ],
            )

            config = PptBatchConfig(
                template_path=template_path,
                input_dir=input_dir,
                output_ppt=output_path,
                llm_notes=LlmNotesConfig(
                    enabled=True,
                    env_path=env_path,
                    system_role_path=system_role_path,
                    target_chars=300,
                    temperature=0.2,
                    max_tokens=400,
                    checkpoint_chars=20,
                ),
            )

            stdout_buffer = io.StringIO()
            with self.assertRaises(KeyboardInterrupt):
                with redirect_stdout(stdout_buffer):
                    generate_presentation(config, llm_client_factory=FakeInterruptedOpenAI)

            self.assertTrue(partial_path.exists())
            partial_presentation = Presentation(partial_path)
            self.assertEqual(len(partial_presentation.slides), 2)
            self.assertIn(
                "第一页分析已经完成",
                partial_presentation.slides[0].notes_slide.notes_text_frame.text,
            )
            self.assertIn(
                "第二页分析进行到一半",
                partial_presentation.slides[1].notes_slide.notes_text_frame.text,
            )

            progress_output = stdout_buffer.getvalue()
            self.assertIn("[1/2] 已保存检查点：notes-report.partial.pptx", progress_output)
            self.assertIn("生成中断，已保存当前检查点", progress_output)


if __name__ == "__main__":
    unittest.main()
