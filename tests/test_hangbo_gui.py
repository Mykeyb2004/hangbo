from __future__ import annotations

import tempfile
import unittest
from pathlib import Path
from unittest import mock

import pandas as pd
from pptx import Presentation

from hangbo_gui import (
    CustomerTypePreviewRow,
    CustomerTypePreviewStatus,
    GuiBatchConfig,
    MainWorkflowSelection,
    SavedBatchProfile,
    StatsPreviewSummary,
    SurveyPlatformApp,
    ThemePalette,
    WorkflowRunController,
    WorkflowRunStatus,
    WorkflowMode,
    batch_profile_storage_path,
    build_ppt_thumbnail_cache_dir,
    build_gui_batch_config_text,
    build_category_intro_slides_text,
    build_stats_preview_summary_text,
    build_survey_stats_command,
    build_workflow_status_text,
    build_main_workflow_step_keys,
    build_merge_command,
    build_ppt_config_text,
    build_stats_preview_summary,
    build_task_commands,
    build_survey_stats_config_text,
    delete_batch_profile,
    default_selected_customer_types,
    discover_ppt_slide_previews,
    generate_ppt_slide_thumbnail_images,
    load_gui_batch_config,
    load_gui_session,
    load_saved_batch_profiles,
    ordered_selected_customer_types,
    parse_category_intro_slides_text,
    save_batch_profile,
    save_gui_session,
)


class GuiBatchConfigTests(unittest.TestCase):
    def test_effective_input_dir_uses_single_dir_for_single_mode(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年3月",
            workflow_mode=WorkflowMode.SINGLE,
            single_input_dir=Path("/tmp/datas/3月"),
        )

        self.assertEqual(config.effective_input_dir(), Path("/tmp/datas/3月"))

    def test_effective_input_dir_uses_merged_dir_for_merged_mode(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年Q1",
            workflow_mode=WorkflowMode.MERGED,
            single_input_dir=Path("/tmp/datas/3月"),
            merge_output_dir=Path("/tmp/datas/合并结果"),
        )

        self.assertEqual(config.effective_input_dir(), Path("/tmp/datas/合并结果"))

    def test_build_merge_command_returns_none_for_single_mode(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年3月",
            workflow_mode=WorkflowMode.SINGLE,
            single_input_dir=Path("/tmp/datas/3月"),
        )

        self.assertIsNone(build_merge_command(config))

    def test_build_merge_command_includes_all_input_dirs_for_merged_mode(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年Q1",
            workflow_mode=WorkflowMode.MERGED,
            merge_input_dirs=(
                Path("/tmp/datas/1月"),
                Path("/tmp/datas/2月"),
                Path("/tmp/datas/3月"),
            ),
            merge_output_dir=Path("/tmp/datas/合并结果"),
        )

        command = build_merge_command(config)

        self.assertIsNotNone(command)
        assert command is not None
        self.assertEqual(command[1], str(Path("/Users/zhangqijin/PycharmProjects/hangbo/merge_questionnaire_workbooks.py")))
        self.assertEqual(command.count("--input-dir"), 3)
        self.assertIn(str(Path("/tmp/datas/合并结果")), command)


class GuiWorkflowTests(unittest.TestCase):
    def test_single_month_workflow_order_matches_business_sequence(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年3月",
            workflow_mode=WorkflowMode.SINGLE,
            single_input_dir=Path("/tmp/datas/3月"),
        )
        selection = MainWorkflowSelection()

        step_keys = build_main_workflow_step_keys(config, selection)

        self.assertEqual(
            step_keys,
            ("phase_preprocess", "survey_stats", "summary_table", "generate_ppt"),
        )

    def test_merged_workflow_puts_merge_first_only_when_selected(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年Q1",
            workflow_mode=WorkflowMode.MERGED,
            merge_input_dirs=(Path("/tmp/datas/1月"), Path("/tmp/datas/2月")),
            merge_output_dir=Path("/tmp/datas/合并结果"),
        )
        selection = MainWorkflowSelection(include_merge=True)

        step_keys = build_main_workflow_step_keys(config, selection)

        self.assertEqual(
            step_keys,
            ("merge_workbooks", "phase_preprocess", "survey_stats", "summary_table", "generate_ppt"),
        )

    def test_fill_year_month_can_be_inserted_as_optional_preprocess_step(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年3月",
            workflow_mode=WorkflowMode.SINGLE,
            single_input_dir=Path("/tmp/datas/3月"),
        )
        selection = MainWorkflowSelection(include_fill_year_month=True)

        step_keys = build_main_workflow_step_keys(config, selection)

        self.assertEqual(
            step_keys,
            ("phase_preprocess", "fill_year_month", "survey_stats", "summary_table", "generate_ppt"),
        )

    def test_build_task_commands_use_clear_preprocess_titles(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年3月",
            workflow_mode=WorkflowMode.SINGLE,
            single_input_dir=Path("/tmp/datas/3月"),
        )
        selection = MainWorkflowSelection(
            include_fill_year_month=True,
            include_survey_stats=False,
            include_summary=False,
            include_ppt=False,
        )

        with mock.patch(
            "hangbo_gui.build_phase_preprocess_command",
            return_value=["python", "phase_column_preprocess.py"],
        ):
            with mock.patch(
                "hangbo_gui.build_fill_year_month_command",
                return_value=["python", "fill_year_month_columns.py"],
            ):
                tasks = build_task_commands(config, selection)

        self.assertEqual(
            [task.title for task in tasks],
            [
                "兼容新版调查问卷数据结构",
                "在数据源中加入年份+月份",
            ],
        )


class GuiPreprocessCopyTests(unittest.TestCase):
    def test_build_preprocess_page_uses_business_friendly_copy(self) -> None:
        app = object.__new__(SurveyPlatformApp)
        app.palette = ThemePalette()
        app.year_value_var = mock.Mock()
        app.month_value_var = mock.Mock()
        app.run_phase_preprocess_task = mock.Mock()
        app.run_fill_year_month_task = mock.Mock()
        app._register_start_button = mock.Mock()

        label_texts: list[str] = []
        frame_texts: list[str] = []
        button_texts: list[str] = []

        def fake_widget(*_: object, **__: object) -> mock.Mock:
            widget = mock.Mock()
            widget.grid = mock.Mock()
            widget.pack = mock.Mock()
            widget.columnconfigure = mock.Mock()
            return widget

        def fake_label(*_: object, **kwargs: object) -> mock.Mock:
            text = kwargs.get("text")
            if isinstance(text, str):
                label_texts.append(text)
            return fake_widget()

        def fake_label_frame(*_: object, **kwargs: object) -> mock.Mock:
            text = kwargs.get("text")
            if isinstance(text, str):
                frame_texts.append(text)
            return fake_widget()

        def fake_button(*_: object, **kwargs: object) -> mock.Mock:
            text = kwargs.get("text")
            if isinstance(text, str):
                button_texts.append(text)
            return fake_widget()

        with mock.patch("hangbo_gui.ttk.Frame", side_effect=fake_widget):
            with mock.patch("hangbo_gui.ttk.LabelFrame", side_effect=fake_label_frame):
                with mock.patch("hangbo_gui.ttk.Label", side_effect=fake_label):
                    with mock.patch("hangbo_gui.ttk.Entry", side_effect=fake_widget):
                        with mock.patch("hangbo_gui.ttk.Button", side_effect=fake_button):
                            SurveyPlatformApp._build_preprocess_page(app, mock.Mock())

        self.assertIn("兼容新版调查问卷数据结构", frame_texts)
        self.assertIn("在数据源中加入年份+月份", frame_texts)
        self.assertIn(
            "说明：如果新版调查问卷在第三列增加了“一期/二期”等期次字段，这一步会自动把该列移到最后，避免后续统计错位。",
            label_texts,
        )
        self.assertIn(
            "说明：给问卷数据补写“年份”“月份”两列，方便后续合并文件。",
            label_texts,
        )
        self.assertIn(
            "执行兼容新版结构（phase_column_preprocess.py）",
            button_texts,
        )
        self.assertIn(
            "执行补写年份+月份（fill_year_month_columns.py）",
            button_texts,
        )


class GuiConfigRenderingTests(unittest.TestCase):
    def test_build_ppt_thumbnail_cache_dir_uses_file_metadata(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            ppt_path = Path(temp_dir) / "chapter sample.pptx"
            ppt_path.write_bytes(b"fake-ppt-content")
            thumbnail_root = Path(temp_dir) / "thumb-cache"
            with mock.patch("hangbo_gui.GUI_THUMBNAIL_DIR", thumbnail_root):
                cache_dir = build_ppt_thumbnail_cache_dir(ppt_path)

        self.assertIn("chapter_sample", cache_dir.name)

    def test_generate_ppt_slide_thumbnail_images_reuses_existing_cache(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            ppt_path = Path(temp_dir) / "chapter.pptx"
            ppt_path.write_bytes(b"fake-ppt-content")
            thumbnail_root = Path(temp_dir) / "thumb-cache"
            with mock.patch("hangbo_gui.GUI_THUMBNAIL_DIR", thumbnail_root):
                cache_dir = build_ppt_thumbnail_cache_dir(ppt_path)
                cache_dir.mkdir(parents=True, exist_ok=True)
                cached_thumbnail = cache_dir / "chapter-001.png"
                cached_thumbnail.write_bytes(b"\x89PNG\r\n\x1a\n")

                with mock.patch("hangbo_gui.subprocess.run") as run_mock:
                    thumbnails = generate_ppt_slide_thumbnail_images(ppt_path)

        self.assertEqual(thumbnails, (cached_thumbnail,))
        run_mock.assert_not_called()

    def test_generate_ppt_slide_thumbnail_images_runs_conversion_pipeline(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            ppt_path = Path(temp_dir) / "chapter.pptx"
            ppt_path.write_bytes(b"fake-ppt-content")
            thumbnail_root = Path(temp_dir) / "thumb-cache"
            with mock.patch("hangbo_gui.GUI_THUMBNAIL_DIR", thumbnail_root):
                cache_dir = build_ppt_thumbnail_cache_dir(ppt_path)

                def fake_run(command: list[str], **_: object) -> mock.Mock:
                    if command[0] == "soffice":
                        (cache_dir / "chapter.pdf").write_bytes(b"%PDF-1.4")
                    elif command[0] == "gs":
                        (cache_dir / "chapter-001.png").write_bytes(b"\x89PNG\r\n\x1a\n")
                        (cache_dir / "chapter-002.png").write_bytes(b"\x89PNG\r\n\x1a\n")
                    return mock.Mock(stdout="", stderr="")

                with mock.patch(
                    "hangbo_gui.subprocess.run",
                    side_effect=fake_run,
                ) as run_mock:
                    thumbnails = generate_ppt_slide_thumbnail_images(ppt_path)

        self.assertEqual(
            thumbnails,
            (
                cache_dir / "chapter-001.png",
                cache_dir / "chapter-002.png",
            ),
        )
        self.assertEqual(run_mock.call_count, 2)

    def test_discover_ppt_slide_previews_reads_slide_titles(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            ppt_path = Path(temp_dir) / "chapter.pptx"
            presentation = Presentation()
            first_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            first_slide.shapes.title.text = "一、会展客户"
            second_slide = presentation.slides.add_slide(presentation.slide_layouts[0])
            second_slide.shapes.title.text = "五、酒店客户封面"
            presentation.save(ppt_path)

            previews = discover_ppt_slide_previews(ppt_path)

        self.assertEqual(len(previews), 2)
        self.assertEqual(previews[0].slide_number, 1)
        self.assertEqual(previews[0].title, "一、会展客户")
        self.assertEqual(previews[0].label, "1. 一、会展客户")
        self.assertEqual(previews[1].label, "2. 五、酒店客户封面")

    def test_discover_ppt_slide_previews_falls_back_when_no_text(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            ppt_path = Path(temp_dir) / "blank.pptx"
            presentation = Presentation()
            presentation.slides.add_slide(presentation.slide_layouts[6])
            presentation.save(ppt_path)

            previews = discover_ppt_slide_previews(ppt_path)

        self.assertEqual(len(previews), 1)
        self.assertEqual(previews[0].label, "1. 未识别到文字")

    def test_build_category_intro_slides_text_round_trip(self) -> None:
        text = build_category_intro_slides_text(
            [
                ("一、会展客户", "templates/chapter.pptx", 3),
                ("五、酒店客户", "templates/chapter.pptx", 5),
            ]
        )

        self.assertEqual(
            parse_category_intro_slides_text(text),
            (
                ("一、会展客户", "templates/chapter.pptx", 3),
                ("五、酒店客户", "templates/chapter.pptx", 5),
            ),
        )

    def test_build_survey_stats_config_text_uses_effective_input_dir(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年Q1",
            workflow_mode=WorkflowMode.MERGED,
            merge_output_dir=Path("/tmp/datas/合并结果"),
            stats_output_dir=Path("/tmp/output/Q1"),
            calculation_mode="summary",
            sheet_name="问卷数据",
        )

        text = build_survey_stats_config_text(config)

        self.assertIn('input_dir = "/tmp/datas/合并结果"', text)
        self.assertIn('output_dir = "/tmp/output/Q1"', text)
        self.assertIn('calculation_mode = "summary"', text)
        self.assertIn('sheet_name = "问卷数据"', text)

    def test_build_ppt_config_text_uses_stats_output_dir_as_input(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年3月",
            workflow_mode=WorkflowMode.SINGLE,
            single_input_dir=Path("/tmp/datas/3月"),
            stats_output_dir=Path("/tmp/output/3月"),
            ppt_template_path=Path("/tmp/templates/template.pptx"),
            output_ppt_path=Path("/tmp/output/3月满意度报告.pptx"),
            ppt_section_mode="auto",
        )

        text = build_ppt_config_text(config)

        self.assertIn('template_path = "/tmp/templates/template.pptx"', text)
        self.assertIn('input_dir = "/tmp/output/3月"', text)
        self.assertIn('output_ppt = "/tmp/output/3月满意度报告.pptx"', text)
        self.assertIn('section_mode = "auto"', text)

    def test_build_ppt_config_text_includes_advanced_sections(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年Q1",
            workflow_mode=WorkflowMode.SINGLE,
            single_input_dir=Path("/tmp/datas/Q1"),
            stats_output_dir=Path("/tmp/output/Q1"),
            ppt_template_path=Path("/tmp/templates/template.pptx"),
            output_ppt_path=Path("/tmp/output/Q1报告.pptx"),
            ppt_file_pattern="*.xlsm",
            ppt_sheet_name_mode="named",
            ppt_sheet_name="结果页",
            ppt_title_suffix="（季度版）",
            ppt_blank_display="-",
            ppt_max_single_table_rows="20",
            ppt_max_split_table_rows="22",
            ppt_body_font_size_pt="11.5",
            ppt_header_font_size_pt="12",
            ppt_summary_font_size_pt="13",
            ppt_template_slide_index="2",
            ppt_chart_page_enabled=True,
            ppt_chart_placeholder_text="图表说明文案",
            ppt_chart_image_dpi="300",
            ppt_llm_notes_enabled=True,
            ppt_llm_env_path=".env.prod",
            ppt_llm_system_role_path="roles/system_role.md",
            ppt_llm_target_chars="360",
            ppt_llm_temperature="0.6",
            ppt_llm_max_tokens="700",
            ppt_llm_checkpoint_chars="120",
            ppt_category_intro_slides_text=(
                "一、会展客户|templates/chapter.pptx|3\n"
                "五、酒店客户|templates/chapter.pptx|5"
            ),
            ppt_layout_summary_table_left="0.9",
            ppt_layout_chart_textbox_width="5.8",
        )

        text = build_ppt_config_text(config)

        self.assertIn('file_pattern = "*.xlsm"', text)
        self.assertIn('sheet_name_mode = "named"', text)
        self.assertIn('sheet_name = "结果页"', text)
        self.assertIn('title_suffix = "（季度版）"', text)
        self.assertIn("max_single_table_rows = 20", text)
        self.assertIn("template_slide_index = 2", text)
        self.assertIn('[category_intro_slides."一、会展客户"]', text)
        self.assertIn("[chart_page]", text)
        self.assertIn("enabled = true", text)
        self.assertIn('placeholder_text = "图表说明文案"', text)
        self.assertIn("[llm_notes]", text)
        self.assertIn('env_path = ".env.prod"', text)
        self.assertIn("[layout.summary_table]", text)
        self.assertIn("left = 0.9", text)
        self.assertIn("[layout.chart_textbox]", text)
        self.assertIn("width = 5.8", text)

    def test_build_survey_stats_command_appends_selected_jobs_in_order(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年3月",
            workflow_mode=WorkflowMode.SINGLE,
            single_input_dir=Path("/tmp/datas/3月"),
            stats_output_dir=Path("/tmp/output/3月"),
        )

        with mock.patch(
            "hangbo_gui.write_runtime_config",
            return_value=Path("/tmp/gui_runtime/survey_stats_001.toml"),
        ):
            command = build_survey_stats_command(
                config,
                dry_run=True,
                selected_job_names=("展览主承办", "散客"),
            )

        selected_job_names = [
            command[index + 1]
            for index, value in enumerate(command[:-1])
            if value == "--job"
        ]
        self.assertEqual(selected_job_names, ["展览主承办", "散客"])
        self.assertIn("--dry-run", command)

    def test_build_gui_batch_config_text_round_trip_preserves_merge_inputs(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年Q1",
            workflow_mode=WorkflowMode.MERGED,
            single_input_dir=Path("/tmp/datas/3月"),
            merge_input_dirs=(Path("/tmp/datas/1月"), Path("/tmp/datas/2月")),
            merge_output_dir=Path("/tmp/datas/合并结果"),
            year_value="2026",
            month_value="03",
            stats_output_dir=Path("/tmp/output/Q1"),
            summary_output_dir=Path("/tmp/summary/Q1"),
            output_ppt_path=Path("/tmp/output/Q1报告.pptx"),
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            profile_path = Path(temp_dir) / "batch.toml"
            profile_path.write_text(
                build_gui_batch_config_text(config),
                encoding="utf-8",
            )

            loaded = load_gui_batch_config(profile_path)

        self.assertEqual(loaded.batch_name, "2026年Q1")
        self.assertEqual(loaded.workflow_mode, WorkflowMode.MERGED)
        self.assertEqual(
            loaded.merge_input_dirs,
            (Path("/tmp/datas/1月").resolve(), Path("/tmp/datas/2月").resolve()),
        )
        self.assertEqual(loaded.output_ppt_path, Path("/tmp/output/Q1报告.pptx").resolve())
        self.assertEqual(loaded.ppt_sheet_name_mode, "first")

    def test_build_gui_batch_config_text_round_trip_preserves_ppt_advanced_fields(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年Q1",
            ppt_file_pattern="*.xlsm",
            ppt_sheet_name_mode="named",
            ppt_sheet_name="结果页",
            ppt_title_suffix="（季度版）",
            ppt_chart_page_enabled=True,
            ppt_chart_placeholder_text="图表说明\n第二行",
            ppt_llm_notes_enabled=True,
            ppt_llm_env_path=".env.prod",
            ppt_category_intro_slides_text="一、会展客户|templates/chapter.pptx|3",
            ppt_layout_chart_textbox_width="5.8",
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            profile_path = Path(temp_dir) / "batch.toml"
            profile_path.write_text(
                build_gui_batch_config_text(config),
                encoding="utf-8",
            )

            loaded = load_gui_batch_config(profile_path)

        self.assertEqual(loaded.ppt_file_pattern, "*.xlsm")
        self.assertEqual(loaded.ppt_sheet_name_mode, "named")
        self.assertEqual(loaded.ppt_sheet_name, "结果页")
        self.assertEqual(loaded.ppt_title_suffix, "（季度版）")
        self.assertTrue(loaded.ppt_chart_page_enabled)
        self.assertEqual(loaded.ppt_chart_placeholder_text, "图表说明\n第二行")
        self.assertTrue(loaded.ppt_llm_notes_enabled)
        self.assertEqual(loaded.ppt_llm_env_path, ".env.prod")
        self.assertEqual(
            loaded.ppt_category_intro_slides_text,
            "一、会展客户|templates/chapter.pptx|3",
        )
        self.assertEqual(loaded.ppt_layout_chart_textbox_width, "5.8")


class GuiThemeTests(unittest.TestCase):
    def test_apply_theme_uses_surface_background_for_radio_and_check_controls(self) -> None:
        app = object.__new__(SurveyPlatformApp)
        app.palette = ThemePalette()
        app.configure = mock.Mock()

        with mock.patch("hangbo_gui.ttk.Style") as style_factory:
            SurveyPlatformApp._apply_theme(app)

        style = style_factory.return_value
        style.configure.assert_any_call(
            "TRadiobutton",
            background=app.palette.surface,
            foreground=app.palette.text,
            font=("PingFang SC", 10),
        )
        style.configure.assert_any_call(
            "TCheckbutton",
            background=app.palette.surface,
            foreground=app.palette.text,
            font=("PingFang SC", 10),
        )
        style.map.assert_any_call("TRadiobutton", background=[("active", app.palette.surface)])
        style.map.assert_any_call("TCheckbutton", background=[("active", app.palette.surface)])

    def test_apply_theme_defines_root_label_styles_for_background_sections(self) -> None:
        app = object.__new__(SurveyPlatformApp)
        app.palette = ThemePalette()
        app.configure = mock.Mock()

        with mock.patch("hangbo_gui.ttk.Style") as style_factory:
            SurveyPlatformApp._apply_theme(app)

        style = style_factory.return_value
        style.configure.assert_any_call(
            "Root.SubHeader.TLabel",
            background=app.palette.background,
            foreground=app.palette.text,
            font=("PingFang SC", 12, "bold"),
        )
        style.configure.assert_any_call(
            "Root.Body.TLabel",
            background=app.palette.background,
            foreground=app.palette.text,
            font=("PingFang SC", 10),
        )
        style.configure.assert_any_call(
            "Root.Muted.TLabel",
            background=app.palette.background,
            foreground=app.palette.muted_text,
            font=("PingFang SC", 10),
        )


class GuiBatchPersistenceTests(unittest.TestCase):
    def test_save_and_list_batch_profiles_round_trip(self) -> None:
        config_a = GuiBatchConfig(
            batch_name="2026年3月",
            workflow_mode=WorkflowMode.SINGLE,
            single_input_dir=Path("/tmp/datas/3月"),
        )
        config_b = GuiBatchConfig(
            batch_name="2026年Q1",
            workflow_mode=WorkflowMode.MERGED,
            merge_input_dirs=(Path("/tmp/datas/1月"), Path("/tmp/datas/2月")),
            merge_output_dir=Path("/tmp/datas/合并结果"),
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            batch_dir = Path(temp_dir) / "batches"

            path_a = save_batch_profile(config_a, batch_dir=batch_dir)
            path_b = save_batch_profile(config_b, batch_dir=batch_dir)
            profiles = load_saved_batch_profiles(batch_dir=batch_dir)

        self.assertTrue(path_a.name.startswith("batch_"))
        self.assertTrue(path_b.name.startswith("batch_"))
        self.assertEqual([profile.batch_name for profile in profiles], ["2026年3月", "2026年Q1"])
        self.assertIsInstance(profiles[0], SavedBatchProfile)
        self.assertEqual(profiles[1].config.workflow_mode, WorkflowMode.MERGED)

    def test_delete_batch_profile_removes_saved_file(self) -> None:
        config = GuiBatchConfig(batch_name="2026年3月")

        with tempfile.TemporaryDirectory() as temp_dir:
            batch_dir = Path(temp_dir) / "batches"
            save_batch_profile(config, batch_dir=batch_dir)

            deleted = delete_batch_profile("2026年3月", batch_dir=batch_dir)
            remaining_profiles = load_saved_batch_profiles(batch_dir=batch_dir)

        self.assertTrue(deleted)
        self.assertEqual(remaining_profiles, ())

    def test_save_and_load_gui_session_preserves_active_saved_batch_name(self) -> None:
        config = GuiBatchConfig(
            batch_name="2026年Q1",
            workflow_mode=WorkflowMode.MERGED,
            merge_input_dirs=(Path("/tmp/datas/1月"), Path("/tmp/datas/2月")),
            merge_output_dir=Path("/tmp/datas/合并结果"),
        )

        with tempfile.TemporaryDirectory() as temp_dir:
            session_path = Path(temp_dir) / "last_session.toml"

            save_gui_session(
                config,
                active_saved_batch_name="2026年Q1",
                session_path=session_path,
            )
            loaded_config, active_saved_batch_name = load_gui_session(session_path)

        self.assertEqual(loaded_config.batch_name, "2026年Q1")
        self.assertEqual(loaded_config.workflow_mode, WorkflowMode.MERGED)
        self.assertEqual(active_saved_batch_name, "2026年Q1")

    def test_batch_profile_storage_path_is_stable_for_same_batch_name(self) -> None:
        batch_dir = Path("/tmp/gui_profiles/batches")

        path_a = batch_profile_storage_path("2026年3月", batch_dir=batch_dir)
        path_b = batch_profile_storage_path("2026年3月", batch_dir=batch_dir)

        self.assertEqual(path_a, path_b)


class CustomerTypePreviewTests(unittest.TestCase):
    def create_workbook(self, path: Path, headers: list[str], rows: list[list[object]]) -> None:
        df = pd.DataFrame(rows, columns=headers)
        with pd.ExcelWriter(path, engine="openpyxl") as writer:
            df.to_excel(writer, sheet_name="问卷数据", index=False)

    def test_build_stats_preview_summary_reports_ready_and_missing_items(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            input_dir = Path(temp_dir) / "datas" / "3月"
            input_dir.mkdir(parents=True)

            self.create_workbook(
                input_dir / "展览.xlsx",
                headers=["A", "B", "C", "D", "E"],
                rows=[
                    ["x", "x", "x", "x", "展览主承办"],
                    ["x", "x", "x", "x", "参展商"],
                ],
            )
            self.create_workbook(
                input_dir / "酒店.xlsx",
                headers=["A", "B", "C"],
                rows=[
                    ["x", "x", "散客"],
                ],
            )

            config = GuiBatchConfig(
                batch_name="2026年3月",
                workflow_mode=WorkflowMode.SINGLE,
                single_input_dir=input_dir,
                stats_output_dir=Path(temp_dir) / "output",
            )

            summary = build_stats_preview_summary(config)

            by_name = {row.customer_type_name: row for row in summary.rows}
            self.assertEqual(by_name["展览主承办"].status, CustomerTypePreviewStatus.READY)
            self.assertEqual(by_name["参展商"].status, CustomerTypePreviewStatus.READY)
            self.assertEqual(by_name["专业观众"].status, CustomerTypePreviewStatus.MISSING_ROLE_DATA)
            self.assertEqual(by_name["散客"].status, CustomerTypePreviewStatus.READY)
            self.assertEqual(by_name["住宿团队"].status, CustomerTypePreviewStatus.MISSING_ROLE_DATA)
            self.assertEqual(by_name["会展服务商"].status, CustomerTypePreviewStatus.MISSING_SOURCE_FILE)
            self.assertEqual(summary.ready_count, 3)
            self.assertGreater(summary.missing_source_count, 0)
            self.assertGreater(summary.missing_role_count, 0)

    def test_build_stats_preview_summary_uses_merged_output_dir_in_merged_mode(self) -> None:
        with tempfile.TemporaryDirectory() as temp_dir:
            single_dir = Path(temp_dir) / "datas" / "3月"
            merged_dir = Path(temp_dir) / "datas" / "合并结果"
            single_dir.mkdir(parents=True)
            merged_dir.mkdir(parents=True)

            self.create_workbook(
                merged_dir / "展览.xlsx",
                headers=["A", "B", "C", "D", "E"],
                rows=[["x", "x", "x", "x", "展览主承办"]],
            )

            config = GuiBatchConfig(
                batch_name="2026年Q1",
                workflow_mode=WorkflowMode.MERGED,
                single_input_dir=single_dir,
                merge_output_dir=merged_dir,
            )

            summary = build_stats_preview_summary(config)

            by_name = {row.customer_type_name: row for row in summary.rows}
            self.assertEqual(by_name["展览主承办"].status, CustomerTypePreviewStatus.READY)

    def test_default_selected_customer_types_only_includes_ready_rows(self) -> None:
        summary = StatsPreviewSummary(
            rows=(
                CustomerTypePreviewRow(
                    template_name="organizer",
                    customer_type_name="展览主承办",
                    document_display_name=None,
                    source_file_name="展览.xlsx",
                    output_name="展览主承办.xlsx",
                    status=CustomerTypePreviewStatus.READY,
                    detail="ok",
                ),
                CustomerTypePreviewRow(
                    template_name="visitor",
                    customer_type_name="专业观众",
                    document_display_name=None,
                    source_file_name="展览.xlsx",
                    output_name="专业观众.xlsx",
                    status=CustomerTypePreviewStatus.MISSING_ROLE_DATA,
                    detail="missing",
                ),
            ),
            input_dir=Path("/tmp/datas/3月"),
        )

        selected = default_selected_customer_types(summary)

        self.assertEqual(selected, frozenset({"展览主承办"}))

    def test_ordered_selected_customer_types_follows_preview_row_order(self) -> None:
        rows = (
            CustomerTypePreviewRow(
                template_name="organizer",
                customer_type_name="展览主承办",
                document_display_name=None,
                source_file_name="展览.xlsx",
                output_name="展览主承办.xlsx",
                status=CustomerTypePreviewStatus.READY,
                detail="ok",
            ),
            CustomerTypePreviewRow(
                template_name="hotel_individual_guest",
                customer_type_name="散客",
                document_display_name=None,
                source_file_name="酒店.xlsx",
                output_name="散客.xlsx",
                status=CustomerTypePreviewStatus.READY,
                detail="ok",
            ),
            CustomerTypePreviewRow(
                template_name="exhibitor",
                customer_type_name="参展商",
                document_display_name=None,
                source_file_name="展览.xlsx",
                output_name="参展商.xlsx",
                status=CustomerTypePreviewStatus.READY,
                detail="ok",
            ),
        )

        ordered = ordered_selected_customer_types(rows, frozenset({"参展商", "展览主承办"}))

        self.assertEqual(ordered, ("展览主承办", "参展商"))

    def test_build_stats_preview_summary_text_includes_selected_count(self) -> None:
        summary = StatsPreviewSummary(
            rows=(
                CustomerTypePreviewRow(
                    template_name="organizer",
                    customer_type_name="展览主承办",
                    document_display_name=None,
                    source_file_name="展览.xlsx",
                    output_name="展览主承办.xlsx",
                    status=CustomerTypePreviewStatus.READY,
                    detail="ok",
                ),
                CustomerTypePreviewRow(
                    template_name="visitor",
                    customer_type_name="专业观众",
                    document_display_name=None,
                    source_file_name="展览.xlsx",
                    output_name="专业观众.xlsx",
                    status=CustomerTypePreviewStatus.MISSING_ROLE_DATA,
                    detail="missing",
                ),
            ),
            input_dir=Path("/tmp/datas/3月"),
        )

        summary_text = build_stats_preview_summary_text(summary, frozenset({"展览主承办"}))

        self.assertIn("已选 1 个", summary_text)
        self.assertIn("可生成 1 个", summary_text)


class WorkflowRunControllerTests(unittest.TestCase):
    def test_controller_disables_start_and_enables_terminate_while_running(self) -> None:
        controller = WorkflowRunController()

        controller.begin(("survey_stats", "summary_table", "generate_ppt"))
        controller.mark_task_started("survey_stats")

        self.assertEqual(controller.status, WorkflowRunStatus.RUNNING)
        self.assertFalse(controller.start_enabled)
        self.assertTrue(controller.terminate_enabled)
        self.assertEqual(controller.active_step_key, "survey_stats")

    def test_controller_transitions_to_cancelled_after_terminate_request(self) -> None:
        controller = WorkflowRunController()

        controller.begin(("survey_stats", "summary_table"))
        controller.mark_task_started("survey_stats")
        controller.request_cancel()
        controller.finish_run(False)

        self.assertEqual(controller.status, WorkflowRunStatus.CANCELLED)
        self.assertTrue(controller.start_enabled)
        self.assertFalse(controller.terminate_enabled)
        self.assertTrue(controller.cancel_requested)

    def test_controller_transitions_to_failed_without_cancel_request(self) -> None:
        controller = WorkflowRunController()

        controller.begin(("survey_stats", "summary_table"))
        controller.mark_task_started("summary_table")
        controller.mark_task_finished("summary_table", False)
        controller.finish_run(False)

        self.assertEqual(controller.status, WorkflowRunStatus.FAILED)
        self.assertEqual(controller.failed_step_key, "summary_table")

    def test_build_workflow_status_text_reports_running_and_failed_titles(self) -> None:
        controller = WorkflowRunController()
        controller.begin(("survey_stats", "summary_table"))
        controller.mark_task_started("survey_stats")

        running_text = build_workflow_status_text(
            controller,
            {
                "survey_stats": "生成分项统计",
                "summary_table": "生成汇总表",
            },
        )

        self.assertEqual(running_text, "执行中 1/2：生成分项统计")

        controller.mark_task_finished("survey_stats", False)
        controller.finish_run(False)

        failed_text = build_workflow_status_text(
            controller,
            {
                "survey_stats": "生成分项统计",
                "summary_table": "生成汇总表",
            },
        )

        self.assertEqual(failed_text, "执行失败：生成分项统计")


if __name__ == "__main__":
    unittest.main()
