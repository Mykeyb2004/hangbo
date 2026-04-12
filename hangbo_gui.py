from __future__ import annotations

import hashlib
import math
import queue
import shlex
import subprocess
import sys
import tempfile
import threading
import tomllib
from dataclasses import dataclass, replace
from enum import Enum
from pathlib import Path
from tkinter import filedialog, messagebox, scrolledtext
import tkinter as tk
from tkinter import ttk

from pptx import Presentation
from summary_table import SUMMARY_ROW_DEFINITIONS
from survey_customer_mappings import STANDARD_CUSTOMER_TYPE_MAPPINGS
from survey_stats import (
    BatchConfig,
    DIRECTORY_NOTICE_REASON_MISSING_ROLE_DATA,
    DIRECTORY_NOTICE_REASON_MISSING_SOURCE_FILE,
    discover_directory_jobs,
)

PROJECT_ROOT = Path(__file__).resolve().parent
LOG_DIR = PROJECT_ROOT / "logs"
GUI_RUNTIME_DIR = LOG_DIR / "gui_runtime"
GUI_THUMBNAIL_DIR = GUI_RUNTIME_DIR / "ppt_thumbnails"
GUI_PROFILE_DIR = LOG_DIR / "gui_profiles"
GUI_BATCH_DIR = GUI_PROFILE_DIR / "batches"
GUI_SESSION_PATH = GUI_PROFILE_DIR / "last_session.toml"
DEFAULT_SHEET_NAME = "问卷数据"
PHASE_PREPROCESS_SECTION_TITLE = "兼容新版调查问卷数据结构"
PHASE_PREPROCESS_DESCRIPTION = (
    "说明：如果新版调查问卷在第三列增加了“一期/二期”等期次字段，"
    "这一步会自动把该列移到最后，避免后续统计错位。"
)
PHASE_PREPROCESS_BUTTON_TEXT = "执行兼容新版结构（phase_column_preprocess.py）"
PHASE_PREPROCESS_TASK_TITLE = "兼容新版调查问卷数据结构"
PHASE_PREPROCESS_WORKFLOW_TEXT = f"预处理：{PHASE_PREPROCESS_TASK_TITLE}"
FILL_YEAR_MONTH_SECTION_TITLE = "在数据源中加入年份+月份"
FILL_YEAR_MONTH_DESCRIPTION = "说明：给问卷数据补写“年份”“月份”两列，方便后续合并文件。"
FILL_YEAR_MONTH_BUTTON_TEXT = "执行补写年份+月份（fill_year_month_columns.py）"
FILL_YEAR_MONTH_TASK_TITLE = "在数据源中加入年份+月份"
FILL_YEAR_MONTH_WORKFLOW_TEXT = f"预处理：{FILL_YEAR_MONTH_TASK_TITLE}"
PPT_DEFAULT_FILE_PATTERN = "*.xlsx"
PPT_DEFAULT_SHEET_NAME_MODE = "first"
PPT_DEFAULT_TITLE_SUFFIX = ""
PPT_DEFAULT_MAX_SINGLE_TABLE_ROWS = "18"
PPT_DEFAULT_MAX_SPLIT_TABLE_ROWS = "19"
PPT_DEFAULT_BODY_FONT_SIZE_PT = "10.5"
PPT_DEFAULT_HEADER_FONT_SIZE_PT = "11.0"
PPT_DEFAULT_SUMMARY_FONT_SIZE_PT = "12.0"
PPT_DEFAULT_TEMPLATE_SLIDE_INDEX = "0"
PPT_DEFAULT_CHART_PLACEHOLDER_TEXT = (
    "图表分析内容待补充。后续将在此处补充该客户分组二级指标的整体解读、优势项与待提升项。"
)
PPT_DEFAULT_CHART_IMAGE_DPI = "220"
PPT_DEFAULT_LLM_ENV_PATH = ".env"
PPT_DEFAULT_LLM_SYSTEM_ROLE_PATH = "system_role.md"
PPT_DEFAULT_LLM_TARGET_CHARS = "300"
PPT_DEFAULT_LLM_TEMPERATURE = "0.4"
PPT_DEFAULT_LLM_MAX_TOKENS = "500"
PPT_DEFAULT_LLM_CHECKPOINT_CHARS = "80"
PPT_CATEGORY_LABEL_VALUES = tuple(
    dict.fromkeys(definition.category_label for definition in SUMMARY_ROW_DEFINITIONS)
)
PPT_DEFAULT_CATEGORY_INTRO_SLIDE_NUMBER = "1"
PPT_THUMBNAIL_DPI = 96
PPT_THUMBNAIL_MAX_WIDTH = 240
PPT_THUMBNAIL_MAX_HEIGHT = 135
PPT_SHEET_NAME_MODE_VALUES = ("first", "named")
PPT_LAYOUT_REGION_LABELS = (
    ("summary_table", "摘要表"),
    ("detail_single_table", "单列表"),
    ("detail_left_table", "左侧明细表"),
    ("detail_right_table", "右侧明细表"),
    ("chart_image", "图表区"),
    ("chart_textbox", "图表文字框"),
)
PPT_LAYOUT_DEFAULTS: dict[str, tuple[str, str, str, str]] = {
    "summary_table": ("0.73", "1.45", "11.87", "0.56"),
    "detail_single_table": ("0.73", "2.10", "11.87", "4.95"),
    "detail_left_table": ("0.73", "2.10", "5.78", "4.95"),
    "detail_right_table": ("6.82", "2.10", "5.78", "4.95"),
    "chart_image": ("0.78", "1.58", "5.55", "5.10"),
    "chart_textbox": ("6.55", "1.58", "5.50", "5.10"),
}


class WorkflowMode(str, Enum):
    SINGLE = "single"
    MERGED = "merged"


class CustomerTypePreviewStatus(str, Enum):
    READY = "ready"
    MISSING_SOURCE_FILE = "missing_source_file"
    MISSING_ROLE_DATA = "missing_role_data"
    MISSING_INPUT_DIR = "missing_input_dir"


class WorkflowRunStatus(str, Enum):
    IDLE = "idle"
    RUNNING = "running"
    CANCELLING = "cancelling"
    SUCCEEDED = "succeeded"
    FAILED = "failed"
    CANCELLED = "cancelled"


@dataclass(frozen=True)
class GuiBatchConfig:
    batch_name: str = "2026年3月"
    workflow_mode: WorkflowMode = WorkflowMode.SINGLE
    single_input_dir: Path = PROJECT_ROOT / "datas" / "3月"
    merge_input_dirs: tuple[Path, ...] = ()
    merge_output_dir: Path = PROJECT_ROOT / "datas" / "合并结果"
    sheet_name: str = DEFAULT_SHEET_NAME
    year_value: str = ""
    month_value: str = ""
    stats_output_dir: Path = PROJECT_ROOT / "输出结果" / "3月"
    calculation_mode: str = "template"
    output_format: str = "xlsx"
    summary_output_dir: Path = PROJECT_ROOT / "汇总结果" / "3月"
    summary_output_name: str = "3月客户类型满意度汇总表.xlsx"
    ppt_template_path: Path = PROJECT_ROOT / "templates" / "template.pptx"
    output_ppt_path: Path = PROJECT_ROOT / "输出结果" / "3月满意度报告.pptx"
    ppt_file_pattern: str = PPT_DEFAULT_FILE_PATTERN
    ppt_sheet_name_mode: str = PPT_DEFAULT_SHEET_NAME_MODE
    ppt_sheet_name: str = ""
    ppt_section_mode: str = "auto"
    ppt_blank_display: str = ""
    ppt_title_suffix: str = PPT_DEFAULT_TITLE_SUFFIX
    ppt_max_single_table_rows: str = PPT_DEFAULT_MAX_SINGLE_TABLE_ROWS
    ppt_max_split_table_rows: str = PPT_DEFAULT_MAX_SPLIT_TABLE_ROWS
    ppt_sort_files: bool = True
    ppt_body_font_size_pt: str = PPT_DEFAULT_BODY_FONT_SIZE_PT
    ppt_header_font_size_pt: str = PPT_DEFAULT_HEADER_FONT_SIZE_PT
    ppt_summary_font_size_pt: str = PPT_DEFAULT_SUMMARY_FONT_SIZE_PT
    ppt_template_slide_index: str = PPT_DEFAULT_TEMPLATE_SLIDE_INDEX
    ppt_chart_page_enabled: bool = False
    ppt_chart_placeholder_text: str = PPT_DEFAULT_CHART_PLACEHOLDER_TEXT
    ppt_chart_image_dpi: str = PPT_DEFAULT_CHART_IMAGE_DPI
    ppt_llm_notes_enabled: bool = False
    ppt_llm_env_path: str = PPT_DEFAULT_LLM_ENV_PATH
    ppt_llm_system_role_path: str = PPT_DEFAULT_LLM_SYSTEM_ROLE_PATH
    ppt_llm_target_chars: str = PPT_DEFAULT_LLM_TARGET_CHARS
    ppt_llm_temperature: str = PPT_DEFAULT_LLM_TEMPERATURE
    ppt_llm_max_tokens: str = PPT_DEFAULT_LLM_MAX_TOKENS
    ppt_llm_checkpoint_chars: str = PPT_DEFAULT_LLM_CHECKPOINT_CHARS
    ppt_category_intro_slides_text: str = ""
    ppt_layout_summary_table_left: str = PPT_LAYOUT_DEFAULTS["summary_table"][0]
    ppt_layout_summary_table_top: str = PPT_LAYOUT_DEFAULTS["summary_table"][1]
    ppt_layout_summary_table_width: str = PPT_LAYOUT_DEFAULTS["summary_table"][2]
    ppt_layout_summary_table_height: str = PPT_LAYOUT_DEFAULTS["summary_table"][3]
    ppt_layout_detail_single_table_left: str = PPT_LAYOUT_DEFAULTS["detail_single_table"][0]
    ppt_layout_detail_single_table_top: str = PPT_LAYOUT_DEFAULTS["detail_single_table"][1]
    ppt_layout_detail_single_table_width: str = PPT_LAYOUT_DEFAULTS["detail_single_table"][2]
    ppt_layout_detail_single_table_height: str = PPT_LAYOUT_DEFAULTS["detail_single_table"][3]
    ppt_layout_detail_left_table_left: str = PPT_LAYOUT_DEFAULTS["detail_left_table"][0]
    ppt_layout_detail_left_table_top: str = PPT_LAYOUT_DEFAULTS["detail_left_table"][1]
    ppt_layout_detail_left_table_width: str = PPT_LAYOUT_DEFAULTS["detail_left_table"][2]
    ppt_layout_detail_left_table_height: str = PPT_LAYOUT_DEFAULTS["detail_left_table"][3]
    ppt_layout_detail_right_table_left: str = PPT_LAYOUT_DEFAULTS["detail_right_table"][0]
    ppt_layout_detail_right_table_top: str = PPT_LAYOUT_DEFAULTS["detail_right_table"][1]
    ppt_layout_detail_right_table_width: str = PPT_LAYOUT_DEFAULTS["detail_right_table"][2]
    ppt_layout_detail_right_table_height: str = PPT_LAYOUT_DEFAULTS["detail_right_table"][3]
    ppt_layout_chart_image_left: str = PPT_LAYOUT_DEFAULTS["chart_image"][0]
    ppt_layout_chart_image_top: str = PPT_LAYOUT_DEFAULTS["chart_image"][1]
    ppt_layout_chart_image_width: str = PPT_LAYOUT_DEFAULTS["chart_image"][2]
    ppt_layout_chart_image_height: str = PPT_LAYOUT_DEFAULTS["chart_image"][3]
    ppt_layout_chart_textbox_left: str = PPT_LAYOUT_DEFAULTS["chart_textbox"][0]
    ppt_layout_chart_textbox_top: str = PPT_LAYOUT_DEFAULTS["chart_textbox"][1]
    ppt_layout_chart_textbox_width: str = PPT_LAYOUT_DEFAULTS["chart_textbox"][2]
    ppt_layout_chart_textbox_height: str = PPT_LAYOUT_DEFAULTS["chart_textbox"][3]

    def effective_input_dir(self) -> Path:
        if self.workflow_mode == WorkflowMode.MERGED:
            return self.merge_output_dir
        return self.single_input_dir


DEFAULT_GUI_BATCH_CONFIG = GuiBatchConfig()


@dataclass(frozen=True)
class MainWorkflowSelection:
    include_merge: bool = False
    include_phase_preprocess: bool = True
    include_fill_year_month: bool = False
    include_survey_stats: bool = True
    include_summary: bool = True
    include_ppt: bool = True


@dataclass(frozen=True)
class TaskCommand:
    key: str
    title: str
    command: tuple[str, ...]


@dataclass(frozen=True)
class PptSlidePreview:
    slide_number: int
    title: str
    label: str


@dataclass(frozen=True)
class CustomerTypePreviewRow:
    template_name: str
    customer_type_name: str
    document_display_name: str | None
    source_file_name: str
    output_name: str
    status: CustomerTypePreviewStatus
    detail: str
    note: str | None = None


@dataclass(frozen=True)
class StatsPreviewSummary:
    rows: tuple[CustomerTypePreviewRow, ...]
    input_dir: Path

    @property
    def ready_count(self) -> int:
        return sum(1 for row in self.rows if row.status == CustomerTypePreviewStatus.READY)

    @property
    def missing_source_count(self) -> int:
        return sum(
            1 for row in self.rows if row.status == CustomerTypePreviewStatus.MISSING_SOURCE_FILE
        )

    @property
    def missing_role_count(self) -> int:
        return sum(
            1 for row in self.rows if row.status == CustomerTypePreviewStatus.MISSING_ROLE_DATA
        )


@dataclass(frozen=True)
class SavedBatchProfile:
    batch_name: str
    path: Path
    config: GuiBatchConfig


def default_selected_customer_types(summary: StatsPreviewSummary) -> frozenset[str]:
    return frozenset(
        row.customer_type_name
        for row in summary.rows
        if row.status == CustomerTypePreviewStatus.READY
    )


def page_key_for_task(task_key: str) -> str:
    return {
        "merge_workbooks": "data_source",
        "phase_preprocess": "preprocess",
        "fill_year_month": "preprocess",
        "survey_stats": "stats",
        "summary_table": "summary",
        "generate_ppt": "ppt",
    }.get(task_key, "dashboard")


def customer_type_preview_status_text(status: CustomerTypePreviewStatus) -> str:
    return {
        CustomerTypePreviewStatus.READY: "可生成",
        CustomerTypePreviewStatus.MISSING_SOURCE_FILE: "缺少来源文件",
        CustomerTypePreviewStatus.MISSING_ROLE_DATA: "未匹配身份值",
        CustomerTypePreviewStatus.MISSING_INPUT_DIR: "输入目录不存在",
    }[status]


def ordered_selected_customer_types(
    rows: tuple[CustomerTypePreviewRow, ...],
    selected_customer_types: frozenset[str],
) -> tuple[str, ...]:
    return tuple(
        row.customer_type_name
        for row in rows
        if row.customer_type_name in selected_customer_types
    )


def build_stats_preview_summary_text(
    summary: StatsPreviewSummary,
    selected_customer_types: frozenset[str],
) -> str:
    selected_count = len(ordered_selected_customer_types(summary.rows, selected_customer_types))
    return (
        f"输入目录：{summary.input_dir}；"
        f"已选 {selected_count} 个；"
        f"可生成 {summary.ready_count} 个；"
        f"缺少来源 {summary.missing_source_count} 个；"
        f"缺少身份值 {summary.missing_role_count} 个"
    )


def build_workflow_status_text(
    controller: WorkflowRunController,
    task_title_lookup: dict[str, str],
) -> str:
    total = len(controller.planned_step_keys)
    completed = len(controller.completed_step_keys)
    active_key = controller.active_step_key
    active_title = task_title_lookup.get(active_key or "", active_key or "等待启动")
    failed_key = controller.failed_step_key
    failed_title = task_title_lookup.get(failed_key or "", failed_key or "未知步骤")

    if controller.status == WorkflowRunStatus.IDLE:
        return "未开始"
    if controller.status == WorkflowRunStatus.RUNNING:
        return f"执行中 {completed + 1}/{max(total, 1)}：{active_title}"
    if controller.status == WorkflowRunStatus.CANCELLING:
        return f"正在终止：{active_title}"
    if controller.status == WorkflowRunStatus.SUCCEEDED:
        return "执行完成"
    if controller.status == WorkflowRunStatus.CANCELLED:
        return "已取消"
    return f"执行失败：{failed_title}"


def bool_to_toml(value: bool) -> str:
    return "true" if value else "false"


def layout_field_key(region_name: str, field_name: str) -> str:
    return f"ppt_layout_{region_name}_{field_name}"


def build_gui_ppt_layout_lines(config: GuiBatchConfig) -> list[str]:
    lines: list[str] = []
    for region_name, _ in PPT_LAYOUT_REGION_LABELS:
        for field_name in ("left", "top", "width", "height"):
            config_key = layout_field_key(region_name, field_name)
            lines.append(
                f"{config_key} = {toml_quote(getattr(config, config_key))}"
            )
    return lines


def build_gui_batch_config_text(
    config: GuiBatchConfig,
    *,
    active_saved_batch_name: str | None = None,
) -> str:
    lines = [
        f"batch_name = {toml_quote(config.batch_name)}",
        f'workflow_mode = "{config.workflow_mode.value}"',
        f"single_input_dir = {toml_quote(config.single_input_dir)}",
        f"merge_output_dir = {toml_quote(config.merge_output_dir)}",
        f"sheet_name = {toml_quote(config.sheet_name)}",
        f"year_value = {toml_quote(config.year_value)}",
        f"month_value = {toml_quote(config.month_value)}",
        f"stats_output_dir = {toml_quote(config.stats_output_dir)}",
        f'calculation_mode = "{config.calculation_mode}"',
        f'output_format = "{config.output_format}"',
        f"summary_output_dir = {toml_quote(config.summary_output_dir)}",
        f"summary_output_name = {toml_quote(config.summary_output_name)}",
        f"ppt_template_path = {toml_quote(config.ppt_template_path)}",
        f"output_ppt_path = {toml_quote(config.output_ppt_path)}",
        f"ppt_file_pattern = {toml_quote(config.ppt_file_pattern)}",
        f'ppt_sheet_name_mode = "{config.ppt_sheet_name_mode}"',
        f"ppt_sheet_name = {toml_quote(config.ppt_sheet_name)}",
        f'ppt_section_mode = "{config.ppt_section_mode}"',
        f"ppt_blank_display = {toml_quote(config.ppt_blank_display)}",
        f"ppt_title_suffix = {toml_quote(config.ppt_title_suffix)}",
        f"ppt_max_single_table_rows = {toml_quote(config.ppt_max_single_table_rows)}",
        f"ppt_max_split_table_rows = {toml_quote(config.ppt_max_split_table_rows)}",
        f"ppt_sort_files = {bool_to_toml(config.ppt_sort_files)}",
        f"ppt_body_font_size_pt = {toml_quote(config.ppt_body_font_size_pt)}",
        f"ppt_header_font_size_pt = {toml_quote(config.ppt_header_font_size_pt)}",
        f"ppt_summary_font_size_pt = {toml_quote(config.ppt_summary_font_size_pt)}",
        f"ppt_template_slide_index = {toml_quote(config.ppt_template_slide_index)}",
        f"ppt_chart_page_enabled = {bool_to_toml(config.ppt_chart_page_enabled)}",
        f"ppt_chart_placeholder_text = {toml_quote(config.ppt_chart_placeholder_text)}",
        f"ppt_chart_image_dpi = {toml_quote(config.ppt_chart_image_dpi)}",
        f"ppt_llm_notes_enabled = {bool_to_toml(config.ppt_llm_notes_enabled)}",
        f"ppt_llm_env_path = {toml_quote(config.ppt_llm_env_path)}",
        f"ppt_llm_system_role_path = {toml_quote(config.ppt_llm_system_role_path)}",
        f"ppt_llm_target_chars = {toml_quote(config.ppt_llm_target_chars)}",
        f"ppt_llm_temperature = {toml_quote(config.ppt_llm_temperature)}",
        f"ppt_llm_max_tokens = {toml_quote(config.ppt_llm_max_tokens)}",
        f"ppt_llm_checkpoint_chars = {toml_quote(config.ppt_llm_checkpoint_chars)}",
        f"ppt_category_intro_slides_text = {toml_quote(config.ppt_category_intro_slides_text)}",
    ]
    lines.extend(build_gui_ppt_layout_lines(config))
    if active_saved_batch_name is not None:
        lines.append(f"active_saved_batch_name = {toml_quote(active_saved_batch_name)}")
    if config.merge_input_dirs:
        lines.append("")
        lines.append("merge_input_dirs = [")
        for path in config.merge_input_dirs:
            lines.append(f"  {toml_quote(path)},")
        lines.append("]")
    return "\n".join(lines) + "\n"


def _string_value(raw_data: dict[str, object], key: str, default: str) -> str:
    value = raw_data.get(key)
    if value is None:
        return default
    return str(value)


def _path_value(raw_data: dict[str, object], key: str, default: Path) -> Path:
    value = raw_data.get(key)
    if value is None:
        return default
    return Path(str(value)).resolve()


def _bool_value(raw_data: dict[str, object], key: str, default: bool) -> bool:
    value = raw_data.get(key)
    if value is None:
        return default
    if isinstance(value, bool):
        return value
    return str(value).strip().lower() in {"1", "true", "yes", "on"}


def parse_gui_batch_config(raw_data: dict[str, object]) -> GuiBatchConfig:
    raw_merge_input_dirs = raw_data.get("merge_input_dirs", [])
    if raw_merge_input_dirs is None:
        raw_merge_input_dirs = []
    if not isinstance(raw_merge_input_dirs, list):
        raise ValueError("merge_input_dirs 必须是列表。")
    merge_input_dirs = tuple(Path(str(value)).resolve() for value in raw_merge_input_dirs)

    raw_workflow_mode = _string_value(
        raw_data,
        "workflow_mode",
        DEFAULT_GUI_BATCH_CONFIG.workflow_mode.value,
    ).strip() or DEFAULT_GUI_BATCH_CONFIG.workflow_mode.value

    return GuiBatchConfig(
        batch_name=_string_value(raw_data, "batch_name", DEFAULT_GUI_BATCH_CONFIG.batch_name).strip()
        or DEFAULT_GUI_BATCH_CONFIG.batch_name,
        workflow_mode=WorkflowMode(raw_workflow_mode),
        single_input_dir=_path_value(
            raw_data,
            "single_input_dir",
            DEFAULT_GUI_BATCH_CONFIG.single_input_dir,
        ),
        merge_input_dirs=merge_input_dirs,
        merge_output_dir=_path_value(
            raw_data,
            "merge_output_dir",
            DEFAULT_GUI_BATCH_CONFIG.merge_output_dir,
        ),
        sheet_name=_string_value(raw_data, "sheet_name", DEFAULT_GUI_BATCH_CONFIG.sheet_name)
        or DEFAULT_GUI_BATCH_CONFIG.sheet_name,
        year_value=_string_value(raw_data, "year_value", DEFAULT_GUI_BATCH_CONFIG.year_value),
        month_value=_string_value(raw_data, "month_value", DEFAULT_GUI_BATCH_CONFIG.month_value),
        stats_output_dir=_path_value(
            raw_data,
            "stats_output_dir",
            DEFAULT_GUI_BATCH_CONFIG.stats_output_dir,
        ),
        calculation_mode=_string_value(
            raw_data,
            "calculation_mode",
            DEFAULT_GUI_BATCH_CONFIG.calculation_mode,
        )
        or DEFAULT_GUI_BATCH_CONFIG.calculation_mode,
        output_format=_string_value(raw_data, "output_format", DEFAULT_GUI_BATCH_CONFIG.output_format)
        or DEFAULT_GUI_BATCH_CONFIG.output_format,
        summary_output_dir=_path_value(
            raw_data,
            "summary_output_dir",
            DEFAULT_GUI_BATCH_CONFIG.summary_output_dir,
        ),
        summary_output_name=_string_value(
            raw_data,
            "summary_output_name",
            DEFAULT_GUI_BATCH_CONFIG.summary_output_name,
        )
        or DEFAULT_GUI_BATCH_CONFIG.summary_output_name,
        ppt_template_path=_path_value(
            raw_data,
            "ppt_template_path",
            DEFAULT_GUI_BATCH_CONFIG.ppt_template_path,
        ),
        output_ppt_path=_path_value(
            raw_data,
            "output_ppt_path",
            DEFAULT_GUI_BATCH_CONFIG.output_ppt_path,
        ),
        ppt_file_pattern=_string_value(
            raw_data,
            "ppt_file_pattern",
            DEFAULT_GUI_BATCH_CONFIG.ppt_file_pattern,
        )
        or DEFAULT_GUI_BATCH_CONFIG.ppt_file_pattern,
        ppt_sheet_name_mode=_string_value(
            raw_data,
            "ppt_sheet_name_mode",
            DEFAULT_GUI_BATCH_CONFIG.ppt_sheet_name_mode,
        )
        or DEFAULT_GUI_BATCH_CONFIG.ppt_sheet_name_mode,
        ppt_sheet_name=_string_value(
            raw_data,
            "ppt_sheet_name",
            DEFAULT_GUI_BATCH_CONFIG.ppt_sheet_name,
        ),
        ppt_section_mode=_string_value(
            raw_data,
            "ppt_section_mode",
            DEFAULT_GUI_BATCH_CONFIG.ppt_section_mode,
        )
        or DEFAULT_GUI_BATCH_CONFIG.ppt_section_mode,
        ppt_blank_display=_string_value(
            raw_data,
            "ppt_blank_display",
            DEFAULT_GUI_BATCH_CONFIG.ppt_blank_display,
        ),
        ppt_title_suffix=_string_value(
            raw_data,
            "ppt_title_suffix",
            DEFAULT_GUI_BATCH_CONFIG.ppt_title_suffix,
        ),
        ppt_max_single_table_rows=_string_value(
            raw_data,
            "ppt_max_single_table_rows",
            DEFAULT_GUI_BATCH_CONFIG.ppt_max_single_table_rows,
        ),
        ppt_max_split_table_rows=_string_value(
            raw_data,
            "ppt_max_split_table_rows",
            DEFAULT_GUI_BATCH_CONFIG.ppt_max_split_table_rows,
        ),
        ppt_sort_files=_bool_value(
            raw_data,
            "ppt_sort_files",
            DEFAULT_GUI_BATCH_CONFIG.ppt_sort_files,
        ),
        ppt_body_font_size_pt=_string_value(
            raw_data,
            "ppt_body_font_size_pt",
            DEFAULT_GUI_BATCH_CONFIG.ppt_body_font_size_pt,
        ),
        ppt_header_font_size_pt=_string_value(
            raw_data,
            "ppt_header_font_size_pt",
            DEFAULT_GUI_BATCH_CONFIG.ppt_header_font_size_pt,
        ),
        ppt_summary_font_size_pt=_string_value(
            raw_data,
            "ppt_summary_font_size_pt",
            DEFAULT_GUI_BATCH_CONFIG.ppt_summary_font_size_pt,
        ),
        ppt_template_slide_index=_string_value(
            raw_data,
            "ppt_template_slide_index",
            DEFAULT_GUI_BATCH_CONFIG.ppt_template_slide_index,
        ),
        ppt_chart_page_enabled=_bool_value(
            raw_data,
            "ppt_chart_page_enabled",
            DEFAULT_GUI_BATCH_CONFIG.ppt_chart_page_enabled,
        ),
        ppt_chart_placeholder_text=_string_value(
            raw_data,
            "ppt_chart_placeholder_text",
            DEFAULT_GUI_BATCH_CONFIG.ppt_chart_placeholder_text,
        ),
        ppt_chart_image_dpi=_string_value(
            raw_data,
            "ppt_chart_image_dpi",
            DEFAULT_GUI_BATCH_CONFIG.ppt_chart_image_dpi,
        ),
        ppt_llm_notes_enabled=_bool_value(
            raw_data,
            "ppt_llm_notes_enabled",
            DEFAULT_GUI_BATCH_CONFIG.ppt_llm_notes_enabled,
        ),
        ppt_llm_env_path=_string_value(
            raw_data,
            "ppt_llm_env_path",
            DEFAULT_GUI_BATCH_CONFIG.ppt_llm_env_path,
        ),
        ppt_llm_system_role_path=_string_value(
            raw_data,
            "ppt_llm_system_role_path",
            DEFAULT_GUI_BATCH_CONFIG.ppt_llm_system_role_path,
        ),
        ppt_llm_target_chars=_string_value(
            raw_data,
            "ppt_llm_target_chars",
            DEFAULT_GUI_BATCH_CONFIG.ppt_llm_target_chars,
        ),
        ppt_llm_temperature=_string_value(
            raw_data,
            "ppt_llm_temperature",
            DEFAULT_GUI_BATCH_CONFIG.ppt_llm_temperature,
        ),
        ppt_llm_max_tokens=_string_value(
            raw_data,
            "ppt_llm_max_tokens",
            DEFAULT_GUI_BATCH_CONFIG.ppt_llm_max_tokens,
        ),
        ppt_llm_checkpoint_chars=_string_value(
            raw_data,
            "ppt_llm_checkpoint_chars",
            DEFAULT_GUI_BATCH_CONFIG.ppt_llm_checkpoint_chars,
        ),
        ppt_category_intro_slides_text=_string_value(
            raw_data,
            "ppt_category_intro_slides_text",
            DEFAULT_GUI_BATCH_CONFIG.ppt_category_intro_slides_text,
        ),
        ppt_layout_summary_table_left=_string_value(
            raw_data,
            "ppt_layout_summary_table_left",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_summary_table_left,
        ),
        ppt_layout_summary_table_top=_string_value(
            raw_data,
            "ppt_layout_summary_table_top",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_summary_table_top,
        ),
        ppt_layout_summary_table_width=_string_value(
            raw_data,
            "ppt_layout_summary_table_width",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_summary_table_width,
        ),
        ppt_layout_summary_table_height=_string_value(
            raw_data,
            "ppt_layout_summary_table_height",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_summary_table_height,
        ),
        ppt_layout_detail_single_table_left=_string_value(
            raw_data,
            "ppt_layout_detail_single_table_left",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_detail_single_table_left,
        ),
        ppt_layout_detail_single_table_top=_string_value(
            raw_data,
            "ppt_layout_detail_single_table_top",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_detail_single_table_top,
        ),
        ppt_layout_detail_single_table_width=_string_value(
            raw_data,
            "ppt_layout_detail_single_table_width",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_detail_single_table_width,
        ),
        ppt_layout_detail_single_table_height=_string_value(
            raw_data,
            "ppt_layout_detail_single_table_height",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_detail_single_table_height,
        ),
        ppt_layout_detail_left_table_left=_string_value(
            raw_data,
            "ppt_layout_detail_left_table_left",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_detail_left_table_left,
        ),
        ppt_layout_detail_left_table_top=_string_value(
            raw_data,
            "ppt_layout_detail_left_table_top",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_detail_left_table_top,
        ),
        ppt_layout_detail_left_table_width=_string_value(
            raw_data,
            "ppt_layout_detail_left_table_width",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_detail_left_table_width,
        ),
        ppt_layout_detail_left_table_height=_string_value(
            raw_data,
            "ppt_layout_detail_left_table_height",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_detail_left_table_height,
        ),
        ppt_layout_detail_right_table_left=_string_value(
            raw_data,
            "ppt_layout_detail_right_table_left",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_detail_right_table_left,
        ),
        ppt_layout_detail_right_table_top=_string_value(
            raw_data,
            "ppt_layout_detail_right_table_top",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_detail_right_table_top,
        ),
        ppt_layout_detail_right_table_width=_string_value(
            raw_data,
            "ppt_layout_detail_right_table_width",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_detail_right_table_width,
        ),
        ppt_layout_detail_right_table_height=_string_value(
            raw_data,
            "ppt_layout_detail_right_table_height",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_detail_right_table_height,
        ),
        ppt_layout_chart_image_left=_string_value(
            raw_data,
            "ppt_layout_chart_image_left",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_chart_image_left,
        ),
        ppt_layout_chart_image_top=_string_value(
            raw_data,
            "ppt_layout_chart_image_top",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_chart_image_top,
        ),
        ppt_layout_chart_image_width=_string_value(
            raw_data,
            "ppt_layout_chart_image_width",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_chart_image_width,
        ),
        ppt_layout_chart_image_height=_string_value(
            raw_data,
            "ppt_layout_chart_image_height",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_chart_image_height,
        ),
        ppt_layout_chart_textbox_left=_string_value(
            raw_data,
            "ppt_layout_chart_textbox_left",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_chart_textbox_left,
        ),
        ppt_layout_chart_textbox_top=_string_value(
            raw_data,
            "ppt_layout_chart_textbox_top",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_chart_textbox_top,
        ),
        ppt_layout_chart_textbox_width=_string_value(
            raw_data,
            "ppt_layout_chart_textbox_width",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_chart_textbox_width,
        ),
        ppt_layout_chart_textbox_height=_string_value(
            raw_data,
            "ppt_layout_chart_textbox_height",
            DEFAULT_GUI_BATCH_CONFIG.ppt_layout_chart_textbox_height,
        ),
    )


def load_gui_batch_config(profile_path: Path) -> GuiBatchConfig:
    raw_data = tomllib.loads(profile_path.read_text(encoding="utf-8"))
    return parse_gui_batch_config(raw_data)


def load_gui_session(session_path: Path = GUI_SESSION_PATH) -> tuple[GuiBatchConfig, str | None]:
    raw_data = tomllib.loads(session_path.read_text(encoding="utf-8"))
    config = parse_gui_batch_config(raw_data)
    active_saved_batch_name = str(raw_data.get("active_saved_batch_name", "")).strip() or None
    return config, active_saved_batch_name


def batch_profile_storage_path(
    batch_name: str,
    batch_dir: Path = GUI_BATCH_DIR,
) -> Path:
    encoded_name = batch_name.encode("utf-8").hex()
    return batch_dir / f"batch_{encoded_name}.toml"


def save_batch_profile(
    config: GuiBatchConfig,
    *,
    batch_dir: Path = GUI_BATCH_DIR,
) -> Path:
    batch_dir.mkdir(parents=True, exist_ok=True)
    profile_path = batch_profile_storage_path(config.batch_name, batch_dir)
    profile_path.write_text(build_gui_batch_config_text(config), encoding="utf-8")
    return profile_path


def load_saved_batch_profiles(
    batch_dir: Path = GUI_BATCH_DIR,
) -> tuple[SavedBatchProfile, ...]:
    if not batch_dir.exists():
        return ()

    profiles: list[SavedBatchProfile] = []
    for profile_path in sorted(batch_dir.glob("batch_*.toml")):
        if not profile_path.is_file():
            continue
        try:
            config = load_gui_batch_config(profile_path)
        except Exception:
            continue
        profiles.append(
            SavedBatchProfile(
                batch_name=config.batch_name,
                path=profile_path,
                config=config,
            )
        )
    return tuple(sorted(profiles, key=lambda profile: profile.batch_name))


def delete_batch_profile(
    batch_name: str,
    *,
    batch_dir: Path = GUI_BATCH_DIR,
) -> bool:
    profile_path = batch_profile_storage_path(batch_name, batch_dir)
    if not profile_path.exists():
        return False
    profile_path.unlink()
    return True


def save_gui_session(
    config: GuiBatchConfig,
    *,
    active_saved_batch_name: str | None,
    session_path: Path = GUI_SESSION_PATH,
) -> Path:
    session_path.parent.mkdir(parents=True, exist_ok=True)
    session_path.write_text(
        build_gui_batch_config_text(
            config,
            active_saved_batch_name=active_saved_batch_name or "",
        ),
        encoding="utf-8",
    )
    return session_path


class WorkflowRunController:
    def __init__(self) -> None:
        self.status = WorkflowRunStatus.IDLE
        self.planned_step_keys: tuple[str, ...] = ()
        self.active_step_key: str | None = None
        self.completed_step_keys: tuple[str, ...] = ()
        self.failed_step_key: str | None = None
        self.cancel_requested = False

    @property
    def start_enabled(self) -> bool:
        return self.status in {
            WorkflowRunStatus.IDLE,
            WorkflowRunStatus.SUCCEEDED,
            WorkflowRunStatus.FAILED,
            WorkflowRunStatus.CANCELLED,
        }

    @property
    def terminate_enabled(self) -> bool:
        return self.status == WorkflowRunStatus.RUNNING

    def begin(self, step_keys: tuple[str, ...]) -> None:
        self.status = WorkflowRunStatus.RUNNING
        self.planned_step_keys = step_keys
        self.active_step_key = None
        self.completed_step_keys = ()
        self.failed_step_key = None
        self.cancel_requested = False

    def mark_task_started(self, step_key: str) -> None:
        self.active_step_key = step_key
        if self.status != WorkflowRunStatus.CANCELLING:
            self.status = WorkflowRunStatus.RUNNING

    def mark_task_finished(self, step_key: str, success: bool) -> None:
        if success:
            if step_key not in self.completed_step_keys:
                self.completed_step_keys = (*self.completed_step_keys, step_key)
        else:
            self.failed_step_key = step_key
        self.active_step_key = None

    def request_cancel(self) -> None:
        if self.status in {WorkflowRunStatus.RUNNING, WorkflowRunStatus.CANCELLING}:
            self.cancel_requested = True
            self.status = WorkflowRunStatus.CANCELLING

    def finish_run(self, success: bool) -> None:
        self.active_step_key = None
        if self.cancel_requested and not success:
            self.status = WorkflowRunStatus.CANCELLED
            return
        if success:
            self.status = WorkflowRunStatus.SUCCEEDED
            return
        self.status = WorkflowRunStatus.FAILED


@dataclass(frozen=True)
class ThemePalette:
    background: str = "#f6f1ea"
    surface: str = "#fffaf6"
    surface_alt: str = "#f1e3d4"
    border: str = "#d7c2ad"
    primary: str = "#8b3c2c"
    primary_hover: str = "#6f2d21"
    accent: str = "#b8644c"
    text: str = "#2d241f"
    muted_text: str = "#74675d"
    success: str = "#2f7d4a"
    warning: str = "#b06a1c"


def build_directory_batch_config(config: GuiBatchConfig) -> BatchConfig:
    return BatchConfig(
        config_path=(PROJECT_ROOT / "hangbo_gui.py").resolve(),
        output_dir=config.stats_output_dir,
        output_format=config.output_format,
        calculation_mode=config.calculation_mode,
        sheet_name=config.sheet_name,
        input_dir=config.effective_input_dir(),
    )


def build_stats_preview_summary(config: GuiBatchConfig) -> StatsPreviewSummary:
    input_dir = config.effective_input_dir()
    if not input_dir.exists() or not input_dir.is_dir():
        return StatsPreviewSummary(
            rows=tuple(
                CustomerTypePreviewRow(
                    template_name=mapping.template_name,
                    customer_type_name=mapping.template_role_name,
                    document_display_name=mapping.document_display_name,
                    source_file_name=mapping.source_file_name,
                    output_name=f"{mapping.template_role_name}.xlsx",
                    status=CustomerTypePreviewStatus.MISSING_INPUT_DIR,
                    detail=f"输入目录不存在：{input_dir}",
                    note=mapping.note,
                )
                for mapping in STANDARD_CUSTOMER_TYPE_MAPPINGS
            ),
            input_dir=input_dir,
        )

    discovery = discover_directory_jobs(build_directory_batch_config(config))
    job_lookup = {job.name: job for job in discovery.jobs}
    missing_notice_lookup = {
        notice.customer_type_name: notice for notice in discovery.missing_customer_type_notices
    }
    rows: list[CustomerTypePreviewRow] = []
    for mapping in STANDARD_CUSTOMER_TYPE_MAPPINGS:
        job = job_lookup.get(mapping.template_role_name)
        if job is not None:
            rows.append(
                CustomerTypePreviewRow(
                    template_name=mapping.template_name,
                    customer_type_name=mapping.template_role_name,
                    document_display_name=mapping.document_display_name,
                    source_file_name=mapping.source_file_name,
                    output_name=f"{job.output_name}.{config.output_format}",
                    status=CustomerTypePreviewStatus.READY,
                    detail=f"来源文件：{job.path.name}",
                    note=mapping.note,
                )
            )
            continue

        notice = missing_notice_lookup.get(mapping.template_role_name)
        if notice is not None and notice.reason == DIRECTORY_NOTICE_REASON_MISSING_ROLE_DATA:
            status = CustomerTypePreviewStatus.MISSING_ROLE_DATA
            detail = f"来源文件存在，但未找到身份值：{notice.source_reference}"
        else:
            status = CustomerTypePreviewStatus.MISSING_SOURCE_FILE
            source_reference = notice.source_reference if notice is not None else mapping.source_file_name
            detail = f"缺少来源文件：{source_reference}"
        rows.append(
            CustomerTypePreviewRow(
                template_name=mapping.template_name,
                customer_type_name=mapping.template_role_name,
                document_display_name=mapping.document_display_name,
                source_file_name=mapping.source_file_name,
                output_name=f"{mapping.template_role_name}.{config.output_format}",
                status=status,
                detail=detail,
                note=mapping.note,
            )
        )

    return StatsPreviewSummary(rows=tuple(rows), input_dir=input_dir)


def build_main_workflow_step_keys(
    config: GuiBatchConfig,
    selection: MainWorkflowSelection,
) -> tuple[str, ...]:
    step_keys: list[str] = []
    if config.workflow_mode == WorkflowMode.MERGED and selection.include_merge:
        step_keys.append("merge_workbooks")
    if selection.include_phase_preprocess:
        step_keys.append("phase_preprocess")
    if selection.include_fill_year_month:
        step_keys.append("fill_year_month")
    if selection.include_survey_stats:
        step_keys.append("survey_stats")
    if selection.include_summary:
        step_keys.append("summary_table")
    if selection.include_ppt:
        step_keys.append("generate_ppt")
    return tuple(step_keys)


def toml_quote(value: str | Path) -> str:
    text = str(value)
    escaped = (
        text.replace("\\", "\\\\")
        .replace("\n", "\\n")
        .replace("\t", "\\t")
        .replace('"', '\\"')
    )
    return f'"{escaped}"'


def build_survey_stats_config_text(config: GuiBatchConfig) -> str:
    return "\n".join(
        [
            f"output_dir = {toml_quote(config.stats_output_dir)}",
            f'output_format = "{config.output_format}"',
            f'calculation_mode = "{config.calculation_mode}"',
            f"input_dir = {toml_quote(config.effective_input_dir())}",
            f"sheet_name = {toml_quote(config.sheet_name)}",
            "",
        ]
    )


def normalize_ppt_sheet_name_mode(sheet_name_mode: str) -> str:
    normalized = sheet_name_mode.strip().lower() or PPT_DEFAULT_SHEET_NAME_MODE
    if normalized not in PPT_SHEET_NAME_MODE_VALUES:
        raise ValueError(f"PPT sheet_name_mode 仅支持: {', '.join(PPT_SHEET_NAME_MODE_VALUES)}")
    return normalized


def parse_positive_int_text(value: str, field_label: str, *, minimum: int = 1) -> int:
    text = value.strip()
    if not text:
        raise ValueError(f"请填写 {field_label}。")
    try:
        parsed = int(text)
    except ValueError as exc:
        raise ValueError(f"{field_label} 必须是整数。") from exc
    if parsed < minimum:
        raise ValueError(f"{field_label} 必须大于等于 {minimum}。")
    return parsed


def parse_non_negative_int_text(value: str, field_label: str) -> int:
    return parse_positive_int_text(value, field_label, minimum=0)


def parse_positive_float_text(value: str, field_label: str, *, allow_zero: bool = False) -> float:
    text = value.strip()
    if not text:
        raise ValueError(f"请填写 {field_label}。")
    try:
        parsed = float(text)
    except ValueError as exc:
        raise ValueError(f"{field_label} 必须是数字。") from exc
    if allow_zero:
        if parsed < 0:
            raise ValueError(f"{field_label} 必须大于等于 0。")
    elif parsed <= 0:
        raise ValueError(f"{field_label} 必须大于 0。")
    return parsed


def parse_category_intro_slides_text(
    text: str,
) -> tuple[tuple[str, str, int], ...]:
    entries: list[tuple[str, str, int]] = []
    for line_number, raw_line in enumerate(text.splitlines(), start=1):
        line = raw_line.strip()
        if not line or line.startswith("#"):
            continue
        parts = [part.strip() for part in line.split("|")]
        if len(parts) != 3:
            raise ValueError(
                "章节页配置格式应为“客户大类|PPT路径|页码”，"
                f"第 {line_number} 行格式不正确。"
            )
        category_label, ppt_path, slide_number_text = parts
        if not category_label:
            raise ValueError(f"章节页配置第 {line_number} 行缺少客户大类名称。")
        if not ppt_path:
            raise ValueError(f"章节页配置第 {line_number} 行缺少 PPT 路径。")
        slide_number = parse_positive_int_text(
            slide_number_text,
            f"章节页配置第 {line_number} 行页码",
        )
        entries.append((category_label, ppt_path, slide_number))
    return tuple(entries)


def build_category_intro_slides_text(
    entries: tuple[tuple[str, str, int], ...] | list[tuple[str, str, int]],
) -> str:
    lines: list[str] = []
    for line_number, (category_label, ppt_path, slide_number) in enumerate(
        entries,
        start=1,
    ):
        normalized_category = category_label.strip()
        normalized_ppt_path = ppt_path.strip()
        if not normalized_category:
            raise ValueError(f"章节页配置第 {line_number} 行缺少客户大类名称。")
        if not normalized_ppt_path:
            raise ValueError(f"章节页配置第 {line_number} 行缺少 PPT 路径。")
        normalized_slide_number = parse_positive_int_text(
            str(slide_number),
            f"章节页配置第 {line_number} 行页码",
        )
        lines.append(
            f"{normalized_category}|{normalized_ppt_path}|{normalized_slide_number}"
        )
    return "\n".join(lines)


def summarize_ppt_slide_text(raw_text: str, *, max_length: int = 36) -> str:
    condensed = " ".join(raw_text.split())
    if not condensed:
        return "未识别到文字"
    if len(condensed) <= max_length:
        return condensed
    return condensed[: max_length - 1].rstrip() + "…"


def discover_ppt_slide_previews(ppt_path: Path) -> tuple[PptSlidePreview, ...]:
    presentation = Presentation(str(ppt_path))
    previews: list[PptSlidePreview] = []
    for slide_index, slide in enumerate(presentation.slides, start=1):
        text_fragments: list[str] = []
        for shape in slide.shapes:
            if not getattr(shape, "has_text_frame", False):
                continue
            text = shape.text.strip()
            if text:
                text_fragments.append(text)
        title = summarize_ppt_slide_text(" ".join(text_fragments))
        previews.append(
            PptSlidePreview(
                slide_number=slide_index,
                title=title,
                label=f"{slide_index}. {title}",
            )
        )
    return tuple(previews)


def build_ppt_thumbnail_cache_dir(ppt_path: Path) -> Path:
    resolved_path = ppt_path.expanduser().resolve()
    stat = resolved_path.stat()
    cache_key_source = (
        f"{resolved_path}:{stat.st_mtime_ns}:{stat.st_size}".encode("utf-8")
    )
    cache_key = hashlib.sha1(cache_key_source).hexdigest()[:12]
    safe_stem = "".join(
        character if character.isalnum() or character in {"-", "_"} else "_"
        for character in resolved_path.stem
    ).strip("_") or "ppt"
    return GUI_THUMBNAIL_DIR / f"{safe_stem}_{cache_key}"


def list_generated_ppt_thumbnail_images(
    cache_dir: Path,
    stem: str,
) -> tuple[Path, ...]:
    return tuple(sorted(cache_dir.glob(f"{stem}-*.png")))


def generate_ppt_slide_thumbnail_images(
    ppt_path: Path,
) -> tuple[Path, ...]:
    resolved_path = ppt_path.expanduser().resolve()
    cache_dir = build_ppt_thumbnail_cache_dir(resolved_path)
    cache_dir.mkdir(parents=True, exist_ok=True)

    cached_images = list_generated_ppt_thumbnail_images(cache_dir, resolved_path.stem)
    if cached_images:
        return cached_images

    pdf_path = cache_dir / f"{resolved_path.stem}.pdf"
    output_pattern = cache_dir / f"{resolved_path.stem}-%03d.png"

    with tempfile.TemporaryDirectory(
        prefix="lo_profile_",
        dir=str(cache_dir),
    ) as profile_dir:
        profile_uri = Path(profile_dir).resolve().as_uri()
        subprocess.run(
            [
                "soffice",
                f"-env:UserInstallation={profile_uri}",
                "--headless",
                "--convert-to",
                "pdf",
                "--outdir",
                str(cache_dir),
                str(resolved_path),
            ],
            check=True,
            capture_output=True,
            text=True,
        )

    subprocess.run(
        [
            "gs",
            "-dSAFER",
            "-dBATCH",
            "-dNOPAUSE",
            "-sDEVICE=pngalpha",
            f"-r{PPT_THUMBNAIL_DPI}",
            "-o",
            str(output_pattern),
            str(pdf_path),
        ],
        check=True,
        capture_output=True,
        text=True,
    )

    generated_images = list_generated_ppt_thumbnail_images(cache_dir, resolved_path.stem)
    if not generated_images:
        raise RuntimeError("未能生成 PPT 页面缩略图。")
    return generated_images


def build_ppt_layout_lines(config: GuiBatchConfig) -> list[str]:
    lines: list[str] = []
    for region_name, region_label in PPT_LAYOUT_REGION_LABELS:
        left = parse_positive_float_text(
            getattr(config, layout_field_key(region_name, "left")),
            f"{region_label} left",
            allow_zero=True,
        )
        top = parse_positive_float_text(
            getattr(config, layout_field_key(region_name, "top")),
            f"{region_label} top",
            allow_zero=True,
        )
        width = parse_positive_float_text(
            getattr(config, layout_field_key(region_name, "width")),
            f"{region_label} width",
        )
        height = parse_positive_float_text(
            getattr(config, layout_field_key(region_name, "height")),
            f"{region_label} height",
        )
        lines.extend(
            [
                f"[layout.{region_name}]",
                f"left = {left}",
                f"top = {top}",
                f"width = {width}",
                f"height = {height}",
                "",
            ]
        )
    return lines


def build_ppt_config_text(config: GuiBatchConfig) -> str:
    sheet_name_mode = normalize_ppt_sheet_name_mode(config.ppt_sheet_name_mode)
    max_single_table_rows = parse_positive_int_text(
        config.ppt_max_single_table_rows,
        "PPT 单表最大行数",
    )
    max_split_table_rows = parse_positive_int_text(
        config.ppt_max_split_table_rows,
        "PPT 左右双表最大行数",
    )
    body_font_size_pt = parse_positive_float_text(
        config.ppt_body_font_size_pt,
        "PPT 正文字号",
    )
    header_font_size_pt = parse_positive_float_text(
        config.ppt_header_font_size_pt,
        "PPT 表头字号",
    )
    summary_font_size_pt = parse_positive_float_text(
        config.ppt_summary_font_size_pt,
        "PPT 摘要字号",
    )
    template_slide_index = parse_non_negative_int_text(
        config.ppt_template_slide_index,
        "PPT 模板页索引",
    )
    chart_image_dpi = parse_positive_int_text(
        config.ppt_chart_image_dpi,
        "PPT 图表 DPI",
    )
    llm_target_chars = parse_positive_int_text(
        config.ppt_llm_target_chars,
        "备注页目标字数",
    )
    llm_max_tokens = parse_positive_int_text(
        config.ppt_llm_max_tokens,
        "备注页 max_tokens",
    )
    llm_checkpoint_chars = parse_positive_int_text(
        config.ppt_llm_checkpoint_chars,
        "备注页检查点字符数",
    )
    llm_temperature = parse_positive_float_text(
        config.ppt_llm_temperature,
        "备注页 temperature",
        allow_zero=True,
    )

    lines = [
        f"template_path = {toml_quote(config.ppt_template_path)}",
        f"input_dir = {toml_quote(config.stats_output_dir)}",
        f"output_ppt = {toml_quote(config.output_ppt_path)}",
        "",
        f"file_pattern = {toml_quote(config.ppt_file_pattern.strip() or PPT_DEFAULT_FILE_PATTERN)}",
        f'sheet_name_mode = "{sheet_name_mode}"',
    ]
    if sheet_name_mode == "named":
        sheet_name = config.ppt_sheet_name.strip()
        if not sheet_name:
            raise ValueError("sheet_name_mode=named 时请填写 PPT sheet 名。")
        lines.append(f"sheet_name = {toml_quote(sheet_name)}")
    lines.extend(
        [
            f'section_mode = "{config.ppt_section_mode}"',
            f"blank_display = {toml_quote(config.ppt_blank_display)}",
            f"title_suffix = {toml_quote(config.ppt_title_suffix)}",
            f"max_single_table_rows = {max_single_table_rows}",
            f"max_split_table_rows = {max_split_table_rows}",
            f"sort_files = {bool_to_toml(config.ppt_sort_files)}",
            f"body_font_size_pt = {body_font_size_pt}",
            f"header_font_size_pt = {header_font_size_pt}",
            f"summary_font_size_pt = {summary_font_size_pt}",
            f"template_slide_index = {template_slide_index}",
            "",
        ]
    )

    for category_label, ppt_path, slide_number in parse_category_intro_slides_text(
        config.ppt_category_intro_slides_text
    ):
        lines.extend(
            [
                f'[category_intro_slides.{toml_quote(category_label)}]',
                f"ppt_path = {toml_quote(ppt_path)}",
                f"slide_number = {slide_number}",
                "",
            ]
        )

    lines.extend(
        [
            "[chart_page]",
            f"enabled = {bool_to_toml(config.ppt_chart_page_enabled)}",
            f"placeholder_text = {toml_quote(config.ppt_chart_placeholder_text)}",
            f"image_dpi = {chart_image_dpi}",
            "",
            "[llm_notes]",
            f"enabled = {bool_to_toml(config.ppt_llm_notes_enabled)}",
            f"env_path = {toml_quote(config.ppt_llm_env_path)}",
            f"system_role_path = {toml_quote(config.ppt_llm_system_role_path)}",
            f"target_chars = {llm_target_chars}",
            f"temperature = {llm_temperature}",
            f"max_tokens = {llm_max_tokens}",
            f"checkpoint_chars = {llm_checkpoint_chars}",
            "",
        ]
    )
    lines.extend(build_ppt_layout_lines(config))
    return "\n".join(lines).rstrip() + "\n"


def write_runtime_config(prefix: str, content: str) -> Path:
    GUI_RUNTIME_DIR.mkdir(parents=True, exist_ok=True)
    safe_prefix = prefix.replace(" ", "_").replace("/", "_")
    existing_count = len(list(GUI_RUNTIME_DIR.glob(f"{safe_prefix}_*.toml")))
    path = GUI_RUNTIME_DIR / f"{safe_prefix}_{existing_count + 1:03d}.toml"
    path.write_text(content, encoding="utf-8")
    return path


def discover_excel_files(input_dir: Path) -> tuple[Path, ...]:
    if not input_dir.exists() or not input_dir.is_dir():
        return ()
    return tuple(
        sorted(
            path
            for path in input_dir.glob("*.xlsx")
            if path.is_file()
            and not path.name.startswith("~$")
            and not path.name.startswith("._")
        )
    )


def build_python_command(script_name: str, *args: str) -> list[str]:
    return [sys.executable, str(PROJECT_ROOT / script_name), *args]


def build_merge_command(config: GuiBatchConfig) -> list[str] | None:
    if config.workflow_mode != WorkflowMode.MERGED:
        return None
    if not config.merge_input_dirs:
        raise ValueError("多月合并模式下，至少需要一个输入目录。")
    command = build_python_command("merge_questionnaire_workbooks.py")
    for input_dir in config.merge_input_dirs:
        command.extend(["--input-dir", str(input_dir)])
    command.extend(["--output-dir", str(config.merge_output_dir), "--sheet-name", config.sheet_name])
    return command


def build_phase_preprocess_command(config: GuiBatchConfig) -> list[str]:
    excel_files = discover_excel_files(config.effective_input_dir())
    if not excel_files:
        raise ValueError(f"未在 {config.effective_input_dir()} 下找到可处理的 xlsx 文件。")
    command = build_python_command("phase_column_preprocess.py", *(str(path) for path in excel_files))
    command.extend(["--sheet-name", config.sheet_name])
    return command


def build_fill_year_month_command(config: GuiBatchConfig) -> list[str]:
    if not config.year_value.strip():
        raise ValueError("请先填写年份。")
    if not config.month_value.strip():
        raise ValueError("请先填写月份。")
    return build_python_command(
        "fill_year_month_columns.py",
        "--input-dir",
        str(config.effective_input_dir()),
        "--year",
        config.year_value.strip(),
        "--month",
        config.month_value.strip(),
        "--sheet-name",
        config.sheet_name,
    )


def build_survey_stats_command(
    config: GuiBatchConfig,
    *,
    dry_run: bool = False,
    selected_job_names: tuple[str, ...] = (),
) -> list[str]:
    config_path = write_runtime_config("survey_stats", build_survey_stats_config_text(config))
    command = build_python_command("survey_stats.py", "--config", str(config_path))
    for job_name in selected_job_names:
        command.extend(["--job", job_name])
    if dry_run:
        command.append("--dry-run")
    return command


def build_summary_command(config: GuiBatchConfig) -> list[str]:
    return build_python_command(
        "summary_table.py",
        "--input-dir",
        str(config.stats_output_dir),
        "--output-dir",
        str(config.summary_output_dir),
        "--output-name",
        config.summary_output_name,
    )


def build_ppt_command(config: GuiBatchConfig, *, dry_run: bool = False) -> list[str]:
    config_path = write_runtime_config("generate_ppt", build_ppt_config_text(config))
    command = build_python_command("generate_ppt.py", "--config", str(config_path))
    if dry_run:
        command.append("--dry-run")
    return command


def build_task_commands(
    config: GuiBatchConfig,
    selection: MainWorkflowSelection,
    *,
    selected_stats_job_names: tuple[str, ...] = (),
) -> tuple[TaskCommand, ...]:
    commands: list[TaskCommand] = []
    for step_key in build_main_workflow_step_keys(config, selection):
        if step_key == "merge_workbooks":
            merge_command = build_merge_command(config)
            if merge_command is not None:
                commands.append(TaskCommand(step_key, "合并多月问卷", tuple(merge_command)))
        elif step_key == "phase_preprocess":
            commands.append(
                TaskCommand(
                    step_key,
                    PHASE_PREPROCESS_TASK_TITLE,
                    tuple(build_phase_preprocess_command(config)),
                )
            )
        elif step_key == "fill_year_month":
            commands.append(
                TaskCommand(
                    step_key,
                    FILL_YEAR_MONTH_TASK_TITLE,
                    tuple(build_fill_year_month_command(config)),
                )
            )
        elif step_key == "survey_stats":
            commands.append(
                TaskCommand(
                    step_key,
                    "生成分项统计",
                    tuple(
                        build_survey_stats_command(
                            config,
                            selected_job_names=selected_stats_job_names,
                        )
                    ),
                )
            )
        elif step_key == "summary_table":
            commands.append(
                TaskCommand(step_key, "生成汇总表", tuple(build_summary_command(config)))
            )
        elif step_key == "generate_ppt":
            commands.append(
                TaskCommand(step_key, "生成PPT", tuple(build_ppt_command(config)))
            )
    return tuple(commands)


class BackgroundTaskRunner:
    def __init__(self, root: tk.Misc, app: "SurveyPlatformApp") -> None:
        self.root = root
        self.app = app
        self._queue: queue.Queue[tuple[str, object]] = queue.Queue()
        self._thread: threading.Thread | None = None
        self._process: subprocess.Popen[str] | None = None
        self.root.after(120, self._poll)

    @property
    def is_running(self) -> bool:
        return self._thread is not None and self._thread.is_alive()

    def run_tasks(self, tasks: tuple[TaskCommand, ...]) -> None:
        if self.is_running:
            raise RuntimeError("当前已有任务正在执行。")
        self._thread = threading.Thread(target=self._run_worker, args=(tasks,), daemon=True)
        self._thread.start()

    def terminate(self) -> None:
        if self._process is not None and self._process.poll() is None:
            self._process.terminate()

    def _run_worker(self, tasks: tuple[TaskCommand, ...]) -> None:
        success = True
        for task in tasks:
            self._queue.put(("task_started", task))
            try:
                self._process = subprocess.Popen(
                    list(task.command),
                    cwd=PROJECT_ROOT,
                    stdout=subprocess.PIPE,
                    stderr=subprocess.STDOUT,
                    text=True,
                    bufsize=1,
                )
            except OSError as exc:
                self._queue.put(("log", f"[ERROR] 无法启动任务 {task.title}：{exc}"))
                self._queue.put(("task_finished", (task, False)))
                success = False
                break

            assert self._process.stdout is not None
            for line in self._process.stdout:
                self._queue.put(("log", line.rstrip("\n")))
            return_code = self._process.wait()
            if return_code != 0:
                self._queue.put(("log", f"[ERROR] 任务失败：{task.title}（退出码 {return_code}）"))
                self._queue.put(("task_finished", (task, False)))
                success = False
                break
            self._queue.put(("task_finished", (task, True)))
        self._process = None
        self._queue.put(("all_done", success))

    def _poll(self) -> None:
        while True:
            try:
                event, payload = self._queue.get_nowait()
            except queue.Empty:
                break
            self._dispatch(event, payload)
        self.root.after(120, self._poll)

    def _dispatch(self, event: str, payload: object) -> None:
        if event == "log":
            self.app.append_log(str(payload))
            return
        if event == "task_started":
            task = payload
            assert isinstance(task, TaskCommand)
            self.app.on_task_started(task)
            return
        if event == "task_finished":
            task, success = payload
            assert isinstance(task, TaskCommand)
            assert isinstance(success, bool)
            self.app.on_task_finished(task, success)
            return
        if event == "all_done":
            assert isinstance(payload, bool)
            self.app.on_all_tasks_finished(payload)


class SurveyPlatformApp(tk.Tk):
    PAGE_TITLES = {
        "dashboard": "工作台总览",
        "data_source": "数据源管理",
        "preprocess": "预处理",
        "stats": "分项统计",
        "summary": "汇总统计",
        "ppt": "PPT生成",
        "logs": "任务日志",
    }

    def __init__(self) -> None:
        super().__init__()
        self.title("杭博问卷统计平台")
        self.geometry("1520x920")
        self.minsize(1320, 820)
        self.protocol("WM_DELETE_WINDOW", self.on_app_close)
        self.palette = ThemePalette()
        self.workflow_controller = WorkflowRunController()
        self.active_saved_batch_name: str | None = None
        self._saved_batch_profiles: tuple[SavedBatchProfile, ...] = ()
        self._start_action_buttons: list[ttk.Button] = []
        self._terminate_action_buttons: list[ttk.Button] = []
        self._create_variables()
        self._step_status_vars = {
            "data_source": tk.StringVar(value="待设置"),
            "preprocess": tk.StringVar(value="待执行"),
            "stats": tk.StringVar(value="待执行"),
            "summary": tk.StringVar(value="待执行"),
            "ppt": tk.StringVar(value="待执行"),
        }
        self._task_status_label_var = tk.StringVar(value="未开始")
        self._next_action_var = tk.StringVar(value="请先确认数据源")
        self._raw_file_count_var = tk.StringVar(value="0")
        self._stats_file_count_var = tk.StringVar(value="0")
        self._summary_file_var = tk.StringVar(value="未生成")
        self._ppt_file_var = tk.StringVar(value="未生成")
        self._current_page_title_var = tk.StringVar(value=self.PAGE_TITLES["dashboard"])
        self.stats_preview_summary_var = tk.StringVar(value="未扫描客群")
        self._task_title_lookup: dict[str, str] = {}
        self.stats_preview_summary: StatsPreviewSummary | None = None
        self.stats_preview_rows: tuple[CustomerTypePreviewRow, ...] = ()
        self.selected_customer_types: frozenset[str] = frozenset()
        self._workflow_progress_dialog: tk.Toplevel | None = None
        self._workflow_progress_tree: ttk.Treeview | None = None
        self._workflow_progress_summary_var = tk.StringVar(value="未启动主流程")
        self._workflow_progress_close_var = tk.StringVar(value="关闭")
        self._workflow_progress_item_ids: dict[str, str] = {}
        self._workflow_progress_last_page_key: str | None = None
        self._ppt_advanced_frame: ttk.Frame | None = None
        self._ppt_advanced_toggle_var = tk.StringVar(value="展开高级 PPT 配置")
        self._ppt_advanced_visible = False
        self.ppt_chart_placeholder_text_widget: scrolledtext.ScrolledText | None = None
        self.ppt_category_intro_tree: ttk.Treeview | None = None
        self.ppt_category_intro_entries: list[tuple[str, str, int]] = []
        self.ppt_category_intro_slide_preview_combobox: ttk.Combobox | None = None
        self.ppt_category_intro_slide_previews: tuple[PptSlidePreview, ...] = ()
        self._ppt_thumbnail_dialog: tk.Toplevel | None = None
        self._ppt_thumbnail_photo_refs: list[tk.PhotoImage] = []
        self._apply_theme()
        self._build_layout()
        self._restore_persistent_state()
        self.runner = BackgroundTaskRunner(self, self)
        self.refresh_all_status_views()
        self._sync_action_button_states()
        self.show_page("dashboard")

    def _create_variables(self) -> None:
        self.batch_name_var = tk.StringVar(value="2026年3月")
        self.workflow_mode_var = tk.StringVar(value=WorkflowMode.SINGLE.value)
        self.single_input_dir_var = tk.StringVar(value=str(PROJECT_ROOT / "datas" / "3月"))
        self.merge_output_dir_var = tk.StringVar(value=str(PROJECT_ROOT / "datas" / "合并结果"))
        self.sheet_name_var = tk.StringVar(value=DEFAULT_SHEET_NAME)
        self.year_value_var = tk.StringVar(value="2026")
        self.month_value_var = tk.StringVar(value="03")
        self.stats_output_dir_var = tk.StringVar(value=str(PROJECT_ROOT / "输出结果" / "3月"))
        self.calculation_mode_var = tk.StringVar(value="template")
        self.output_format_var = tk.StringVar(value="xlsx")
        self.summary_output_dir_var = tk.StringVar(value=str(PROJECT_ROOT / "汇总结果" / "3月"))
        self.summary_output_name_var = tk.StringVar(value="3月客户类型满意度汇总表.xlsx")
        self.ppt_template_path_var = tk.StringVar(value=str(PROJECT_ROOT / "templates" / "template.pptx"))
        self.output_ppt_path_var = tk.StringVar(value=str(PROJECT_ROOT / "输出结果" / "3月满意度报告.pptx"))
        self.ppt_file_pattern_var = tk.StringVar(value=PPT_DEFAULT_FILE_PATTERN)
        self.ppt_sheet_name_mode_var = tk.StringVar(value=PPT_DEFAULT_SHEET_NAME_MODE)
        self.ppt_sheet_name_var = tk.StringVar(value="")
        self.ppt_section_mode_var = tk.StringVar(value="auto")
        self.ppt_blank_display_var = tk.StringVar(value="")
        self.ppt_title_suffix_var = tk.StringVar(value=PPT_DEFAULT_TITLE_SUFFIX)
        self.ppt_max_single_table_rows_var = tk.StringVar(value=PPT_DEFAULT_MAX_SINGLE_TABLE_ROWS)
        self.ppt_max_split_table_rows_var = tk.StringVar(value=PPT_DEFAULT_MAX_SPLIT_TABLE_ROWS)
        self.ppt_sort_files_var = tk.BooleanVar(value=True)
        self.ppt_body_font_size_pt_var = tk.StringVar(value=PPT_DEFAULT_BODY_FONT_SIZE_PT)
        self.ppt_header_font_size_pt_var = tk.StringVar(value=PPT_DEFAULT_HEADER_FONT_SIZE_PT)
        self.ppt_summary_font_size_pt_var = tk.StringVar(value=PPT_DEFAULT_SUMMARY_FONT_SIZE_PT)
        self.ppt_template_slide_index_var = tk.StringVar(value=PPT_DEFAULT_TEMPLATE_SLIDE_INDEX)
        self.ppt_chart_page_enabled_var = tk.BooleanVar(value=False)
        self.ppt_chart_placeholder_text_var = tk.StringVar(value=PPT_DEFAULT_CHART_PLACEHOLDER_TEXT)
        self.ppt_chart_image_dpi_var = tk.StringVar(value=PPT_DEFAULT_CHART_IMAGE_DPI)
        self.ppt_llm_notes_enabled_var = tk.BooleanVar(value=False)
        self.ppt_llm_env_path_var = tk.StringVar(value=PPT_DEFAULT_LLM_ENV_PATH)
        self.ppt_llm_system_role_path_var = tk.StringVar(value=PPT_DEFAULT_LLM_SYSTEM_ROLE_PATH)
        self.ppt_llm_target_chars_var = tk.StringVar(value=PPT_DEFAULT_LLM_TARGET_CHARS)
        self.ppt_llm_temperature_var = tk.StringVar(value=PPT_DEFAULT_LLM_TEMPERATURE)
        self.ppt_llm_max_tokens_var = tk.StringVar(value=PPT_DEFAULT_LLM_MAX_TOKENS)
        self.ppt_llm_checkpoint_chars_var = tk.StringVar(value=PPT_DEFAULT_LLM_CHECKPOINT_CHARS)
        self.ppt_category_intro_slides_text_var = tk.StringVar(value="")
        default_category_label = PPT_CATEGORY_LABEL_VALUES[0] if PPT_CATEGORY_LABEL_VALUES else ""
        self.ppt_category_intro_category_var = tk.StringVar(value=default_category_label)
        self.ppt_category_intro_ppt_path_var = tk.StringVar(value="")
        self.ppt_category_intro_slide_number_var = tk.StringVar(
            value=PPT_DEFAULT_CATEGORY_INTRO_SLIDE_NUMBER
        )
        self.ppt_category_intro_slide_preview_var = tk.StringVar(value="")
        self.ppt_category_intro_slide_status_var = tk.StringVar(
            value="选择封面 PPT 后可读取页面列表。"
        )
        self.ppt_layout_vars: dict[str, dict[str, tk.StringVar]] = {}
        for region_name, _ in PPT_LAYOUT_REGION_LABELS:
            defaults = PPT_LAYOUT_DEFAULTS[region_name]
            self.ppt_layout_vars[region_name] = {
                "left": tk.StringVar(value=defaults[0]),
                "top": tk.StringVar(value=defaults[1]),
                "width": tk.StringVar(value=defaults[2]),
                "height": tk.StringVar(value=defaults[3]),
            }
        self.saved_batch_var = tk.StringVar(value="")
        self.merge_input_list: list[str] = []

    def _apply_theme(self) -> None:
        style = ttk.Style(self)
        style.theme_use("clam")
        self.configure(bg=self.palette.background)
        style.configure(".", background=self.palette.background, foreground=self.palette.text)
        style.configure("Root.TFrame", background=self.palette.background)
        style.configure("Surface.TFrame", background=self.palette.surface)
        style.configure("Card.TFrame", background=self.palette.surface, relief="flat")
        style.configure(
            "Sidebar.TFrame",
            background=self.palette.surface_alt,
            bordercolor=self.palette.border,
            relief="solid",
        )
        style.configure("Header.TLabel", background=self.palette.background, foreground=self.palette.text, font=("PingFang SC", 18, "bold"))
        style.configure("Root.SubHeader.TLabel", background=self.palette.background, foreground=self.palette.text, font=("PingFang SC", 12, "bold"))
        style.configure("Root.Body.TLabel", background=self.palette.background, foreground=self.palette.text, font=("PingFang SC", 10))
        style.configure("Root.Muted.TLabel", background=self.palette.background, foreground=self.palette.muted_text, font=("PingFang SC", 10))
        style.configure("SubHeader.TLabel", background=self.palette.surface, foreground=self.palette.text, font=("PingFang SC", 12, "bold"))
        style.configure("Body.TLabel", background=self.palette.surface, foreground=self.palette.text, font=("PingFang SC", 10))
        style.configure("Muted.TLabel", background=self.palette.surface, foreground=self.palette.muted_text, font=("PingFang SC", 10))
        style.configure("TRadiobutton", background=self.palette.surface, foreground=self.palette.text, font=("PingFang SC", 10))
        style.map("TRadiobutton", background=[("active", self.palette.surface)])
        style.configure("TCheckbutton", background=self.palette.surface, foreground=self.palette.text, font=("PingFang SC", 10))
        style.map("TCheckbutton", background=[("active", self.palette.surface)])
        style.configure("Sidebar.TButton", background=self.palette.surface_alt, foreground=self.palette.text, padding=(16, 10), anchor="w", relief="flat")
        style.map("Sidebar.TButton", background=[("active", self.palette.surface), ("pressed", self.palette.surface)])
        style.configure("Primary.TButton", background=self.palette.primary, foreground="#ffffff", padding=(14, 8), relief="flat")
        style.map(
            "Primary.TButton",
            background=[("active", self.palette.primary_hover), ("pressed", self.palette.primary_hover)],
            foreground=[("disabled", "#f5e8e3")],
        )
        style.configure("Secondary.TButton", background=self.palette.surface, foreground=self.palette.primary, padding=(12, 8), relief="flat", borderwidth=1)
        style.map("Secondary.TButton", background=[("active", self.palette.surface_alt)])
        style.configure("Treeview", background="#fffdfb", fieldbackground="#fffdfb", rowheight=28, foreground=self.palette.text, bordercolor=self.palette.border)
        style.configure("Treeview.Heading", background=self.palette.surface_alt, foreground=self.palette.text, font=("PingFang SC", 10, "bold"))
        style.configure("TEntry", fieldbackground="#fffdfb", bordercolor=self.palette.border)
        style.configure("TCombobox", fieldbackground="#fffdfb")
        style.configure("TLabelframe", background=self.palette.surface, foreground=self.palette.text)
        style.configure("TLabelframe.Label", background=self.palette.surface, foreground=self.palette.text, font=("PingFang SC", 10, "bold"))

    def _build_layout(self) -> None:
        root = ttk.Frame(self, style="Root.TFrame", padding=16)
        root.pack(fill="both", expand=True)
        root.columnconfigure(1, weight=1)
        root.rowconfigure(1, weight=1)

        header = ttk.Frame(root, style="Root.TFrame")
        header.grid(row=0, column=0, columnspan=3, sticky="ew", pady=(0, 12))
        header.columnconfigure(1, weight=1)
        header.columnconfigure(2, weight=0)
        ttk.Label(header, text="杭博问卷统计平台", style="Header.TLabel").grid(row=0, column=0, sticky="w")

        batch_manager = ttk.Frame(header, style="Root.TFrame")
        batch_manager.grid(row=0, column=1, sticky="ew", padx=(20, 12))
        batch_manager.columnconfigure(1, weight=1)
        batch_manager.columnconfigure(3, weight=1)
        ttk.Label(batch_manager, text="批次名称", style="Root.Body.TLabel").grid(row=0, column=0, sticky="w")
        ttk.Entry(batch_manager, textvariable=self.batch_name_var, width=18).grid(
            row=0,
            column=1,
            sticky="ew",
            padx=(8, 12),
        )
        ttk.Label(batch_manager, text="已保存批次", style="Root.Body.TLabel").grid(row=0, column=2, sticky="w")
        self.saved_batch_combobox = ttk.Combobox(
            batch_manager,
            textvariable=self.saved_batch_var,
            state="readonly",
            values=(),
            width=22,
        )
        self.saved_batch_combobox.grid(row=0, column=3, sticky="ew", padx=(8, 8))
        ttk.Button(batch_manager, text="加载", style="Secondary.TButton", command=self.load_selected_batch).grid(row=0, column=4, padx=(0, 8))
        ttk.Button(batch_manager, text="新建", style="Secondary.TButton", command=self.create_new_batch).grid(row=0, column=5)

        header_meta = ttk.Frame(header, style="Root.TFrame")
        header_meta.grid(row=0, column=2, sticky="e")
        ttk.Button(header_meta, text="保存批次", style="Secondary.TButton", command=self.save_current_batch).grid(row=0, column=0, padx=(0, 8))
        ttk.Button(header_meta, text="删除批次", style="Secondary.TButton", command=self.delete_selected_batch).grid(row=0, column=1, padx=(0, 8))
        terminate_button = ttk.Button(header_meta, text="终止当前任务", style="Secondary.TButton", command=self.terminate_current_task)
        terminate_button.grid(row=0, column=2, padx=(0, 8))
        self._register_terminate_button(terminate_button)
        main_workflow_button = ttk.Button(header_meta, text="一键执行主流程", style="Primary.TButton", command=self.open_main_workflow_dialog)
        main_workflow_button.grid(row=0, column=3)
        self._register_start_button(main_workflow_button)

        sidebar = ttk.Frame(root, style="Sidebar.TFrame", padding=12)
        sidebar.grid(row=1, column=0, sticky="nsw")

        content_wrap = ttk.Frame(root, style="Root.TFrame")
        content_wrap.grid(row=1, column=1, sticky="nsew", padx=(12, 12))
        content_wrap.columnconfigure(0, weight=1)
        content_wrap.rowconfigure(1, weight=1)

        log_wrap = ttk.Frame(root, style="Surface.TFrame", padding=12)
        log_wrap.grid(row=1, column=2, sticky="nsew")
        root.columnconfigure(2, weight=0)

        page_header = ttk.Frame(content_wrap, style="Root.TFrame")
        page_header.grid(row=0, column=0, sticky="ew", pady=(0, 8))
        page_header.columnconfigure(0, weight=1)
        ttk.Label(page_header, textvariable=self._current_page_title_var, style="Root.SubHeader.TLabel").grid(row=0, column=0, sticky="w")
        ttk.Label(page_header, textvariable=self._task_status_label_var, style="Root.Muted.TLabel").grid(row=0, column=1, sticky="e")

        self.page_container = ttk.Frame(content_wrap, style="Surface.TFrame", padding=12)
        self.page_container.grid(row=1, column=0, sticky="nsew")
        self.page_container.columnconfigure(0, weight=1)
        self.page_container.rowconfigure(0, weight=1)

        self._build_sidebar(sidebar)
        self._build_log_panel(log_wrap)
        self._build_pages()

    def _build_sidebar(self, parent: ttk.Frame) -> None:
        nav_items = (
            ("dashboard", "1. 工作台总览"),
            ("data_source", "2. 数据源管理"),
            ("preprocess", "3. 预处理"),
            ("stats", "4. 分项统计"),
            ("summary", "5. 汇总统计"),
            ("ppt", "6. PPT生成"),
            ("logs", "7. 任务日志"),
        )
        for index, (page_key, label) in enumerate(nav_items):
            ttk.Button(
                parent,
                text=label,
                style="Sidebar.TButton",
                command=lambda value=page_key: self.show_page(value),
            ).grid(row=index, column=0, sticky="ew", pady=4)
        parent.columnconfigure(0, weight=1)

    def _build_log_panel(self, parent: ttk.Frame) -> None:
        parent.columnconfigure(0, weight=1)
        parent.rowconfigure(1, weight=1)
        ttk.Label(parent, text="运行日志", style="SubHeader.TLabel").grid(row=0, column=0, sticky="w", pady=(0, 8))
        self.log_text = scrolledtext.ScrolledText(
            parent,
            wrap="word",
            width=42,
            font=("SF Mono", 10),
            bg="#fffdfb",
            fg=self.palette.text,
            relief="flat",
            padx=10,
            pady=10,
        )
        self.log_text.grid(row=1, column=0, sticky="nsew")
        self.log_text.configure(state="disabled")
        ttk.Button(parent, text="清空日志", style="Secondary.TButton", command=self.clear_log).grid(row=2, column=0, sticky="e", pady=(8, 0))

    def _build_pages(self) -> None:
        self.pages: dict[str, ttk.Frame] = {}
        builders = {
            "dashboard": self._build_dashboard_page,
            "data_source": self._build_data_source_page,
            "preprocess": self._build_preprocess_page,
            "stats": self._build_stats_page,
            "summary": self._build_summary_page,
            "ppt": self._build_ppt_page,
            "logs": self._build_logs_page,
        }
        for key, builder in builders.items():
            frame = builder(self.page_container)
            frame.grid(row=0, column=0, sticky="nsew")
            self.pages[key] = frame

    def _register_start_button(self, button: ttk.Button) -> ttk.Button:
        self._start_action_buttons.append(button)
        return button

    def _register_terminate_button(self, button: ttk.Button) -> ttk.Button:
        self._terminate_action_buttons.append(button)
        return button

    def _sync_action_button_states(self) -> None:
        start_state = "normal" if self.workflow_controller.start_enabled else "disabled"
        terminate_state = "normal" if self.workflow_controller.terminate_enabled else "disabled"
        active_start_buttons: list[ttk.Button] = []
        for button in self._start_action_buttons:
            if not button.winfo_exists():
                continue
            button.configure(state=start_state)
            active_start_buttons.append(button)
        self._start_action_buttons = active_start_buttons
        active_terminate_buttons: list[ttk.Button] = []
        for button in self._terminate_action_buttons:
            if not button.winfo_exists():
                continue
            button.configure(state=terminate_state)
            active_terminate_buttons.append(button)
        self._terminate_action_buttons = active_terminate_buttons

    def _build_dashboard_page(self, parent: ttk.Frame) -> ttk.Frame:
        frame = ttk.Frame(parent, style="Surface.TFrame")
        frame.columnconfigure(0, weight=1)
        ttk.Label(frame, text="当前批次与主流程状态", style="SubHeader.TLabel").grid(row=0, column=0, sticky="w", pady=(0, 10))

        top_info = ttk.Frame(frame, style="Surface.TFrame")
        top_info.grid(row=1, column=0, sticky="ew")
        top_info.columnconfigure(1, weight=1)
        info_pairs = (
            ("批次名称", self.batch_name_var),
            ("数据模式", self.workflow_mode_var),
            ("原始数据", self.single_input_dir_var),
            ("分项输出", self.stats_output_dir_var),
            ("汇总输出", self.summary_output_dir_var),
            ("PPT输出", self.output_ppt_path_var),
        )
        for index, (label, variable) in enumerate(info_pairs):
            ttk.Label(top_info, text=f"{label}：", style="Body.TLabel").grid(row=index, column=0, sticky="nw", pady=2)
            ttk.Label(top_info, textvariable=variable, style="Muted.TLabel").grid(row=index, column=1, sticky="nw", pady=2)

        step_frame = ttk.Frame(frame, style="Surface.TFrame")
        step_frame.grid(row=2, column=0, sticky="ew", pady=(18, 12))
        steps = (
            ("数据源", "data_source"),
            ("预处理", "preprocess"),
            ("分项统计", "stats"),
            ("汇总表", "summary"),
            ("PPT", "ppt"),
        )
        for index, (label, key) in enumerate(steps):
            card = ttk.Frame(step_frame, style="Card.TFrame", padding=12)
            card.grid(row=0, column=index, sticky="nsew", padx=(0 if index == 0 else 8, 0))
            ttk.Label(card, text=label, style="SubHeader.TLabel").pack(anchor="w")
            ttk.Label(card, textvariable=self._step_status_vars[key], style="Muted.TLabel").pack(anchor="w", pady=(4, 0))
            step_frame.columnconfigure(index, weight=1)

        card_row = ttk.Frame(frame, style="Surface.TFrame")
        card_row.grid(row=3, column=0, sticky="ew", pady=(6, 12))
        card_row.columnconfigure((0, 1, 2, 3), weight=1)
        self._build_metric_card(card_row, 0, "原始文件数", self._raw_file_count_var)
        self._build_metric_card(card_row, 1, "分项文件数", self._stats_file_count_var)
        self._build_metric_card(card_row, 2, "汇总表状态", self._summary_file_var)
        self._build_metric_card(card_row, 3, "PPT状态", self._ppt_file_var)

        action_box = ttk.Frame(frame, style="Card.TFrame", padding=12)
        action_box.grid(row=4, column=0, sticky="ew")
        action_box.columnconfigure(1, weight=1)
        ttk.Label(action_box, text="当前建议操作", style="SubHeader.TLabel").grid(row=0, column=0, sticky="w")
        ttk.Label(action_box, textvariable=self._next_action_var, style="Muted.TLabel").grid(row=1, column=0, sticky="w", pady=(4, 10))
        action_buttons = (
            ("导入数据源", lambda: self.show_page("data_source")),
            ("执行预处理", lambda: self.show_page("preprocess")),
            ("生成分项统计", lambda: self.show_page("stats")),
            ("生成汇总表", lambda: self.show_page("summary")),
            ("生成PPT", lambda: self.show_page("ppt")),
        )
        button_row = ttk.Frame(action_box, style="Card.TFrame")
        button_row.grid(row=2, column=0, sticky="ew")
        for index, (label, callback) in enumerate(action_buttons):
            style_name = "Primary.TButton" if label == "生成分项统计" else "Secondary.TButton"
            ttk.Button(button_row, text=label, style=style_name, command=callback).grid(row=0, column=index, padx=(0, 8))
        return frame

    def _build_metric_card(self, parent: ttk.Frame, column: int, title: str, variable: tk.StringVar) -> None:
        card = ttk.Frame(parent, style="Card.TFrame", padding=12)
        card.grid(row=0, column=column, sticky="ew", padx=(0 if column == 0 else 8, 0))
        ttk.Label(card, text=title, style="SubHeader.TLabel").pack(anchor="w")
        ttk.Label(card, textvariable=variable, style="Muted.TLabel").pack(anchor="w", pady=(6, 0))

    def _build_data_source_page(self, parent: ttk.Frame) -> ttk.Frame:
        frame = ttk.Frame(parent, style="Surface.TFrame")
        frame.columnconfigure(0, weight=1)

        mode_box = ttk.LabelFrame(frame, text="处理方式", padding=12)
        mode_box.grid(row=0, column=0, sticky="ew", pady=(0, 12))
        ttk.Radiobutton(mode_box, text="单个月份 / 单个批次", value=WorkflowMode.SINGLE.value, variable=self.workflow_mode_var, command=self.refresh_all_status_views).grid(row=0, column=0, sticky="w", pady=2)
        ttk.Radiobutton(mode_box, text="合并多个月份后处理", value=WorkflowMode.MERGED.value, variable=self.workflow_mode_var, command=self.refresh_all_status_views).grid(row=1, column=0, sticky="w", pady=2)

        single_box = ttk.LabelFrame(frame, text="单月处理配置", padding=12)
        single_box.grid(row=1, column=0, sticky="ew", pady=(0, 12))
        single_box.columnconfigure(1, weight=1)
        ttk.Label(single_box, text="原始数据目录", style="Body.TLabel").grid(row=0, column=0, sticky="w")
        ttk.Entry(single_box, textvariable=self.single_input_dir_var).grid(row=0, column=1, sticky="ew", padx=8)
        ttk.Button(single_box, text="浏览", style="Secondary.TButton", command=lambda: self.choose_directory(self.single_input_dir_var)).grid(row=0, column=2)

        merge_box = ttk.LabelFrame(frame, text="多月合并配置", padding=12)
        merge_box.grid(row=2, column=0, sticky="ew", pady=(0, 12))
        merge_box.columnconfigure(0, weight=1)
        self.merge_listbox = tk.Listbox(merge_box, height=5, relief="flat", bg="#fffdfb", fg=self.palette.text, selectbackground=self.palette.surface_alt)
        self.merge_listbox.grid(row=0, column=0, columnspan=3, sticky="ew")
        list_actions = ttk.Frame(merge_box, style="Surface.TFrame")
        list_actions.grid(row=1, column=0, columnspan=3, sticky="ew", pady=(8, 8))
        ttk.Button(list_actions, text="添加目录", style="Secondary.TButton", command=self.add_merge_input_dir).grid(row=0, column=0, padx=(0, 8))
        ttk.Button(list_actions, text="删除所选", style="Secondary.TButton", command=self.remove_selected_merge_dir).grid(row=0, column=1)
        ttk.Label(merge_box, text="合并输出目录", style="Body.TLabel").grid(row=2, column=0, sticky="w")
        ttk.Entry(merge_box, textvariable=self.merge_output_dir_var).grid(row=3, column=0, sticky="ew", pady=(4, 0))
        ttk.Button(merge_box, text="浏览", style="Secondary.TButton", command=lambda: self.choose_directory(self.merge_output_dir_var)).grid(row=3, column=1, padx=8, sticky="w")
        merge_button = ttk.Button(merge_box, text="执行合并", style="Primary.TButton", command=self.run_merge_task)
        merge_button.grid(row=3, column=2, sticky="e")
        self._register_start_button(merge_button)

        scan_box = ttk.LabelFrame(frame, text="当前输入目录与扫描结果", padding=12)
        scan_box.grid(row=3, column=0, sticky="nsew")
        scan_box.columnconfigure(0, weight=1)
        ttk.Label(scan_box, text="当前生效输入目录", style="Body.TLabel").grid(row=0, column=0, sticky="w")
        self.effective_input_var = tk.StringVar()
        ttk.Label(scan_box, textvariable=self.effective_input_var, style="Muted.TLabel").grid(row=1, column=0, sticky="w", pady=(2, 8))
        ttk.Button(scan_box, text="刷新扫描结果", style="Secondary.TButton", command=self.refresh_all_status_views).grid(row=2, column=0, sticky="w", pady=(0, 8))
        self.data_source_tree = ttk.Treeview(scan_box, columns=("path", "exists"), show="headings", height=8)
        self.data_source_tree.heading("path", text="路径")
        self.data_source_tree.heading("exists", text="状态")
        self.data_source_tree.column("path", width=620)
        self.data_source_tree.column("exists", width=120, anchor="center")
        self.data_source_tree.grid(row=3, column=0, sticky="nsew")
        return frame

    def _build_preprocess_page(self, parent: ttk.Frame) -> ttk.Frame:
        frame = ttk.Frame(parent, style="Surface.TFrame")
        frame.columnconfigure(0, weight=1)

        phase_box = ttk.LabelFrame(frame, text=PHASE_PREPROCESS_SECTION_TITLE, padding=12)
        phase_box.grid(row=0, column=0, sticky="ew", pady=(0, 12))
        ttk.Label(phase_box, text=PHASE_PREPROCESS_DESCRIPTION, style="Muted.TLabel", wraplength=900, justify="left").grid(row=0, column=0, sticky="w")
        phase_button = ttk.Button(phase_box, text=PHASE_PREPROCESS_BUTTON_TEXT, style="Primary.TButton", command=self.run_phase_preprocess_task)
        phase_button.grid(row=1, column=0, sticky="w", pady=(10, 0))
        self._register_start_button(phase_button)

        fill_box = ttk.LabelFrame(frame, text=FILL_YEAR_MONTH_SECTION_TITLE, padding=12)
        fill_box.grid(row=1, column=0, sticky="ew", pady=(0, 12))
        ttk.Label(fill_box, text=FILL_YEAR_MONTH_DESCRIPTION, style="Muted.TLabel", wraplength=900, justify="left").grid(row=0, column=0, columnspan=3, sticky="w", pady=(0, 10))
        ttk.Label(fill_box, text="年份", style="Body.TLabel").grid(row=1, column=0, sticky="w")
        ttk.Entry(fill_box, textvariable=self.year_value_var, width=12).grid(row=2, column=0, sticky="w", pady=(4, 0))
        ttk.Label(fill_box, text="月份", style="Body.TLabel").grid(row=1, column=1, sticky="w", padx=(12, 0))
        ttk.Entry(fill_box, textvariable=self.month_value_var, width=12).grid(row=2, column=1, sticky="w", padx=(12, 0), pady=(4, 0))
        fill_button = ttk.Button(fill_box, text=FILL_YEAR_MONTH_BUTTON_TEXT, style="Secondary.TButton", command=self.run_fill_year_month_task)
        fill_button.grid(row=2, column=2, sticky="w", padx=(12, 0))
        self._register_start_button(fill_button)

        note_box = ttk.LabelFrame(frame, text="说明", padding=12)
        note_box.grid(row=2, column=0, sticky="ew")
        note_text = (
            "月份检查脚本文档已存在，但当前仓库中缺少 check_start_time_month.py。"
            f" 第一版 GUI 先保留预处理主链：{PHASE_PREPROCESS_TASK_TITLE}、{FILL_YEAR_MONTH_TASK_TITLE}。"
        )
        ttk.Label(note_box, text=note_text, style="Muted.TLabel", wraplength=900, justify="left").grid(row=0, column=0, sticky="w")
        return frame

    def _build_stats_page(self, parent: ttk.Frame) -> ttk.Frame:
        frame = ttk.Frame(parent, style="Surface.TFrame")
        frame.columnconfigure(0, weight=1)

        config_box = ttk.LabelFrame(frame, text="分项统计配置", padding=12)
        config_box.grid(row=0, column=0, sticky="ew", pady=(0, 12))
        config_box.columnconfigure(1, weight=1)
        ttk.Label(config_box, text="输入目录", style="Body.TLabel").grid(row=0, column=0, sticky="w")
        self.stats_effective_input_var = tk.StringVar()
        ttk.Label(config_box, textvariable=self.stats_effective_input_var, style="Muted.TLabel").grid(row=0, column=1, sticky="w", padx=8)
        ttk.Label(config_box, text="输出目录", style="Body.TLabel").grid(row=1, column=0, sticky="w")
        ttk.Entry(config_box, textvariable=self.stats_output_dir_var).grid(row=1, column=1, sticky="ew", padx=8, pady=6)
        ttk.Button(config_box, text="浏览", style="Secondary.TButton", command=lambda: self.choose_directory(self.stats_output_dir_var)).grid(row=1, column=2, sticky="w")
        ttk.Label(config_box, text="sheet", style="Body.TLabel").grid(row=2, column=0, sticky="w")
        ttk.Entry(config_box, textvariable=self.sheet_name_var, width=16).grid(row=2, column=1, sticky="w", padx=8)
        ttk.Label(config_box, text="计算口径", style="Body.TLabel").grid(row=2, column=2, sticky="w")
        ttk.Combobox(config_box, textvariable=self.calculation_mode_var, values=("template", "summary"), width=12, state="readonly").grid(row=2, column=3, sticky="w", padx=8)
        button_row = ttk.Frame(config_box, style="Surface.TFrame")
        button_row.grid(row=3, column=0, columnspan=4, sticky="w", pady=(10, 0))
        scan_button = ttk.Button(button_row, text="扫描可统计客群", style="Secondary.TButton", command=self.scan_stats_preview)
        scan_button.grid(row=0, column=0, padx=(0, 8))
        self._register_start_button(scan_button)
        dry_run_button = ttk.Button(button_row, text="dry-run 校验", style="Secondary.TButton", command=self.run_survey_stats_dry_run_task)
        dry_run_button.grid(row=0, column=1, padx=(0, 8))
        self._register_start_button(dry_run_button)
        stats_button = ttk.Button(button_row, text="开始生成分项统计", style="Primary.TButton", command=self.run_survey_stats_task)
        stats_button.grid(row=0, column=2, padx=(0, 8))
        self._register_start_button(stats_button)
        ttk.Button(button_row, text="保存当前批次", style="Secondary.TButton", command=self.save_current_batch).grid(row=0, column=3)

        summary_box = ttk.LabelFrame(frame, text="当前状态", padding=12)
        summary_box.grid(row=1, column=0, sticky="ew", pady=(0, 12))
        self.stats_summary_var = tk.StringVar(value="待扫描")
        ttk.Label(summary_box, textvariable=self.stats_summary_var, style="Muted.TLabel").grid(row=0, column=0, sticky="w")

        preview_box = ttk.LabelFrame(frame, text="客群扫描预览", padding=12)
        preview_box.grid(row=2, column=0, sticky="nsew", pady=(0, 12))
        preview_box.columnconfigure(0, weight=1)
        ttk.Label(preview_box, textvariable=self.stats_preview_summary_var, style="Muted.TLabel").grid(row=0, column=0, sticky="w", pady=(0, 8))
        ttk.Label(
            preview_box,
            text="双击某一行可切换勾选；正式执行和 dry-run 都只处理当前已勾选客群。",
            style="Muted.TLabel",
        ).grid(row=1, column=0, sticky="w", pady=(0, 8))
        preview_actions = ttk.Frame(preview_box, style="Surface.TFrame")
        preview_actions.grid(row=2, column=0, sticky="w", pady=(0, 8))
        ttk.Button(preview_actions, text="全选", style="Secondary.TButton", command=self.select_all_customer_types).grid(row=0, column=0, padx=(0, 8))
        ttk.Button(preview_actions, text="清空", style="Secondary.TButton", command=self.clear_selected_customer_types).grid(row=0, column=1, padx=(0, 8))
        ttk.Button(preview_actions, text="仅选可生成", style="Secondary.TButton", command=self.select_ready_customer_types).grid(row=0, column=2)
        self.stats_preview_tree = ttk.Treeview(
            preview_box,
            columns=("selected", "customer_type", "source_file", "status", "output_name", "detail"),
            show="headings",
            height=10,
        )
        self.stats_preview_tree.heading("selected", text="选择")
        self.stats_preview_tree.heading("customer_type", text="客群")
        self.stats_preview_tree.heading("source_file", text="来源文件")
        self.stats_preview_tree.heading("status", text="状态")
        self.stats_preview_tree.heading("output_name", text="输出文件")
        self.stats_preview_tree.heading("detail", text="详情")
        self.stats_preview_tree.column("selected", width=70, anchor="center")
        self.stats_preview_tree.column("customer_type", width=180)
        self.stats_preview_tree.column("source_file", width=140)
        self.stats_preview_tree.column("status", width=130, anchor="center")
        self.stats_preview_tree.column("output_name", width=180)
        self.stats_preview_tree.column("detail", width=420)
        self.stats_preview_tree.grid(row=3, column=0, sticky="nsew")
        self.stats_preview_tree.bind("<Double-1>", self.on_stats_preview_row_double_click)

        result_box = ttk.LabelFrame(frame, text="输出目录文件预览", padding=12)
        result_box.grid(row=3, column=0, sticky="nsew")
        result_box.columnconfigure(0, weight=1)
        self.stats_tree = ttk.Treeview(result_box, columns=("name", "path"), show="headings", height=10)
        self.stats_tree.heading("name", text="文件名")
        self.stats_tree.heading("path", text="路径")
        self.stats_tree.column("name", width=220)
        self.stats_tree.column("path", width=650)
        self.stats_tree.grid(row=0, column=0, sticky="nsew")
        return frame

    def _build_summary_page(self, parent: ttk.Frame) -> ttk.Frame:
        frame = ttk.Frame(parent, style="Surface.TFrame")
        frame.columnconfigure(0, weight=1)
        box = ttk.LabelFrame(frame, text="汇总统计配置", padding=12)
        box.grid(row=0, column=0, sticky="ew", pady=(0, 12))
        box.columnconfigure(1, weight=1)
        ttk.Label(box, text="分项结果目录", style="Body.TLabel").grid(row=0, column=0, sticky="w")
        ttk.Label(box, textvariable=self.stats_output_dir_var, style="Muted.TLabel").grid(row=0, column=1, sticky="w", padx=8)
        ttk.Label(box, text="汇总输出目录", style="Body.TLabel").grid(row=1, column=0, sticky="w")
        ttk.Entry(box, textvariable=self.summary_output_dir_var).grid(row=1, column=1, sticky="ew", padx=8, pady=6)
        ttk.Button(box, text="浏览", style="Secondary.TButton", command=lambda: self.choose_directory(self.summary_output_dir_var)).grid(row=1, column=2, sticky="w")
        ttk.Label(box, text="输出文件名", style="Body.TLabel").grid(row=2, column=0, sticky="w")
        ttk.Entry(box, textvariable=self.summary_output_name_var).grid(row=2, column=1, sticky="ew", padx=8)
        summary_button = ttk.Button(box, text="生成汇总表", style="Primary.TButton", command=self.run_summary_task)
        summary_button.grid(row=3, column=0, sticky="w", pady=(10, 0))
        self._register_start_button(summary_button)

        result_box = ttk.LabelFrame(frame, text="结果状态", padding=12)
        result_box.grid(row=1, column=0, sticky="ew")
        self.summary_status_var = tk.StringVar(value="待执行")
        ttk.Label(result_box, textvariable=self.summary_status_var, style="Muted.TLabel").grid(row=0, column=0, sticky="w")
        return frame

    def _toggle_ppt_advanced_config(self) -> None:
        if self._ppt_advanced_frame is None:
            return
        self._ppt_advanced_visible = not self._ppt_advanced_visible
        if self._ppt_advanced_visible:
            self._ppt_advanced_frame.grid()
            self._ppt_advanced_toggle_var.set("收起高级 PPT 配置")
            return
        self._ppt_advanced_frame.grid_remove()
        self._ppt_advanced_toggle_var.set("展开高级 PPT 配置")

    def _build_ppt_general_tab(self, notebook: ttk.Notebook) -> None:
        tab = ttk.Frame(notebook, style="Surface.TFrame", padding=12)
        tab.columnconfigure(1, weight=1)
        tab.columnconfigure(3, weight=1)
        rows = (
            ("文件匹配", self.ppt_file_pattern_var, 0, 0),
            ("sheet_name_mode", self.ppt_sheet_name_mode_var, 0, 2),
            ("sheet_name", self.ppt_sheet_name_var, 1, 0),
            ("标题后缀", self.ppt_title_suffix_var, 1, 2),
            ("单表最大行数", self.ppt_max_single_table_rows_var, 2, 0),
            ("双表单侧最大行数", self.ppt_max_split_table_rows_var, 2, 2),
            ("正文字号", self.ppt_body_font_size_pt_var, 3, 0),
            ("表头字号", self.ppt_header_font_size_pt_var, 3, 2),
            ("摘要字号", self.ppt_summary_font_size_pt_var, 4, 0),
            ("模板页索引", self.ppt_template_slide_index_var, 4, 2),
        )
        for label, variable, row, column in rows:
            ttk.Label(tab, text=label, style="Body.TLabel").grid(row=row, column=column, sticky="w", pady=4)
            if label == "sheet_name_mode":
                ttk.Combobox(
                    tab,
                    textvariable=variable,
                    values=PPT_SHEET_NAME_MODE_VALUES,
                    state="readonly",
                    width=16,
                ).grid(row=row, column=column + 1, sticky="ew", padx=8, pady=4)
            else:
                ttk.Entry(tab, textvariable=variable).grid(row=row, column=column + 1, sticky="ew", padx=8, pady=4)
        ttk.Checkbutton(
            tab,
            text="按文件名排序后生成",
            variable=self.ppt_sort_files_var,
        ).grid(row=5, column=0, columnspan=2, sticky="w", pady=(8, 0))
        notebook.add(tab, text="基础高级")

    def _build_ppt_chart_tab(self, notebook: ttk.Notebook) -> None:
        tab = ttk.Frame(notebook, style="Surface.TFrame", padding=12)
        tab.columnconfigure(1, weight=1)
        ttk.Checkbutton(
            tab,
            text="启用图表页（每个客户分组数据页后追加图表页）",
            variable=self.ppt_chart_page_enabled_var,
        ).grid(row=0, column=0, columnspan=3, sticky="w", pady=(0, 8))
        ttk.Label(tab, text="图表 DPI", style="Body.TLabel").grid(row=1, column=0, sticky="w")
        ttk.Entry(tab, textvariable=self.ppt_chart_image_dpi_var, width=14).grid(row=1, column=1, sticky="w", padx=8)
        ttk.Label(
            tab,
            text="回退占位文案",
            style="Body.TLabel",
        ).grid(row=2, column=0, sticky="nw", pady=(10, 4))
        self.ppt_chart_placeholder_text_widget = scrolledtext.ScrolledText(
            tab,
            height=5,
            wrap="word",
            font=("PingFang SC", 10),
            bg="#fffdfb",
            fg=self.palette.text,
            relief="flat",
        )
        self.ppt_chart_placeholder_text_widget.grid(row=2, column=1, columnspan=2, sticky="ew", pady=(10, 4))
        self.ppt_chart_placeholder_text_widget.insert("1.0", self.ppt_chart_placeholder_text_var.get())
        notebook.add(tab, text="图表页")

    def _build_ppt_notes_tab(self, notebook: ttk.Notebook) -> None:
        tab = ttk.Frame(notebook, style="Surface.TFrame", padding=12)
        tab.columnconfigure(1, weight=1)
        ttk.Checkbutton(
            tab,
            text="启用备注页 LLM 分析",
            variable=self.ppt_llm_notes_enabled_var,
        ).grid(row=0, column=0, columnspan=3, sticky="w", pady=(0, 8))
        ttk.Label(tab, text="env_path", style="Body.TLabel").grid(row=1, column=0, sticky="w")
        ttk.Entry(tab, textvariable=self.ppt_llm_env_path_var).grid(row=1, column=1, sticky="ew", padx=8, pady=4)
        ttk.Button(
            tab,
            text="浏览",
            style="Secondary.TButton",
            command=lambda: self.choose_file(
                self.ppt_llm_env_path_var,
                [("Env", ".env"), ("All Files", "*")],
            ),
        ).grid(row=1, column=2, sticky="w")
        ttk.Label(tab, text="system_role_path", style="Body.TLabel").grid(row=2, column=0, sticky="w")
        ttk.Entry(tab, textvariable=self.ppt_llm_system_role_path_var).grid(row=2, column=1, sticky="ew", padx=8, pady=4)
        ttk.Button(
            tab,
            text="浏览",
            style="Secondary.TButton",
            command=lambda: self.choose_file(
                self.ppt_llm_system_role_path_var,
                [("Markdown", "*.md"), ("All Files", "*")],
            ),
        ).grid(row=2, column=2, sticky="w")
        scalar_rows = (
            ("目标字数", self.ppt_llm_target_chars_var, 3, 0),
            ("temperature", self.ppt_llm_temperature_var, 3, 2),
            ("max_tokens", self.ppt_llm_max_tokens_var, 4, 0),
            ("checkpoint_chars", self.ppt_llm_checkpoint_chars_var, 4, 2),
        )
        tab.columnconfigure(3, weight=1)
        for label, variable, row, column in scalar_rows:
            ttk.Label(tab, text=label, style="Body.TLabel").grid(row=row, column=column, sticky="w", pady=4)
            ttk.Entry(tab, textvariable=variable, width=14).grid(row=row, column=column + 1, sticky="ew", padx=8, pady=4)
        notebook.add(tab, text="备注页")

    def _build_ppt_category_tab(self, notebook: ttk.Notebook) -> None:
        tab = ttk.Frame(notebook, style="Surface.TFrame", padding=12)
        tab.columnconfigure(0, weight=1)
        tab.rowconfigure(1, weight=1)
        ttk.Label(
            tab,
            text=(
                "按客户大类配置封面页。选中上方列表可修改；未配置的客户大类不会插入封面。"
            ),
            style="Muted.TLabel",
            wraplength=860,
            justify="left",
        ).grid(row=0, column=0, sticky="w", pady=(0, 8))
        tree_wrap = ttk.Frame(tab, style="Surface.TFrame")
        tree_wrap.grid(row=1, column=0, sticky="nsew")
        tree_wrap.columnconfigure(0, weight=1)
        tree_wrap.rowconfigure(0, weight=1)
        self.ppt_category_intro_tree = ttk.Treeview(
            tree_wrap,
            columns=("category", "ppt_path", "slide_number"),
            show="headings",
            height=6,
        )
        self.ppt_category_intro_tree.heading("category", text="客户大类")
        self.ppt_category_intro_tree.heading("ppt_path", text="封面 PPT 路径")
        self.ppt_category_intro_tree.heading("slide_number", text="页码")
        self.ppt_category_intro_tree.column("category", width=180, anchor="w")
        self.ppt_category_intro_tree.column("ppt_path", width=470, anchor="w")
        self.ppt_category_intro_tree.column("slide_number", width=90, anchor="center")
        self.ppt_category_intro_tree.grid(row=0, column=0, sticky="nsew")
        self.ppt_category_intro_tree.bind(
            "<<TreeviewSelect>>",
            self._load_selected_ppt_category_intro_entry,
        )
        tree_scrollbar = ttk.Scrollbar(
            tree_wrap,
            orient="vertical",
            command=self.ppt_category_intro_tree.yview,
        )
        tree_scrollbar.grid(row=0, column=1, sticky="ns")
        self.ppt_category_intro_tree.configure(yscrollcommand=tree_scrollbar.set)

        form = ttk.LabelFrame(tab, text="封面配置", padding=12)
        form.grid(row=2, column=0, sticky="ew", pady=(12, 0))
        form.columnconfigure(1, weight=1)
        ttk.Label(form, text="客户大类", style="Body.TLabel").grid(
            row=0, column=0, sticky="w"
        )
        ttk.Combobox(
            form,
            textvariable=self.ppt_category_intro_category_var,
            values=PPT_CATEGORY_LABEL_VALUES,
        ).grid(row=0, column=1, sticky="ew", padx=8, pady=4)
        ttk.Label(form, text="页码", style="Body.TLabel").grid(
            row=0, column=2, sticky="w"
        )
        ttk.Entry(
            form,
            textvariable=self.ppt_category_intro_slide_number_var,
            width=10,
        ).grid(row=0, column=3, sticky="w", padx=8, pady=4)
        ttk.Label(form, text="封面 PPT", style="Body.TLabel").grid(
            row=1, column=0, sticky="w"
        )
        ttk.Entry(
            form,
            textvariable=self.ppt_category_intro_ppt_path_var,
        ).grid(row=1, column=1, columnspan=3, sticky="ew", padx=8, pady=4)
        ttk.Button(
            form,
            text="浏览",
            style="Secondary.TButton",
            command=self.choose_ppt_category_intro_file,
        ).grid(row=1, column=4, sticky="w")
        ttk.Label(form, text="页面预览", style="Body.TLabel").grid(
            row=2, column=0, sticky="w"
        )
        self.ppt_category_intro_slide_preview_combobox = ttk.Combobox(
            form,
            textvariable=self.ppt_category_intro_slide_preview_var,
            state="readonly",
        )
        self.ppt_category_intro_slide_preview_combobox.grid(
            row=2, column=1, columnspan=3, sticky="ew", padx=8, pady=4
        )
        self.ppt_category_intro_slide_preview_combobox.bind(
            "<<ComboboxSelected>>",
            self._on_ppt_category_intro_slide_preview_selected,
        )
        ttk.Button(
            form,
            text="读取页列表",
            style="Secondary.TButton",
            command=lambda: self._load_ppt_category_intro_slide_previews(
                show_message=True
            ),
        ).grid(row=2, column=4, sticky="w")
        ttk.Button(
            form,
            text="按缩略图选择",
            style="Secondary.TButton",
            command=self.open_ppt_category_intro_thumbnail_dialog,
        ).grid(row=3, column=4, sticky="w", pady=(8, 0))
        ttk.Label(
            form,
            textvariable=self.ppt_category_intro_slide_status_var,
            style="Muted.TLabel",
            wraplength=760,
            justify="left",
        ).grid(row=3, column=0, columnspan=4, sticky="w", pady=(4, 0))

        action_row = ttk.Frame(form, style="Surface.TFrame")
        action_row.grid(row=4, column=0, columnspan=5, sticky="w", pady=(10, 0))
        ttk.Button(
            action_row,
            text="新增",
            style="Secondary.TButton",
            command=self.add_ppt_category_intro_entry,
        ).grid(row=0, column=0, padx=(0, 8))
        ttk.Button(
            action_row,
            text="更新选中",
            style="Secondary.TButton",
            command=self.update_selected_ppt_category_intro_entry,
        ).grid(row=0, column=1, padx=(0, 8))
        ttk.Button(
            action_row,
            text="删除选中",
            style="Secondary.TButton",
            command=self.delete_selected_ppt_category_intro_entry,
        ).grid(row=0, column=2, padx=(0, 8))
        ttk.Button(
            action_row,
            text="上移",
            style="Secondary.TButton",
            command=lambda: self.move_selected_ppt_category_intro_entry(-1),
        ).grid(row=0, column=3, padx=(0, 8))
        ttk.Button(
            action_row,
            text="下移",
            style="Secondary.TButton",
            command=lambda: self.move_selected_ppt_category_intro_entry(1),
        ).grid(row=0, column=4, padx=(0, 8))
        ttk.Button(
            action_row,
            text="清空表单",
            style="Secondary.TButton",
            command=self.clear_ppt_category_intro_form,
        ).grid(row=0, column=5)
        ttk.Label(
            form,
            text="页码按 PowerPoint 可见页码填写，从 1 开始。",
            style="Muted.TLabel",
        ).grid(row=5, column=0, columnspan=5, sticky="w", pady=(8, 0))
        notebook.add(tab, text="客户大类封面")

    def _build_ppt_layout_tab(self, notebook: ttk.Notebook) -> None:
        tab = ttk.Frame(notebook, style="Surface.TFrame", padding=12)
        for column in range(1, 5):
            tab.columnconfigure(column, weight=1)
        headers = ("区域", "left", "top", "width", "height")
        for column, header in enumerate(headers):
            ttk.Label(tab, text=header, style="SubHeader.TLabel").grid(row=0, column=column, sticky="w", pady=(0, 8))
        for row_index, (region_name, region_label) in enumerate(PPT_LAYOUT_REGION_LABELS, start=1):
            ttk.Label(tab, text=region_label, style="Body.TLabel").grid(row=row_index, column=0, sticky="w", pady=4)
            for column_index, field_name in enumerate(("left", "top", "width", "height"), start=1):
                ttk.Entry(
                    tab,
                    textvariable=self.ppt_layout_vars[region_name][field_name],
                    width=10,
                ).grid(row=row_index, column=column_index, sticky="ew", padx=6, pady=4)
        notebook.add(tab, text="布局")

    def _build_ppt_advanced_panel(self, parent: ttk.Frame) -> ttk.Frame:
        frame = ttk.Frame(parent, style="Surface.TFrame")
        frame.columnconfigure(0, weight=1)
        notebook = ttk.Notebook(frame)
        notebook.grid(row=0, column=0, sticky="ew")
        self._build_ppt_general_tab(notebook)
        self._build_ppt_chart_tab(notebook)
        self._build_ppt_notes_tab(notebook)
        self._build_ppt_category_tab(notebook)
        self._build_ppt_layout_tab(notebook)
        return frame

    def _build_ppt_page(self, parent: ttk.Frame) -> ttk.Frame:
        frame = ttk.Frame(parent, style="Surface.TFrame")
        frame.columnconfigure(0, weight=1)
        box = ttk.LabelFrame(frame, text="PPT生成配置", padding=12)
        box.grid(row=0, column=0, sticky="ew", pady=(0, 12))
        box.columnconfigure(1, weight=1)
        ttk.Label(box, text="统计结果目录", style="Body.TLabel").grid(row=0, column=0, sticky="w")
        ttk.Label(box, textvariable=self.stats_output_dir_var, style="Muted.TLabel").grid(row=0, column=1, sticky="w", padx=8)
        ttk.Label(box, text="模板PPT", style="Body.TLabel").grid(row=1, column=0, sticky="w")
        ttk.Entry(box, textvariable=self.ppt_template_path_var).grid(row=1, column=1, sticky="ew", padx=8, pady=6)
        ttk.Button(box, text="浏览", style="Secondary.TButton", command=lambda: self.choose_file(self.ppt_template_path_var, [("PowerPoint", "*.pptx")])).grid(row=1, column=2)
        ttk.Label(box, text="输出PPT", style="Body.TLabel").grid(row=2, column=0, sticky="w")
        ttk.Entry(box, textvariable=self.output_ppt_path_var).grid(row=2, column=1, sticky="ew", padx=8, pady=6)
        ttk.Button(box, text="浏览", style="Secondary.TButton", command=lambda: self.choose_save_file(self.output_ppt_path_var, ".pptx", [("PowerPoint", "*.pptx")])).grid(row=2, column=2)
        ttk.Label(box, text="section_mode", style="Body.TLabel").grid(row=3, column=0, sticky="w")
        ttk.Combobox(box, textvariable=self.ppt_section_mode_var, values=("auto", "template", "summary"), state="readonly", width=12).grid(row=3, column=1, sticky="w", padx=8)
        ttk.Label(box, text="空值显示", style="Body.TLabel").grid(row=4, column=0, sticky="w")
        ttk.Entry(box, textvariable=self.ppt_blank_display_var, width=18).grid(row=4, column=1, sticky="w", padx=8, pady=6)
        ttk.Button(
            box,
            textvariable=self._ppt_advanced_toggle_var,
            style="Secondary.TButton",
            command=self._toggle_ppt_advanced_config,
        ).grid(row=5, column=0, sticky="w", pady=(10, 0))
        self._ppt_advanced_frame = self._build_ppt_advanced_panel(box)
        self._ppt_advanced_frame.grid(row=6, column=0, columnspan=3, sticky="ew", pady=(12, 0))
        self._ppt_advanced_frame.grid_remove()
        button_row = ttk.Frame(box, style="Surface.TFrame")
        button_row.grid(row=7, column=0, columnspan=3, sticky="w", pady=(12, 0))
        ppt_dry_run_button = ttk.Button(button_row, text="dry-run 校验", style="Secondary.TButton", command=self.run_ppt_dry_run_task)
        ppt_dry_run_button.grid(row=0, column=0, padx=(0, 8))
        self._register_start_button(ppt_dry_run_button)
        ppt_button = ttk.Button(button_row, text="生成PPT", style="Primary.TButton", command=self.run_ppt_task)
        ppt_button.grid(row=0, column=1)
        self._register_start_button(ppt_button)

        status_box = ttk.LabelFrame(frame, text="结果状态", padding=12)
        status_box.grid(row=1, column=0, sticky="ew")
        self.ppt_status_var = tk.StringVar(value="待执行")
        ttk.Label(status_box, textvariable=self.ppt_status_var, style="Muted.TLabel").grid(row=0, column=0, sticky="w")
        return frame

    def _build_logs_page(self, parent: ttk.Frame) -> ttk.Frame:
        frame = ttk.Frame(parent, style="Surface.TFrame")
        frame.columnconfigure(0, weight=1)
        info_box = ttk.LabelFrame(frame, text="说明", padding=12)
        info_box.grid(row=0, column=0, sticky="ew", pady=(0, 12))
        ttk.Label(
            info_box,
            text="右侧日志面板会实时显示脚本输出。运行 survey_stats.py 和 generate_ppt.py 时，GUI 会自动生成临时 TOML 配置写入 logs/gui_runtime。",
            style="Muted.TLabel",
            wraplength=900,
            justify="left",
        ).grid(row=0, column=0, sticky="w")
        runtime_box = ttk.LabelFrame(frame, text="运行目录", padding=12)
        runtime_box.grid(row=1, column=0, sticky="ew")
        ttk.Label(runtime_box, text=str(GUI_RUNTIME_DIR), style="Muted.TLabel").grid(row=0, column=0, sticky="w")
        return frame

    def show_page(self, page_key: str) -> None:
        self._current_page_title_var.set(self.PAGE_TITLES[page_key])
        frame = self.pages[page_key]
        frame.tkraise()

    def _refresh_merge_listbox(self) -> None:
        self.merge_listbox.delete(0, "end")
        for item in self.merge_input_list:
            self.merge_listbox.insert("end", item)

    def _read_text_widget_value(
        self,
        widget: scrolledtext.ScrolledText | None,
        fallback_var: tk.StringVar,
    ) -> str:
        if widget is None:
            return fallback_var.get()
        return widget.get("1.0", "end-1c")

    def _set_text_widget_value(
        self,
        widget: scrolledtext.ScrolledText | None,
        value: str,
        fallback_var: tk.StringVar,
    ) -> None:
        fallback_var.set(value)
        if widget is None:
            return
        widget.delete("1.0", "end")
        widget.insert("1.0", value)

    def _selected_ppt_category_intro_index(self) -> int | None:
        if self.ppt_category_intro_tree is None:
            return None
        selection = self.ppt_category_intro_tree.selection()
        if not selection:
            return None
        try:
            return int(selection[0])
        except ValueError:
            return None

    def _refresh_ppt_category_intro_slide_preview_combobox(self) -> None:
        values = tuple(preview.label for preview in self.ppt_category_intro_slide_previews)
        if self.ppt_category_intro_slide_preview_combobox is not None:
            self.ppt_category_intro_slide_preview_combobox.configure(values=values)
        if not values:
            self.ppt_category_intro_slide_preview_var.set("")

    def _sync_ppt_category_intro_slide_preview_selection(
        self,
        *,
        preferred_slide_number: int | None = None,
    ) -> None:
        if not self.ppt_category_intro_slide_previews:
            self.ppt_category_intro_slide_preview_var.set("")
            return
        target_slide_number = preferred_slide_number
        if target_slide_number is None:
            try:
                target_slide_number = int(self.ppt_category_intro_slide_number_var.get())
            except ValueError:
                target_slide_number = None
        for preview in self.ppt_category_intro_slide_previews:
            if preview.slide_number == target_slide_number:
                self.ppt_category_intro_slide_preview_var.set(preview.label)
                return
        first_preview = self.ppt_category_intro_slide_previews[0]
        self.ppt_category_intro_slide_preview_var.set(first_preview.label)
        self.ppt_category_intro_slide_number_var.set(str(first_preview.slide_number))

    def _load_ppt_category_intro_slide_previews(
        self,
        *,
        preferred_slide_number: int | None = None,
        show_message: bool = False,
    ) -> bool:
        ppt_path_text = self.ppt_category_intro_ppt_path_var.get().strip()
        if not ppt_path_text:
            self.ppt_category_intro_slide_previews = ()
            self._refresh_ppt_category_intro_slide_preview_combobox()
            self.ppt_category_intro_slide_status_var.set("请先选择封面 PPT。")
            if show_message:
                messagebox.showinfo("未选择文件", "请先选择封面 PPT 文件。")
            return False
        ppt_path = Path(ppt_path_text).expanduser()
        if not ppt_path.is_absolute():
            ppt_path = (PROJECT_ROOT / ppt_path).resolve()
        if not ppt_path.exists():
            self.ppt_category_intro_slide_previews = ()
            self._refresh_ppt_category_intro_slide_preview_combobox()
            self.ppt_category_intro_slide_status_var.set("封面 PPT 不存在，请检查路径。")
            if show_message:
                messagebox.showerror("文件不存在", f"未找到封面 PPT：\n{ppt_path}")
            return False
        try:
            previews = discover_ppt_slide_previews(ppt_path)
        except Exception as exc:
            self.ppt_category_intro_slide_previews = ()
            self._refresh_ppt_category_intro_slide_preview_combobox()
            self.ppt_category_intro_slide_status_var.set("读取页面列表失败。")
            if show_message:
                messagebox.showerror("读取失败", f"封面 PPT 读取失败：\n{exc}")
            return False
        self.ppt_category_intro_slide_previews = previews
        self._refresh_ppt_category_intro_slide_preview_combobox()
        if previews:
            self._sync_ppt_category_intro_slide_preview_selection(
                preferred_slide_number=preferred_slide_number
            )
            self.ppt_category_intro_slide_status_var.set(
                f"已读取 {len(previews)} 页，可直接从下拉框选择。"
            )
        else:
            self.ppt_category_intro_slide_status_var.set("该 PPT 没有可读取的页面。")
        return True

    def _on_ppt_category_intro_slide_preview_selected(
        self,
        _event: object | None = None,
    ) -> None:
        selected_label = self.ppt_category_intro_slide_preview_var.get().strip()
        for preview in self.ppt_category_intro_slide_previews:
            if preview.label == selected_label:
                self.ppt_category_intro_slide_number_var.set(str(preview.slide_number))
                return

    def choose_ppt_category_intro_file(self) -> None:
        initial = self.ppt_category_intro_ppt_path_var.get().strip() or str(PROJECT_ROOT)
        selected = filedialog.askopenfilename(
            initialdir=str(Path(initial).expanduser().resolve().parent),
            filetypes=[("PowerPoint", "*.pptx"), ("All Files", "*")],
        )
        if not selected:
            return
        self.ppt_category_intro_ppt_path_var.set(selected)
        self._load_ppt_category_intro_slide_previews(show_message=True)
        self.refresh_all_status_views()

    def _thumbnail_select_slide(
        self,
        preview: PptSlidePreview,
    ) -> None:
        self.ppt_category_intro_slide_number_var.set(str(preview.slide_number))
        self.ppt_category_intro_slide_preview_var.set(preview.label)
        self._close_ppt_thumbnail_dialog()

    def _close_ppt_thumbnail_dialog(self) -> None:
        if self._ppt_thumbnail_dialog is not None:
            dialog = self._ppt_thumbnail_dialog
            self._ppt_thumbnail_dialog = None
            dialog.destroy()
        self._ppt_thumbnail_photo_refs = []

    def open_ppt_category_intro_thumbnail_dialog(self) -> None:
        if not self._load_ppt_category_intro_slide_previews(show_message=True):
            return

        ppt_path = Path(self.ppt_category_intro_ppt_path_var.get().strip()).expanduser()
        if not ppt_path.is_absolute():
            ppt_path = (PROJECT_ROOT / ppt_path).resolve()

        try:
            thumbnail_paths = generate_ppt_slide_thumbnail_images(ppt_path)
        except FileNotFoundError as exc:
            messagebox.showerror(
                "缺少缩略图依赖",
                f"无法生成缩略图，缺少命令：{exc.filename}",
            )
            return
        except subprocess.CalledProcessError as exc:
            error_message = exc.stderr.strip() or exc.stdout.strip() or str(exc)
            messagebox.showerror(
                "缩略图生成失败",
                f"封面页缩略图生成失败：\n{error_message}",
            )
            return
        except Exception as exc:
            messagebox.showerror("缩略图生成失败", f"封面页缩略图生成失败：\n{exc}")
            return

        previews = self.ppt_category_intro_slide_previews
        total_cards = min(len(previews), len(thumbnail_paths))
        if total_cards == 0:
            messagebox.showinfo("无可选页面", "当前封面 PPT 没有可供选择的页面。")
            return

        if self._ppt_thumbnail_dialog is not None:
            self._close_ppt_thumbnail_dialog()

        dialog = tk.Toplevel(self)
        dialog.title("按缩略图选择封面页")
        dialog.transient(self)
        dialog.configure(bg=self.palette.background)
        dialog.geometry("980x720")
        dialog.minsize(860, 560)
        dialog.protocol("WM_DELETE_WINDOW", self._close_ppt_thumbnail_dialog)

        outer = ttk.Frame(dialog, style="Surface.TFrame", padding=16)
        outer.pack(fill="both", expand=True)
        outer.columnconfigure(0, weight=1)
        outer.rowconfigure(1, weight=1)

        ttk.Label(
            outer,
            text="点击缩略图即可选择封面页。首次打开会生成缓存，后续会更快。",
            style="Muted.TLabel",
            wraplength=900,
            justify="left",
        ).grid(row=0, column=0, sticky="w")

        canvas = tk.Canvas(
            outer,
            bg=self.palette.surface,
            highlightthickness=0,
            bd=0,
        )
        canvas.grid(row=1, column=0, sticky="nsew", pady=(12, 0))
        scrollbar = ttk.Scrollbar(outer, orient="vertical", command=canvas.yview)
        scrollbar.grid(row=1, column=1, sticky="ns", pady=(12, 0))
        canvas.configure(yscrollcommand=scrollbar.set)

        cards_frame = ttk.Frame(canvas, style="Surface.TFrame")
        canvas_window = canvas.create_window((0, 0), window=cards_frame, anchor="nw")

        def _sync_canvas_width(event: tk.Event) -> None:
            canvas.itemconfigure(canvas_window, width=event.width)

        def _sync_scroll_region(_event: tk.Event) -> None:
            canvas.configure(scrollregion=canvas.bbox("all"))

        canvas.bind("<Configure>", _sync_canvas_width)
        cards_frame.bind("<Configure>", _sync_scroll_region)

        photo_refs: list[tk.PhotoImage] = []
        for card_index in range(total_cards):
            preview = previews[card_index]
            thumbnail_path = thumbnail_paths[card_index]
            image = tk.PhotoImage(file=str(thumbnail_path))
            scale_factor = max(
                1,
                math.ceil(image.width() / PPT_THUMBNAIL_MAX_WIDTH),
                math.ceil(image.height() / PPT_THUMBNAIL_MAX_HEIGHT),
            )
            if scale_factor > 1:
                image = image.subsample(scale_factor, scale_factor)
            photo_refs.append(image)

            card = ttk.Frame(cards_frame, style="Card.TFrame", padding=10)
            row_index = card_index // 3
            column_index = card_index % 3
            card.grid(
                row=row_index,
                column=column_index,
                sticky="n",
                padx=8,
                pady=8,
            )

            image_button = ttk.Button(
                card,
                image=image,
                style="Secondary.TButton",
                command=lambda selected_preview=preview: self._thumbnail_select_slide(
                    selected_preview
                ),
            )
            image_button.grid(row=0, column=0, sticky="nsew")
            ttk.Label(
                card,
                text=preview.label,
                style="Body.TLabel",
                wraplength=240,
                justify="center",
            ).grid(row=1, column=0, sticky="ew", pady=(8, 0))

        footer = ttk.Frame(outer, style="Surface.TFrame")
        footer.grid(row=2, column=0, columnspan=2, sticky="ew", pady=(12, 0))
        ttk.Button(
            footer,
            text="关闭",
            style="Secondary.TButton",
            command=self._close_ppt_thumbnail_dialog,
        ).grid(row=0, column=0, sticky="e")

        self._ppt_thumbnail_dialog = dialog
        self._ppt_thumbnail_photo_refs = photo_refs

    def _sync_ppt_category_intro_text_var(self) -> None:
        self.ppt_category_intro_slides_text_var.set(
            build_category_intro_slides_text(self.ppt_category_intro_entries)
        )

    def _refresh_ppt_category_intro_tree(
        self,
        *,
        selected_index: int | None = None,
    ) -> None:
        self._sync_ppt_category_intro_text_var()
        if self.ppt_category_intro_tree is None:
            return
        for item_id in self.ppt_category_intro_tree.get_children():
            self.ppt_category_intro_tree.delete(item_id)
        for index, (category_label, ppt_path, slide_number) in enumerate(
            self.ppt_category_intro_entries
        ):
            self.ppt_category_intro_tree.insert(
                "",
                "end",
                iid=str(index),
                values=(category_label, ppt_path, slide_number),
            )
        if (
            selected_index is not None
            and 0 <= selected_index < len(self.ppt_category_intro_entries)
        ):
            selected_item = str(selected_index)
            self.ppt_category_intro_tree.selection_set((selected_item,))
            self.ppt_category_intro_tree.focus(selected_item)
        else:
            for item_id in self.ppt_category_intro_tree.selection():
                self.ppt_category_intro_tree.selection_remove(item_id)

    def _load_ppt_category_intro_entries_from_text(self, text: str) -> None:
        self.ppt_category_intro_entries = list(parse_category_intro_slides_text(text))
        self._refresh_ppt_category_intro_tree()
        self.clear_ppt_category_intro_form()

    def _load_selected_ppt_category_intro_entry(self, _event: object | None = None) -> None:
        selected_index = self._selected_ppt_category_intro_index()
        if selected_index is None:
            return
        category_label, ppt_path, slide_number = self.ppt_category_intro_entries[
            selected_index
        ]
        self.ppt_category_intro_category_var.set(category_label)
        self.ppt_category_intro_ppt_path_var.set(ppt_path)
        self.ppt_category_intro_slide_number_var.set(str(slide_number))
        self._load_ppt_category_intro_slide_previews(
            preferred_slide_number=slide_number
        )

    def _read_ppt_category_intro_form_entry(self) -> tuple[str, str, int]:
        category_label = self.ppt_category_intro_category_var.get().strip()
        ppt_path = self.ppt_category_intro_ppt_path_var.get().strip()
        if not category_label:
            raise ValueError("请先选择或填写客户大类。")
        if not ppt_path:
            raise ValueError("请先选择封面 PPT 路径。")
        slide_number = parse_positive_int_text(
            self.ppt_category_intro_slide_number_var.get(),
            "封面页码",
        )
        return (category_label, ppt_path, slide_number)

    def _ensure_unique_ppt_category_intro_label(
        self,
        category_label: str,
        *,
        ignore_index: int | None = None,
    ) -> None:
        for index, (existing_label, _, _) in enumerate(self.ppt_category_intro_entries):
            if ignore_index is not None and index == ignore_index:
                continue
            if existing_label == category_label:
                raise ValueError(f"客户大类“{category_label}”已经配置过封面。")

    def add_ppt_category_intro_entry(self) -> None:
        try:
            entry = self._read_ppt_category_intro_form_entry()
            self._ensure_unique_ppt_category_intro_label(entry[0])
        except ValueError as exc:
            messagebox.showerror("新增封面失败", str(exc))
            return
        self.ppt_category_intro_entries.append(entry)
        self._refresh_ppt_category_intro_tree(
            selected_index=len(self.ppt_category_intro_entries) - 1
        )
        self.refresh_all_status_views()

    def update_selected_ppt_category_intro_entry(self) -> None:
        selected_index = self._selected_ppt_category_intro_index()
        if selected_index is None:
            messagebox.showinfo("未选择条目", "请先在上方列表中选择要更新的客户大类封面。")
            return
        try:
            entry = self._read_ppt_category_intro_form_entry()
            self._ensure_unique_ppt_category_intro_label(
                entry[0],
                ignore_index=selected_index,
            )
        except ValueError as exc:
            messagebox.showerror("更新封面失败", str(exc))
            return
        self.ppt_category_intro_entries[selected_index] = entry
        self._refresh_ppt_category_intro_tree(selected_index=selected_index)
        self.refresh_all_status_views()

    def delete_selected_ppt_category_intro_entry(self) -> None:
        selected_index = self._selected_ppt_category_intro_index()
        if selected_index is None:
            messagebox.showinfo("未选择条目", "请先选择要删除的客户大类封面。")
            return
        del self.ppt_category_intro_entries[selected_index]
        next_index = min(selected_index, len(self.ppt_category_intro_entries) - 1)
        self._refresh_ppt_category_intro_tree(
            selected_index=next_index if next_index >= 0 else None
        )
        if next_index >= 0:
            self._load_selected_ppt_category_intro_entry()
        else:
            self.clear_ppt_category_intro_form()
        self.refresh_all_status_views()

    def move_selected_ppt_category_intro_entry(self, offset: int) -> None:
        selected_index = self._selected_ppt_category_intro_index()
        if selected_index is None:
            messagebox.showinfo("未选择条目", "请先选择要移动的客户大类封面。")
            return
        target_index = selected_index + offset
        if target_index < 0 or target_index >= len(self.ppt_category_intro_entries):
            return
        entry = self.ppt_category_intro_entries.pop(selected_index)
        self.ppt_category_intro_entries.insert(target_index, entry)
        self._refresh_ppt_category_intro_tree(selected_index=target_index)
        self._load_selected_ppt_category_intro_entry()
        self.refresh_all_status_views()

    def clear_ppt_category_intro_form(self) -> None:
        default_category_label = (
            PPT_CATEGORY_LABEL_VALUES[0] if PPT_CATEGORY_LABEL_VALUES else ""
        )
        self.ppt_category_intro_category_var.set(default_category_label)
        self.ppt_category_intro_ppt_path_var.set("")
        self.ppt_category_intro_slide_number_var.set(
            PPT_DEFAULT_CATEGORY_INTRO_SLIDE_NUMBER
        )
        self.ppt_category_intro_slide_previews = ()
        self._refresh_ppt_category_intro_slide_preview_combobox()
        self.ppt_category_intro_slide_status_var.set("选择封面 PPT 后可读取页面列表。")
        if self.ppt_category_intro_tree is not None:
            for item_id in self.ppt_category_intro_tree.selection():
                self.ppt_category_intro_tree.selection_remove(item_id)

    def _reset_stats_preview_state(self) -> None:
        self.stats_preview_summary = None
        self.stats_preview_rows = ()
        self.selected_customer_types = frozenset()
        self._refresh_stats_preview_tree()
        self._refresh_stats_preview_summary()

    def _refresh_saved_batch_profiles(self, *, selected_name: str | None = None) -> None:
        self._saved_batch_profiles = load_saved_batch_profiles()
        saved_batch_names = tuple(profile.batch_name for profile in self._saved_batch_profiles)
        self.saved_batch_combobox.configure(values=saved_batch_names)

        target_name = selected_name
        if target_name is None and self.active_saved_batch_name in saved_batch_names:
            target_name = self.active_saved_batch_name
        if target_name in saved_batch_names:
            self.saved_batch_var.set(target_name or "")
            return
        self.saved_batch_var.set("")

    def _persist_session_state(self) -> None:
        save_gui_session(
            self.current_config(),
            active_saved_batch_name=self.active_saved_batch_name,
        )

    def _apply_config_to_form(self, config: GuiBatchConfig) -> None:
        self.batch_name_var.set(config.batch_name)
        self.workflow_mode_var.set(config.workflow_mode.value)
        self.single_input_dir_var.set(str(config.single_input_dir))
        self.merge_input_list = [str(path) for path in config.merge_input_dirs]
        self._refresh_merge_listbox()
        self.merge_output_dir_var.set(str(config.merge_output_dir))
        self.sheet_name_var.set(config.sheet_name)
        self.year_value_var.set(config.year_value)
        self.month_value_var.set(config.month_value)
        self.stats_output_dir_var.set(str(config.stats_output_dir))
        self.calculation_mode_var.set(config.calculation_mode)
        self.output_format_var.set(config.output_format)
        self.summary_output_dir_var.set(str(config.summary_output_dir))
        self.summary_output_name_var.set(config.summary_output_name)
        self.ppt_template_path_var.set(str(config.ppt_template_path))
        self.output_ppt_path_var.set(str(config.output_ppt_path))
        self.ppt_file_pattern_var.set(config.ppt_file_pattern)
        self.ppt_sheet_name_mode_var.set(config.ppt_sheet_name_mode)
        self.ppt_sheet_name_var.set(config.ppt_sheet_name)
        self.ppt_section_mode_var.set(config.ppt_section_mode)
        self.ppt_blank_display_var.set(config.ppt_blank_display)
        self.ppt_title_suffix_var.set(config.ppt_title_suffix)
        self.ppt_max_single_table_rows_var.set(config.ppt_max_single_table_rows)
        self.ppt_max_split_table_rows_var.set(config.ppt_max_split_table_rows)
        self.ppt_sort_files_var.set(config.ppt_sort_files)
        self.ppt_body_font_size_pt_var.set(config.ppt_body_font_size_pt)
        self.ppt_header_font_size_pt_var.set(config.ppt_header_font_size_pt)
        self.ppt_summary_font_size_pt_var.set(config.ppt_summary_font_size_pt)
        self.ppt_template_slide_index_var.set(config.ppt_template_slide_index)
        self.ppt_chart_page_enabled_var.set(config.ppt_chart_page_enabled)
        self._set_text_widget_value(
            self.ppt_chart_placeholder_text_widget,
            config.ppt_chart_placeholder_text,
            self.ppt_chart_placeholder_text_var,
        )
        self.ppt_chart_image_dpi_var.set(config.ppt_chart_image_dpi)
        self.ppt_llm_notes_enabled_var.set(config.ppt_llm_notes_enabled)
        self.ppt_llm_env_path_var.set(config.ppt_llm_env_path)
        self.ppt_llm_system_role_path_var.set(config.ppt_llm_system_role_path)
        self.ppt_llm_target_chars_var.set(config.ppt_llm_target_chars)
        self.ppt_llm_temperature_var.set(config.ppt_llm_temperature)
        self.ppt_llm_max_tokens_var.set(config.ppt_llm_max_tokens)
        self.ppt_llm_checkpoint_chars_var.set(config.ppt_llm_checkpoint_chars)
        self._load_ppt_category_intro_entries_from_text(
            config.ppt_category_intro_slides_text
        )
        for region_name, _ in PPT_LAYOUT_REGION_LABELS:
            for field_name in ("left", "top", "width", "height"):
                self.ppt_layout_vars[region_name][field_name].set(
                    getattr(config, layout_field_key(region_name, field_name))
                )
        self._reset_stats_preview_state()
        self.refresh_all_status_views()

    def _restore_persistent_state(self) -> None:
        self._refresh_saved_batch_profiles()
        config = DEFAULT_GUI_BATCH_CONFIG
        active_saved_batch_name: str | None = None
        if GUI_SESSION_PATH.exists():
            try:
                config, active_saved_batch_name = load_gui_session()
            except Exception as exc:
                self.append_log(f"[WARN] 最近会话恢复失败：{exc}")
        self._apply_config_to_form(config)
        saved_batch_names = {profile.batch_name for profile in self._saved_batch_profiles}
        if active_saved_batch_name in saved_batch_names:
            self.active_saved_batch_name = active_saved_batch_name
        else:
            self.active_saved_batch_name = None
        self._refresh_saved_batch_profiles(selected_name=self.active_saved_batch_name)

    def _build_new_batch_config(self) -> GuiBatchConfig:
        return replace(self.current_config(), batch_name="未命名批次")

    def _ensure_batch_switch_allowed(self) -> bool:
        if not self.workflow_controller.start_enabled or self.runner.is_running:
            messagebox.showwarning("任务执行中", "请等待当前任务完成后再切换或删除批次。")
            return False
        return True

    def load_selected_batch(self) -> None:
        if not self._ensure_batch_switch_allowed():
            return
        selected_batch_name = self.saved_batch_var.get().strip()
        if not selected_batch_name:
            messagebox.showinfo("未选择批次", "请先从下拉列表中选择一个已保存批次。")
            return
        profile = next(
            (item for item in self._saved_batch_profiles if item.batch_name == selected_batch_name),
            None,
        )
        if profile is None:
            messagebox.showerror("批次不存在", "所选批次配置不存在，请刷新后重试。")
            self._refresh_saved_batch_profiles()
            return
        self.active_saved_batch_name = profile.batch_name
        self._apply_config_to_form(profile.config)
        self._refresh_saved_batch_profiles(selected_name=profile.batch_name)
        self._persist_session_state()
        self.append_log(f"[INFO] 已加载批次：{profile.batch_name}")

    def create_new_batch(self) -> None:
        if not self._ensure_batch_switch_allowed():
            return
        self.active_saved_batch_name = None
        self._apply_config_to_form(self._build_new_batch_config())
        self._refresh_saved_batch_profiles()
        self._persist_session_state()
        self.append_log("[INFO] 已创建新的批次草稿。")

    def save_current_batch(self) -> None:
        config = self.current_config()
        if not config.batch_name.strip():
            messagebox.showerror("批次名称为空", "请先填写批次名称后再保存。")
            return
        profile_path = save_batch_profile(config)
        self.active_saved_batch_name = config.batch_name
        self._refresh_saved_batch_profiles(selected_name=config.batch_name)
        self._persist_session_state()
        messagebox.showinfo("已保存", f"批次配置已保存到：\n{profile_path}")
        self.append_log(f"[INFO] 已保存批次：{config.batch_name}")

    def delete_selected_batch(self) -> None:
        if not self._ensure_batch_switch_allowed():
            return
        selected_batch_name = self.saved_batch_var.get().strip() or self.active_saved_batch_name or ""
        if not selected_batch_name:
            messagebox.showinfo("无可删除批次", "当前没有选中的已保存批次。")
            return
        confirmed = messagebox.askyesno(
            "删除批次",
            f"确认删除批次“{selected_batch_name}”吗？\n该操作不会删除原始数据文件。",
        )
        if not confirmed:
            return
        deleted = delete_batch_profile(selected_batch_name)
        if not deleted:
            messagebox.showerror("删除失败", "未找到对应批次配置文件，可能已经被删除。")
            self._refresh_saved_batch_profiles()
            return
        if self.active_saved_batch_name == selected_batch_name:
            self.active_saved_batch_name = None
        self._refresh_saved_batch_profiles()
        self._persist_session_state()
        self.append_log(f"[INFO] 已删除批次：{selected_batch_name}")

    def current_config(self) -> GuiBatchConfig:
        merge_dirs = tuple(Path(item) for item in self.merge_input_list)
        mode = WorkflowMode(self.workflow_mode_var.get())
        return GuiBatchConfig(
            batch_name=self.batch_name_var.get().strip() or "未命名批次",
            workflow_mode=mode,
            single_input_dir=Path(self.single_input_dir_var.get().strip() or ".").resolve(),
            merge_input_dirs=merge_dirs,
            merge_output_dir=Path(self.merge_output_dir_var.get().strip() or ".").resolve(),
            sheet_name=self.sheet_name_var.get().strip() or DEFAULT_SHEET_NAME,
            year_value=self.year_value_var.get().strip(),
            month_value=self.month_value_var.get().strip(),
            stats_output_dir=Path(self.stats_output_dir_var.get().strip() or ".").resolve(),
            calculation_mode=self.calculation_mode_var.get().strip() or "template",
            output_format=self.output_format_var.get().strip() or "xlsx",
            summary_output_dir=Path(self.summary_output_dir_var.get().strip() or ".").resolve(),
            summary_output_name=self.summary_output_name_var.get().strip() or "客户类型满意度汇总表.xlsx",
            ppt_template_path=Path(self.ppt_template_path_var.get().strip() or ".").resolve(),
            output_ppt_path=Path(self.output_ppt_path_var.get().strip() or ".").resolve(),
            ppt_file_pattern=self.ppt_file_pattern_var.get().strip() or PPT_DEFAULT_FILE_PATTERN,
            ppt_sheet_name_mode=self.ppt_sheet_name_mode_var.get().strip() or PPT_DEFAULT_SHEET_NAME_MODE,
            ppt_sheet_name=self.ppt_sheet_name_var.get(),
            ppt_section_mode=self.ppt_section_mode_var.get().strip() or "auto",
            ppt_blank_display=self.ppt_blank_display_var.get(),
            ppt_title_suffix=self.ppt_title_suffix_var.get(),
            ppt_max_single_table_rows=self.ppt_max_single_table_rows_var.get().strip()
            or PPT_DEFAULT_MAX_SINGLE_TABLE_ROWS,
            ppt_max_split_table_rows=self.ppt_max_split_table_rows_var.get().strip()
            or PPT_DEFAULT_MAX_SPLIT_TABLE_ROWS,
            ppt_sort_files=self.ppt_sort_files_var.get(),
            ppt_body_font_size_pt=self.ppt_body_font_size_pt_var.get().strip()
            or PPT_DEFAULT_BODY_FONT_SIZE_PT,
            ppt_header_font_size_pt=self.ppt_header_font_size_pt_var.get().strip()
            or PPT_DEFAULT_HEADER_FONT_SIZE_PT,
            ppt_summary_font_size_pt=self.ppt_summary_font_size_pt_var.get().strip()
            or PPT_DEFAULT_SUMMARY_FONT_SIZE_PT,
            ppt_template_slide_index=self.ppt_template_slide_index_var.get().strip()
            or PPT_DEFAULT_TEMPLATE_SLIDE_INDEX,
            ppt_chart_page_enabled=self.ppt_chart_page_enabled_var.get(),
            ppt_chart_placeholder_text=self._read_text_widget_value(
                self.ppt_chart_placeholder_text_widget,
                self.ppt_chart_placeholder_text_var,
            ),
            ppt_chart_image_dpi=self.ppt_chart_image_dpi_var.get().strip()
            or PPT_DEFAULT_CHART_IMAGE_DPI,
            ppt_llm_notes_enabled=self.ppt_llm_notes_enabled_var.get(),
            ppt_llm_env_path=self.ppt_llm_env_path_var.get().strip() or PPT_DEFAULT_LLM_ENV_PATH,
            ppt_llm_system_role_path=self.ppt_llm_system_role_path_var.get().strip()
            or PPT_DEFAULT_LLM_SYSTEM_ROLE_PATH,
            ppt_llm_target_chars=self.ppt_llm_target_chars_var.get().strip()
            or PPT_DEFAULT_LLM_TARGET_CHARS,
            ppt_llm_temperature=self.ppt_llm_temperature_var.get().strip()
            or PPT_DEFAULT_LLM_TEMPERATURE,
            ppt_llm_max_tokens=self.ppt_llm_max_tokens_var.get().strip()
            or PPT_DEFAULT_LLM_MAX_TOKENS,
            ppt_llm_checkpoint_chars=self.ppt_llm_checkpoint_chars_var.get().strip()
            or PPT_DEFAULT_LLM_CHECKPOINT_CHARS,
            ppt_category_intro_slides_text=build_category_intro_slides_text(
                self.ppt_category_intro_entries
            ),
            ppt_layout_summary_table_left=self.ppt_layout_vars["summary_table"]["left"].get().strip()
            or PPT_LAYOUT_DEFAULTS["summary_table"][0],
            ppt_layout_summary_table_top=self.ppt_layout_vars["summary_table"]["top"].get().strip()
            or PPT_LAYOUT_DEFAULTS["summary_table"][1],
            ppt_layout_summary_table_width=self.ppt_layout_vars["summary_table"]["width"].get().strip()
            or PPT_LAYOUT_DEFAULTS["summary_table"][2],
            ppt_layout_summary_table_height=self.ppt_layout_vars["summary_table"]["height"].get().strip()
            or PPT_LAYOUT_DEFAULTS["summary_table"][3],
            ppt_layout_detail_single_table_left=self.ppt_layout_vars["detail_single_table"]["left"].get().strip()
            or PPT_LAYOUT_DEFAULTS["detail_single_table"][0],
            ppt_layout_detail_single_table_top=self.ppt_layout_vars["detail_single_table"]["top"].get().strip()
            or PPT_LAYOUT_DEFAULTS["detail_single_table"][1],
            ppt_layout_detail_single_table_width=self.ppt_layout_vars["detail_single_table"]["width"].get().strip()
            or PPT_LAYOUT_DEFAULTS["detail_single_table"][2],
            ppt_layout_detail_single_table_height=self.ppt_layout_vars["detail_single_table"]["height"].get().strip()
            or PPT_LAYOUT_DEFAULTS["detail_single_table"][3],
            ppt_layout_detail_left_table_left=self.ppt_layout_vars["detail_left_table"]["left"].get().strip()
            or PPT_LAYOUT_DEFAULTS["detail_left_table"][0],
            ppt_layout_detail_left_table_top=self.ppt_layout_vars["detail_left_table"]["top"].get().strip()
            or PPT_LAYOUT_DEFAULTS["detail_left_table"][1],
            ppt_layout_detail_left_table_width=self.ppt_layout_vars["detail_left_table"]["width"].get().strip()
            or PPT_LAYOUT_DEFAULTS["detail_left_table"][2],
            ppt_layout_detail_left_table_height=self.ppt_layout_vars["detail_left_table"]["height"].get().strip()
            or PPT_LAYOUT_DEFAULTS["detail_left_table"][3],
            ppt_layout_detail_right_table_left=self.ppt_layout_vars["detail_right_table"]["left"].get().strip()
            or PPT_LAYOUT_DEFAULTS["detail_right_table"][0],
            ppt_layout_detail_right_table_top=self.ppt_layout_vars["detail_right_table"]["top"].get().strip()
            or PPT_LAYOUT_DEFAULTS["detail_right_table"][1],
            ppt_layout_detail_right_table_width=self.ppt_layout_vars["detail_right_table"]["width"].get().strip()
            or PPT_LAYOUT_DEFAULTS["detail_right_table"][2],
            ppt_layout_detail_right_table_height=self.ppt_layout_vars["detail_right_table"]["height"].get().strip()
            or PPT_LAYOUT_DEFAULTS["detail_right_table"][3],
            ppt_layout_chart_image_left=self.ppt_layout_vars["chart_image"]["left"].get().strip()
            or PPT_LAYOUT_DEFAULTS["chart_image"][0],
            ppt_layout_chart_image_top=self.ppt_layout_vars["chart_image"]["top"].get().strip()
            or PPT_LAYOUT_DEFAULTS["chart_image"][1],
            ppt_layout_chart_image_width=self.ppt_layout_vars["chart_image"]["width"].get().strip()
            or PPT_LAYOUT_DEFAULTS["chart_image"][2],
            ppt_layout_chart_image_height=self.ppt_layout_vars["chart_image"]["height"].get().strip()
            or PPT_LAYOUT_DEFAULTS["chart_image"][3],
            ppt_layout_chart_textbox_left=self.ppt_layout_vars["chart_textbox"]["left"].get().strip()
            or PPT_LAYOUT_DEFAULTS["chart_textbox"][0],
            ppt_layout_chart_textbox_top=self.ppt_layout_vars["chart_textbox"]["top"].get().strip()
            or PPT_LAYOUT_DEFAULTS["chart_textbox"][1],
            ppt_layout_chart_textbox_width=self.ppt_layout_vars["chart_textbox"]["width"].get().strip()
            or PPT_LAYOUT_DEFAULTS["chart_textbox"][2],
            ppt_layout_chart_textbox_height=self.ppt_layout_vars["chart_textbox"]["height"].get().strip()
            or PPT_LAYOUT_DEFAULTS["chart_textbox"][3],
        )

    def choose_directory(self, variable: tk.StringVar) -> None:
        initial = variable.get().strip() or str(PROJECT_ROOT)
        selected = filedialog.askdirectory(initialdir=initial)
        if selected:
            variable.set(selected)
            self.refresh_all_status_views()

    def choose_file(self, variable: tk.StringVar, filetypes: list[tuple[str, str]]) -> None:
        initial = variable.get().strip() or str(PROJECT_ROOT)
        selected = filedialog.askopenfilename(initialdir=str(Path(initial).parent), filetypes=filetypes)
        if selected:
            variable.set(selected)
            self.refresh_all_status_views()

    def choose_save_file(self, variable: tk.StringVar, default_ext: str, filetypes: list[tuple[str, str]]) -> None:
        initial = variable.get().strip() or str(PROJECT_ROOT / f"output{default_ext}")
        selected = filedialog.asksaveasfilename(initialfile=Path(initial).name, initialdir=str(Path(initial).parent), defaultextension=default_ext, filetypes=filetypes)
        if selected:
            variable.set(selected)
            self.refresh_all_status_views()

    def add_merge_input_dir(self) -> None:
        selected = filedialog.askdirectory(initialdir=str(PROJECT_ROOT))
        if not selected:
            return
        if selected in self.merge_input_list:
            messagebox.showinfo("目录已存在", "该目录已经添加过了。")
            return
        self.merge_input_list.append(selected)
        self.merge_listbox.insert("end", selected)
        self.refresh_all_status_views()

    def remove_selected_merge_dir(self) -> None:
        selection = self.merge_listbox.curselection()
        if not selection:
            return
        for index in reversed(selection):
            del self.merge_input_list[index]
            self.merge_listbox.delete(index)
        self.refresh_all_status_views()

    def refresh_all_status_views(self) -> None:
        config = self.current_config()
        effective_input = config.effective_input_dir()
        stats_files = discover_excel_files(config.stats_output_dir)
        self.effective_input_var.set(str(effective_input))
        self.stats_effective_input_var.set(str(effective_input))
        self._raw_file_count_var.set(str(len(discover_excel_files(effective_input))))
        self._stats_file_count_var.set(str(len(stats_files)))
        summary_path = config.summary_output_dir / config.summary_output_name
        self._summary_file_var.set("已生成" if summary_path.exists() else "未生成")
        self._ppt_file_var.set("已生成" if config.output_ppt_path.exists() else "未生成")

        self._step_status_vars["data_source"].set("已完成" if effective_input.exists() else "待设置")
        if stats_files:
            self._step_status_vars["stats"].set("已完成")
        if summary_path.exists():
            self._step_status_vars["summary"].set("已完成")
        if config.output_ppt_path.exists():
            self._step_status_vars["ppt"].set("已完成")
        self.stats_summary_var.set(
            f"输入目录：{effective_input}；输出目录：{config.stats_output_dir}；当前分项文件数：{len(stats_files)}"
        )
        self.summary_status_var.set(
            f"目标汇总表：{summary_path}"
        )
        self.ppt_status_var.set(
            f"模板：{config.ppt_template_path}；输出：{config.output_ppt_path}"
        )

        self._refresh_treeviews(config)
        self._recompute_next_action(config)

    def _refresh_treeviews(self, config: GuiBatchConfig) -> None:
        for item in self.data_source_tree.get_children():
            self.data_source_tree.delete(item)
        data_source_paths: list[Path] = []
        if config.workflow_mode == WorkflowMode.SINGLE:
            data_source_paths.append(config.single_input_dir)
        else:
            data_source_paths.extend(config.merge_input_dirs)
            data_source_paths.append(config.merge_output_dir)
        for path in data_source_paths:
            status = "存在" if Path(path).exists() else "缺失"
            self.data_source_tree.insert("", "end", values=(str(path), status))

        for item in self.stats_tree.get_children():
            self.stats_tree.delete(item)
        for path in discover_excel_files(config.stats_output_dir):
            self.stats_tree.insert("", "end", values=(path.name, str(path)))

    def _recompute_next_action(self, config: GuiBatchConfig) -> None:
        if not config.effective_input_dir().exists():
            self._next_action_var.set("先配置并确认数据源")
            return
        if self._step_status_vars["preprocess"].get() == "待执行":
            self._next_action_var.set("建议先执行预处理")
            return
        if self._step_status_vars["stats"].get() == "待执行":
            self._next_action_var.set("建议先生成分项统计")
            return
        if self._step_status_vars["summary"].get() == "待执行":
            self._next_action_var.set("建议先生成汇总表")
            return
        if self._step_status_vars["ppt"].get() == "待执行":
            self._next_action_var.set("建议生成PPT")
            return
        self._next_action_var.set("本批次主流程已完成")

    def _update_task_status_text(self) -> None:
        self._task_status_label_var.set(
            build_workflow_status_text(self.workflow_controller, self._task_title_lookup)
        )

    def _set_stats_preview_selection(self, selected_customer_types: frozenset[str]) -> None:
        valid_customer_types = {
            row.customer_type_name
            for row in self.stats_preview_rows
        }
        self.selected_customer_types = frozenset(
            name for name in selected_customer_types if name in valid_customer_types
        )
        self._refresh_stats_preview_tree()
        self._refresh_stats_preview_summary()

    def _apply_stats_preview_summary(self, summary: StatsPreviewSummary) -> None:
        self.stats_preview_summary = summary
        self.stats_preview_rows = summary.rows
        self.selected_customer_types = default_selected_customer_types(summary)
        self._refresh_stats_preview_tree()
        self._refresh_stats_preview_summary()

    def _refresh_stats_preview_summary(self) -> None:
        if self.stats_preview_summary is None:
            self.stats_preview_summary_var.set("未扫描客群")
            return
        self.stats_preview_summary_var.set(
            build_stats_preview_summary_text(
                self.stats_preview_summary,
                self.selected_customer_types,
            )
        )

    def _refresh_stats_preview_tree(self) -> None:
        for item in self.stats_preview_tree.get_children():
            self.stats_preview_tree.delete(item)
        for row in self.stats_preview_rows:
            selected_marker = "✓" if row.customer_type_name in self.selected_customer_types else ""
            self.stats_preview_tree.insert(
                "",
                "end",
                iid=row.customer_type_name,
                values=(
                    selected_marker,
                    row.customer_type_name,
                    row.source_file_name,
                    customer_type_preview_status_text(row.status),
                    row.output_name,
                    row.detail,
                ),
            )

    def select_all_customer_types(self) -> None:
        self._set_stats_preview_selection(
            frozenset(row.customer_type_name for row in self.stats_preview_rows)
        )

    def clear_selected_customer_types(self) -> None:
        self._set_stats_preview_selection(frozenset())

    def select_ready_customer_types(self) -> None:
        if self.stats_preview_summary is None:
            return
        self._set_stats_preview_selection(default_selected_customer_types(self.stats_preview_summary))

    def on_stats_preview_row_double_click(self, event: tk.Event[tk.Misc]) -> None:
        item_id = self.stats_preview_tree.identify_row(event.y)
        if not item_id:
            return
        selected_customer_types = set(self.selected_customer_types)
        if item_id in selected_customer_types:
            selected_customer_types.remove(item_id)
        else:
            selected_customer_types.add(item_id)
        self._set_stats_preview_selection(frozenset(selected_customer_types))

    def _ensure_stats_preview_loaded(self, *, parent: tk.Misc | None = None) -> bool:
        if self.stats_preview_summary is not None:
            return True
        try:
            summary = build_stats_preview_summary(self.current_config())
        except Exception as exc:
            if parent is None:
                messagebox.showerror("扫描失败", str(exc))
            else:
                messagebox.showerror("扫描失败", str(exc), parent=parent)
            return False
        self._apply_stats_preview_summary(summary)
        return True

    def _selected_stats_job_names(self, *, parent: tk.Misc | None = None) -> tuple[str, ...] | None:
        if not self._ensure_stats_preview_loaded(parent=parent):
            return None
        selected_job_names = ordered_selected_customer_types(
            self.stats_preview_rows,
            self.selected_customer_types,
        )
        if selected_job_names:
            return selected_job_names
        if parent is None:
            messagebox.showinfo("未选择客群", "请至少勾选一个客群后再执行分项统计。")
        else:
            messagebox.showinfo("未选择客群", "请至少勾选一个客群后再执行分项统计。", parent=parent)
        return None

    def scan_stats_preview(self) -> None:
        try:
            summary = build_stats_preview_summary(self.current_config())
        except Exception as exc:
            messagebox.showerror("扫描失败", str(exc))
            return
        self._apply_stats_preview_summary(summary)

    def terminate_current_task(self) -> None:
        if not self.workflow_controller.terminate_enabled:
            messagebox.showinfo("无可终止任务", "当前没有正在执行的任务。")
            return
        self.workflow_controller.request_cancel()
        self._update_task_status_text()
        if self.workflow_controller.active_step_key is not None:
            self._set_workflow_progress_row(
                self.workflow_controller.active_step_key,
                "正在终止",
                "已发送终止请求，等待当前脚本退出",
            )
        self._refresh_workflow_progress_summary()
        self._sync_action_button_states()
        self.append_log("[INFO] 已请求终止当前任务。")
        self.runner.terminate()

    def append_log(self, message: str) -> None:
        self.log_text.configure(state="normal")
        self.log_text.insert("end", message + "\n")
        self.log_text.see("end")
        self.log_text.configure(state="disabled")

    def clear_log(self) -> None:
        self.log_text.configure(state="normal")
        self.log_text.delete("1.0", "end")
        self.log_text.configure(state="disabled")

    def _destroy_workflow_progress_dialog(self) -> None:
        dialog = self._workflow_progress_dialog
        self._workflow_progress_dialog = None
        self._workflow_progress_tree = None
        self._workflow_progress_item_ids = {}
        self._workflow_progress_last_page_key = None
        self._workflow_progress_summary_var.set("未启动主流程")
        self._workflow_progress_close_var.set("关闭")
        if dialog is not None and dialog.winfo_exists():
            dialog.destroy()

    def _hide_or_close_workflow_progress_dialog(self) -> None:
        dialog = self._workflow_progress_dialog
        if dialog is None or not dialog.winfo_exists():
            return
        if self.workflow_controller.status in {WorkflowRunStatus.RUNNING, WorkflowRunStatus.CANCELLING}:
            dialog.withdraw()
            return
        self._destroy_workflow_progress_dialog()

    def _open_workflow_progress_dialog(self, tasks: tuple[TaskCommand, ...]) -> None:
        self._destroy_workflow_progress_dialog()
        dialog = tk.Toplevel(self)
        dialog.title("主流程执行进度")
        dialog.transient(self)
        dialog.configure(bg=self.palette.background)
        dialog.geometry("980x420")
        dialog.minsize(860, 360)
        dialog.protocol("WM_DELETE_WINDOW", self._hide_or_close_workflow_progress_dialog)

        wrap = ttk.Frame(dialog, style="Surface.TFrame", padding=16)
        wrap.pack(fill="both", expand=True)
        wrap.columnconfigure(0, weight=1)
        wrap.rowconfigure(1, weight=1)

        ttk.Label(wrap, text="正在执行主流程", style="SubHeader.TLabel").grid(row=0, column=0, sticky="w")
        ttk.Label(wrap, textvariable=self._workflow_progress_summary_var, style="Muted.TLabel").grid(row=0, column=1, sticky="e")

        tree = ttk.Treeview(
            wrap,
            columns=("step", "status", "detail"),
            show="headings",
            height=10,
        )
        tree.heading("step", text="步骤")
        tree.heading("status", text="状态")
        tree.heading("detail", text="进度明细")
        tree.column("step", width=220)
        tree.column("status", width=120, anchor="center")
        tree.column("detail", width=580)
        tree.grid(row=1, column=0, columnspan=2, sticky="nsew", pady=(12, 12))

        footer = ttk.Frame(wrap, style="Surface.TFrame")
        footer.grid(row=2, column=0, columnspan=2, sticky="ew")
        footer.columnconfigure(0, weight=1)
        ttk.Button(footer, text="查看任务日志", style="Secondary.TButton", command=lambda: self.show_page("logs")).grid(row=0, column=0, sticky="w")
        close_button = ttk.Button(
            footer,
            textvariable=self._workflow_progress_close_var,
            style="Secondary.TButton",
            command=self._hide_or_close_workflow_progress_dialog,
        )
        close_button.grid(row=0, column=1, padx=(0, 8))
        terminate_button = ttk.Button(
            footer,
            text="终止当前任务",
            style="Primary.TButton",
            command=self.terminate_current_task,
        )
        terminate_button.grid(row=0, column=2)
        self._register_terminate_button(terminate_button)

        self._workflow_progress_dialog = dialog
        self._workflow_progress_tree = tree
        self._workflow_progress_item_ids = {}
        self._workflow_progress_last_page_key = page_key_for_task(tasks[-1].key) if tasks else None
        for index, task in enumerate(tasks, start=1):
            item_id = task.key
            self._workflow_progress_item_ids[task.key] = item_id
            tree.insert(
                "",
                "end",
                iid=item_id,
                values=(f"{index}. {task.title}", "待执行", "等待前置步骤"),
            )
        self._workflow_progress_summary_var.set(f"共 {len(tasks)} 步，等待启动")
        self._workflow_progress_close_var.set("后台运行")
        dialog.lift()
        self._sync_action_button_states()

    def _set_workflow_progress_row(
        self,
        task_key: str,
        status_text: str,
        detail_text: str,
    ) -> None:
        if self._workflow_progress_tree is None or not self._workflow_progress_tree.winfo_exists():
            return
        item_id = self._workflow_progress_item_ids.get(task_key)
        if item_id is None:
            return
        current_values = self._workflow_progress_tree.item(item_id, "values")
        if not current_values:
            return
        self._workflow_progress_tree.item(
            item_id,
            values=(current_values[0], status_text, detail_text),
        )

    def _refresh_workflow_progress_summary(self) -> None:
        summary = build_workflow_status_text(self.workflow_controller, self._task_title_lookup)
        total = len(self.workflow_controller.planned_step_keys)
        completed = len(self.workflow_controller.completed_step_keys)
        if total:
            summary = f"{summary}；进度 {completed}/{total}"
        self._workflow_progress_summary_var.set(summary)
        if self.workflow_controller.status in {WorkflowRunStatus.RUNNING, WorkflowRunStatus.CANCELLING}:
            self._workflow_progress_close_var.set("后台运行")
        else:
            self._workflow_progress_close_var.set("关闭")

    def _mark_pending_workflow_rows(self, status_text: str, detail_text: str) -> None:
        if self._workflow_progress_tree is None or not self._workflow_progress_tree.winfo_exists():
            return
        for task_key in self.workflow_controller.planned_step_keys:
            if task_key in self.workflow_controller.completed_step_keys:
                continue
            if task_key == self.workflow_controller.failed_step_key:
                continue
            item_id = self._workflow_progress_item_ids.get(task_key)
            if item_id is None:
                continue
            current_values = self._workflow_progress_tree.item(item_id, "values")
            if current_values and current_values[1] == "待执行":
                self._set_workflow_progress_row(task_key, status_text, detail_text)

    def run_task_list(
        self,
        tasks: tuple[TaskCommand, ...],
        *,
        show_workflow_progress: bool = False,
    ) -> None:
        if not tasks:
            messagebox.showinfo("无可执行任务", "当前没有需要执行的任务。")
            return
        if not self.workflow_controller.start_enabled or self.runner.is_running:
            messagebox.showwarning("任务执行中", "请等待当前任务完成后再继续。")
            return
        if show_workflow_progress:
            self._open_workflow_progress_dialog(tasks)
        else:
            self._destroy_workflow_progress_dialog()
        for task in tasks:
            self.append_log(f"[CMD] {task.title}: {shlex.join(task.command)}")
        self._task_title_lookup = {task.key: task.title for task in tasks}
        self.workflow_controller.begin(tuple(task.key for task in tasks))
        self._update_task_status_text()
        self._refresh_workflow_progress_summary()
        self._sync_action_button_states()
        self.runner.run_tasks(tasks)

    def run_merge_task(self) -> None:
        try:
            command = build_merge_command(self.current_config())
            if command is None:
                messagebox.showinfo("无需合并", "当前是单月模式，无需执行合并。")
                return
            self.run_task_list((TaskCommand("merge_workbooks", "合并多月问卷", tuple(command)),))
        except ValueError as exc:
            messagebox.showerror("参数不完整", str(exc))

    def run_phase_preprocess_task(self) -> None:
        self._run_single_builder_task("phase_preprocess", PHASE_PREPROCESS_TASK_TITLE, build_phase_preprocess_command)

    def run_fill_year_month_task(self) -> None:
        self._run_single_builder_task("fill_year_month", FILL_YEAR_MONTH_TASK_TITLE, build_fill_year_month_command)

    def run_survey_stats_task(self) -> None:
        selected_job_names = self._selected_stats_job_names()
        if selected_job_names is None:
            return
        self._run_single_command(
            TaskCommand(
                "survey_stats",
                "生成分项统计",
                tuple(
                    build_survey_stats_command(
                        self.current_config(),
                        selected_job_names=selected_job_names,
                    )
                ),
            )
        )

    def run_survey_stats_dry_run_task(self) -> None:
        selected_job_names = self._selected_stats_job_names()
        if selected_job_names is None:
            return
        self._run_single_command(
            TaskCommand(
                "survey_stats",
                "分项统计 dry-run",
                tuple(
                    build_survey_stats_command(
                        self.current_config(),
                        dry_run=True,
                        selected_job_names=selected_job_names,
                    )
                ),
            )
        )

    def run_summary_task(self) -> None:
        self._run_single_command(
            TaskCommand(
                "summary_table",
                "生成汇总表",
                tuple(build_summary_command(self.current_config())),
            )
        )

    def run_ppt_task(self) -> None:
        self._run_single_command(
            TaskCommand(
                "generate_ppt",
                "生成PPT",
                tuple(build_ppt_command(self.current_config())),
            )
        )

    def run_ppt_dry_run_task(self) -> None:
        self._run_single_command(
            TaskCommand(
                "generate_ppt",
                "PPT dry-run",
                tuple(build_ppt_command(self.current_config(), dry_run=True)),
            )
        )

    def _run_single_builder_task(
        self,
        key: str,
        title: str,
        builder,
    ) -> None:
        try:
            command = builder(self.current_config())
        except ValueError as exc:
            messagebox.showerror("参数不完整", str(exc))
            return
        self._run_single_command(TaskCommand(key, title, tuple(command)))

    def _run_single_command(self, task: TaskCommand) -> None:
        self.run_task_list((task,))

    def on_task_started(self, task: TaskCommand) -> None:
        self.workflow_controller.mark_task_started(task.key)
        self._update_task_status_text()
        self._set_workflow_progress_row(task.key, "执行中", shlex.join(task.command))
        self._refresh_workflow_progress_summary()
        self.append_log(f"[START] {task.title}")

    def on_task_finished(self, task: TaskCommand, success: bool) -> None:
        self.workflow_controller.mark_task_finished(task.key, success)
        if success:
            self.append_log(f"[DONE] {task.title}")
            self._set_workflow_progress_row(task.key, "已完成", "执行完成，可继续下一步")
            if task.key in {"merge_workbooks"}:
                self._step_status_vars["data_source"].set("已完成")
            elif task.key in {"phase_preprocess", "fill_year_month"}:
                self._step_status_vars["preprocess"].set("已完成")
            elif task.key == "survey_stats":
                self._step_status_vars["stats"].set("已完成")
            elif task.key == "summary_table":
                self._step_status_vars["summary"].set("已完成")
            elif task.key == "generate_ppt":
                self._step_status_vars["ppt"].set("已完成")
        else:
            if self.workflow_controller.cancel_requested:
                self.append_log(f"[CANCEL] {task.title}")
                self._set_workflow_progress_row(task.key, "已取消", "用户已终止当前步骤")
            else:
                self.append_log(f"[FAIL] {task.title}")
                self._set_workflow_progress_row(task.key, "失败", "执行失败，请查看日志")
        self._update_task_status_text()
        self._refresh_workflow_progress_summary()
        self.refresh_all_status_views()

    def on_all_tasks_finished(self, success: bool) -> None:
        self.workflow_controller.finish_run(success)
        self._update_task_status_text()
        if self.workflow_controller.status == WorkflowRunStatus.CANCELLED:
            self._mark_pending_workflow_rows("已取消", "主流程已终止，未再执行")
        elif self.workflow_controller.status == WorkflowRunStatus.FAILED:
            self._mark_pending_workflow_rows("未执行", "前序步骤失败，后续步骤未启动")
        self._refresh_workflow_progress_summary()
        self._sync_action_button_states()
        if success:
            self.append_log("[INFO] 所有任务执行完成。")
        elif self.workflow_controller.status == WorkflowRunStatus.CANCELLED:
            self.append_log("[INFO] 主流程已取消。")
        self.refresh_all_status_views()
        if self._workflow_progress_last_page_key is not None:
            self.show_page(self._workflow_progress_last_page_key)

    def save_profile_snapshot(self) -> None:
        self.save_current_batch()

    def on_app_close(self) -> None:
        try:
            self._persist_session_state()
        except Exception as exc:
            self.append_log(f"[WARN] 最近会话保存失败：{exc}")
        self.destroy()

    def open_main_workflow_dialog(self) -> None:
        if not self.workflow_controller.start_enabled:
            messagebox.showwarning("任务执行中", "请先等待当前任务完成或手动终止。")
            return
        config = self.current_config()
        dialog = tk.Toplevel(self)
        dialog.title("一键执行主流程")
        dialog.transient(self)
        dialog.grab_set()
        dialog.configure(bg=self.palette.background)

        include_merge_var = tk.BooleanVar(value=config.workflow_mode == WorkflowMode.MERGED)
        include_phase_var = tk.BooleanVar(value=True)
        include_fill_var = tk.BooleanVar(value=False)
        include_stats_var = tk.BooleanVar(value=True)
        include_summary_var = tk.BooleanVar(value=True)
        include_ppt_var = tk.BooleanVar(value=True)

        wrap = ttk.Frame(dialog, style="Surface.TFrame", padding=16)
        wrap.pack(fill="both", expand=True)
        ttk.Label(wrap, text="当前批次主流程", style="SubHeader.TLabel").pack(anchor="w")
        ttk.Label(wrap, text=f"批次：{config.batch_name}", style="Muted.TLabel").pack(anchor="w", pady=(4, 0))
        ttk.Label(wrap, text=f"输入目录：{config.effective_input_dir()}", style="Muted.TLabel").pack(anchor="w", pady=(0, 12))
        ttk.Label(
            wrap,
            text="若包含“生成分项统计”，将只运行分项统计页当前已勾选的客群。",
            style="Muted.TLabel",
        ).pack(anchor="w", pady=(0, 12))

        options = ttk.LabelFrame(wrap, text="执行步骤", padding=12)
        options.pack(fill="x")
        if config.workflow_mode == WorkflowMode.MERGED:
            ttk.Checkbutton(options, text="先合并多月问卷", variable=include_merge_var).pack(anchor="w")
        ttk.Checkbutton(options, text=PHASE_PREPROCESS_WORKFLOW_TEXT, variable=include_phase_var).pack(anchor="w")
        ttk.Checkbutton(options, text=FILL_YEAR_MONTH_WORKFLOW_TEXT, variable=include_fill_var).pack(anchor="w")
        ttk.Checkbutton(options, text="生成分项统计", variable=include_stats_var).pack(anchor="w")
        ttk.Checkbutton(options, text="生成汇总表", variable=include_summary_var).pack(anchor="w")
        ttk.Checkbutton(options, text="生成PPT", variable=include_ppt_var).pack(anchor="w")

        footer = ttk.Frame(wrap, style="Surface.TFrame")
        footer.pack(fill="x", pady=(16, 0))
        footer.columnconfigure(0, weight=1)

        def start() -> None:
            selection = MainWorkflowSelection(
                include_merge=include_merge_var.get(),
                include_phase_preprocess=include_phase_var.get(),
                include_fill_year_month=include_fill_var.get(),
                include_survey_stats=include_stats_var.get(),
                include_summary=include_summary_var.get(),
                include_ppt=include_ppt_var.get(),
            )
            selected_stats_job_names: tuple[str, ...] = ()
            if selection.include_survey_stats:
                selected_stats_job_names = self._selected_stats_job_names(parent=dialog)
                if selected_stats_job_names is None:
                    return
            try:
                tasks = build_task_commands(
                    self.current_config(),
                    selection,
                    selected_stats_job_names=selected_stats_job_names,
                )
            except ValueError as exc:
                messagebox.showerror("参数不完整", str(exc), parent=dialog)
                return
            dialog.destroy()
            self.run_task_list(tasks, show_workflow_progress=True)

        ttk.Button(footer, text="取消", style="Secondary.TButton", command=dialog.destroy).grid(row=0, column=1, padx=(0, 8))
        ttk.Button(footer, text="开始执行", style="Primary.TButton", command=start).grid(row=0, column=2)


def main() -> None:
    app = SurveyPlatformApp()
    app.mainloop()


if __name__ == "__main__":
    main()
